from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os
from datetime import datetime

class PresentationTemplate:
    def __init__(self, template_name="default"):
        self.template_name = template_name
        self.template_path = os.path.join("templates", "powerpoint", f"{template_name}.pptx")
        
        # Définition des styles par défaut
        self.styles = {
            "default": {
                "colors": {
                    "primary": RGBColor(0, 112, 192),    # Bleu
                    "secondary": RGBColor(255, 255, 255), # Blanc
                    "accent": RGBColor(0, 176, 80),      # Vert
                    "text": RGBColor(0, 0, 0)            # Noir
                },
                "fonts": {
                    "title": "Calibri",
                    "body": "Calibri"
                },
                "sizes": {
                    "title": Pt(44),
                    "subtitle": Pt(32),
                    "heading": Pt(28),
                    "body": Pt(18)
                },
                "logo_path": None
            },
            "olaqin": {
                "colors": {
                    "primary": RGBColor(0, 82, 147),     # Bleu Olaqin
                    "secondary": RGBColor(255, 255, 255), # Blanc
                    "accent": RGBColor(238, 127, 0),     # Orange Olaqin
                    "text": RGBColor(51, 51, 51)         # Gris foncé
                },
                "fonts": {
                    "title": "Montserrat",
                    "body": "Open Sans"
                },
                "sizes": {
                    "title": Pt(44),
                    "subtitle": Pt(32),
                    "heading": Pt(28),
                    "body": Pt(18)
                },
                "logo_path": "templates/powerpoint/olaqin_logo.png"
            }
        }
        
        self.current_style = self.styles.get(template_name, self.styles["default"])
    
    def create_presentation(self):
        """Crée une nouvelle présentation avec le style choisi"""
        prs = Presentation()
        return prs
    
    def apply_title_slide_style(self, slide):
        """Applique le style au slide de titre"""
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        # Style du titre
        title.text_frame.paragraphs[0].font.name = self.current_style["fonts"]["title"]
        title.text_frame.paragraphs[0].font.size = self.current_style["sizes"]["title"]
        title.text_frame.paragraphs[0].font.color.rgb = self.current_style["colors"]["primary"]
        
        # Style du sous-titre
        subtitle.text_frame.paragraphs[0].font.name = self.current_style["fonts"]["body"]
        subtitle.text_frame.paragraphs[0].font.size = self.current_style["sizes"]["subtitle"]
        subtitle.text_frame.paragraphs[0].font.color.rgb = self.current_style["colors"]["text"]
        
        # Ajouter le logo si disponible
        if self.current_style["logo_path"] and os.path.exists(self.current_style["logo_path"]):
            left = Inches(9)  # Position à droite
            top = Inches(0.5) # Position en haut
            slide.shapes.add_picture(self.current_style["logo_path"], left, top, height=Inches(1))
    
    def apply_content_slide_style(self, slide):
        """Applique le style à un slide de contenu"""
        title = slide.shapes.title
        
        # Style du titre
        title.text_frame.paragraphs[0].font.name = self.current_style["fonts"]["title"]
        title.text_frame.paragraphs[0].font.size = self.current_style["sizes"]["heading"]
        title.text_frame.paragraphs[0].font.color.rgb = self.current_style["colors"]["primary"]
    
    def apply_table_style(self, table):
        """Applique le style à un tableau"""
        # Style de l'en-tête (première ligne)
        for cell in table.rows[0].cells:
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.current_style["colors"]["primary"]
            tf = cell.text_frame
            tf.paragraphs[0].font.name = self.current_style["fonts"]["body"]
            tf.paragraphs[0].font.size = self.current_style["sizes"]["body"]
            tf.paragraphs[0].font.color.rgb = self.current_style["colors"]["secondary"]
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Style du contenu (toutes les lignes sauf la première)
        for row_idx in range(1, len(table.rows)):
            for cell in table.rows[row_idx].cells:
                tf = cell.text_frame
                tf.paragraphs[0].font.name = self.current_style["fonts"]["body"]
                tf.paragraphs[0].font.size = self.current_style["sizes"]["body"]
                tf.paragraphs[0].font.color.rgb = self.current_style["colors"]["text"]
                tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    def apply_chart_style(self, chart):
        """Applique le style à un graphique"""
        # À implémenter selon les besoins spécifiques des graphiques
        pass 

class OlaqinTemplate:
    # Couleurs Olaqin
    COLORS = {
        'violet_fonce': RGBColor(107, 44, 145),  # #6B2C91
        'violet_clair': RGBColor(178, 140, 200),  # #B28CC8
        'vert': RGBColor(0, 166, 81),  # #00A651
        'rouge': RGBColor(255, 0, 0),  # #FF0000
        'blanc': RGBColor(255, 255, 255),  # #FFFFFF
        'noir': RGBColor(0, 0, 0),  # #000000
        'gris': RGBColor(128, 128, 128)  # #808080
    }

    def __init__(self):
        self.prs = Presentation()
        self._set_slide_dimensions()
        self._add_layouts()

    def _set_slide_dimensions(self):
        self.prs.slide_width = Inches(13.33)
        self.prs.slide_height = Inches(7.5)

    def _add_layouts(self):
        # Nous garderons les layouts par défaut pour l'instant
        pass

    def add_title_slide(self, title, subtitle=None):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[0])
        
        # Ajouter le logo Olaqin en haut à gauche
        left = Inches(0.5)
        top = Inches(0.5)
        slide.shapes.add_picture("assets/olaqin_logo.png", left, top, height=Inches(0.6))
        
        # Titre
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(8), Inches(1))
        title_frame = title_box.text_frame
        title_para = title_frame.add_paragraph()
        title_para.text = title
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        
        # Sous-titre
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(8), Inches(0.75))
            subtitle_frame = subtitle_box.text_frame
            subtitle_para = subtitle_frame.add_paragraph()
            subtitle_para.text = subtitle
            subtitle_para.font.size = Pt(32)
        
        # Logo hexagonal en filigrane
        slide.shapes.add_picture("assets/olaqin_hex.png", Inches(9), Inches(1), height=Inches(5))
        
        self._add_footer(slide)
        return slide

    def add_content_slide(self, title):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        
        # Titre
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(0.75))
        title_frame = title_box.text_frame
        title_para = title_frame.add_paragraph()
        title_para.text = title
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        
        self._add_footer(slide)
        return slide

    def _add_footer(self, slide):
        # Logo Olaqin
        left = Inches(0.5)
        top = slide.shapes.title.top + Inches(6.5)
        slide.shapes.add_picture("assets/olaqin_logo_small.png", left, top, height=Inches(0.3))
        
        # Date
        date_box = slide.shapes.add_textbox(Inches(2), top, Inches(2), Inches(0.3))
        date_frame = date_box.text_frame
        date_para = date_frame.add_paragraph()
        date_para.text = datetime.now().strftime("%d/%m/%Y")
        date_para.font.size = Pt(10)
        
        # Mention confidentielle
        conf_box = slide.shapes.add_textbox(Inches(4), top, Inches(5), Inches(0.3))
        conf_frame = conf_box.text_frame
        conf_para = conf_frame.add_paragraph()
        conf_para.text = "CONFIDENTIEL – NE PAS DIFFUSER ©Olaqin"
        conf_para.font.size = Pt(10)
        conf_para.font.color.rgb = self.COLORS['violet_clair']
        conf_para.alignment = PP_ALIGN.CENTER
        
        # Logo hexagonal
        slide.shapes.add_picture("assets/olaqin_hex_small.png", Inches(12), top, height=Inches(0.3))

    def create_table(self, slide, rows, cols, left=Inches(0.5), top=Inches(1.5), width=Inches(12)):
        height = Inches(0.4 * rows)
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        
        # Style par défaut
        for cell in table._tbl.iter_tcs():
            tc = cell._tc
            tcPr = tc.get_or_add_tcPr()
            
            # Bordures
            for border in ['left', 'right', 'top', 'bottom']:
                border_element = parse_xml(f'<a:ln w="12700" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"><a:solidFill><a:srgbClr val="808080"/></a:solidFill></a:ln>')
                setattr(tcPr, f'{border}Border', border_element)
        
        return table

    def format_table_header(self, table, row_idx=0):
        for cell in table.rows[row_idx].cells:
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.COLORS['violet_fonce']
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.color.rgb = self.COLORS['blanc']
            paragraph.font.bold = True
            paragraph.alignment = PP_ALIGN.CENTER

    def add_ok_nok(self, cell, is_ok):
        cell.text = "OK" if is_ok else "NOK"
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.color.rgb = self.COLORS['vert'] if is_ok else self.COLORS['rouge']
        paragraph.alignment = PP_ALIGN.CENTER 