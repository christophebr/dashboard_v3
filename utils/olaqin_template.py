from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from datetime import datetime
import os

class OlaqinTemplate:
    # Couleurs Olaqin
    COLORS = {
        'violet_fonce': RGBColor(107, 44, 145),  # #6B2C91
        'violet_clair': RGBColor(178, 140, 200),  # #B28CC8
        'vert': RGBColor(0, 166, 81),  # #00A651
        'rouge': RGBColor(255, 0, 0),  # #FF0000
        'blanc': RGBColor(255, 255, 255),  # #FFFFFF
        'noir': RGBColor(0, 0, 0),  # #000000
        'gris': RGBColor(128, 128, 128)  # #808080
    }

    def __init__(self):
        self.prs = Presentation()
        self._set_slide_dimensions()
        self._add_layouts()
        
        # Vérifier que les ressources existent
        self.assets_path = os.path.join(os.path.dirname(__file__), '..', 'assets')
        if not os.path.exists(self.assets_path):
            os.makedirs(self.assets_path)

    def _set_slide_dimensions(self):
        self.prs.slide_width = Inches(13.33)
        self.prs.slide_height = Inches(7.5)

    def _add_layouts(self):
        # Nous garderons les layouts par défaut pour l'instant
        pass

    def _find_image(self, base_name):
        """Cherche une image dans les formats PNG et JPEG"""
        for ext in ['.png', '.jpeg', '.jpg']:
            path = os.path.join(self.assets_path, f'{base_name}{ext}')
            if os.path.exists(path):
                return path
        return None

    def add_title_slide(self, title, subtitle=None):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[0])
        
        # Ajouter le logo Olaqin en haut à gauche
        logo_path = self._find_image('olaqin_logo')
        if logo_path:
            left = Inches(0.5)
            top = Inches(0.5)
            slide.shapes.add_picture(logo_path, left, top, height=Inches(0.6))
        
        # Titre
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(8), Inches(1))
        title_frame = title_box.text_frame
        title_para = title_frame.add_paragraph()
        title_para.text = title
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        
        # Sous-titre
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(8), Inches(0.75))
            subtitle_frame = subtitle_box.text_frame
            subtitle_para = subtitle_frame.add_paragraph()
            subtitle_para.text = subtitle
            subtitle_para.font.size = Pt(32)
        
        # Logo hexagonal en filigrane
        hex_path = self._find_image('olaqin_hex')
        if hex_path:
            slide.shapes.add_picture(hex_path, Inches(9), Inches(1), height=Inches(5))
        
        self._add_footer(slide)
        return slide

    def add_content_slide(self, title):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        
        # Titre
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(0.75))
        title_frame = title_box.text_frame
        title_para = title_frame.add_paragraph()
        title_para.text = title
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        
        self._add_footer(slide)
        return slide

    def add_bullet_slide(self, title, bullet_points):
        """Ajoute une slide avec des puces"""
        slide = self.add_content_slide(title)
        
        # Ajouter la zone de texte pour les puces
        text_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(11), Inches(5))
        text_frame = text_box.text_frame
        
        for point in bullet_points:
            p = text_frame.add_paragraph()
            p.text = point
            p.level = 0  # Niveau de puce (0 = premier niveau)
            p.font.size = Pt(24)
            
        return slide

    def add_two_column_slide(self, title, left_content=None, right_content=None):
        """Ajoute une slide avec deux colonnes"""
        slide = self.add_content_slide(title)
        
        # Colonne gauche
        if left_content:
            left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.5), Inches(5))
            left_frame = left_box.text_frame
            p = left_frame.add_paragraph()
            p.text = left_content
            p.font.size = Pt(24)
        
        # Colonne droite
        if right_content:
            right_box = slide.shapes.add_textbox(Inches(6.5), Inches(1.5), Inches(5.5), Inches(5))
            right_frame = right_box.text_frame
            p = right_frame.add_paragraph()
            p.text = right_content
            p.font.size = Pt(24)
        
        return slide

    def add_image_slide(self, title, image_path, image_description=None, is_plotly_figure=False, fig=None):
        """Ajoute une slide avec une image centrée. Si is_plotly_figure=True, fig doit être fourni et image_path ignoré."""
        slide = self.add_content_slide(title)
        
        # Ajouter l'image
        if is_plotly_figure and fig is not None:
            import io
            img_bytes = io.BytesIO()
            fig.write_image(img_bytes, format="png")
            img_bytes.seek(0)
            img_width = Inches(10)
            img_height = Inches(4)
            left = (self.prs.slide_width - img_width) / 2
            top = Inches(1.5)
            slide.shapes.add_picture(img_bytes, left, top, width=img_width, height=img_height)
        elif os.path.exists(image_path):
            # Calculer les dimensions pour centrer l'image
            img_width = Inches(10)  # Largeur maximale
            img_height = Inches(4)  # Hauteur maximale
            left = (self.prs.slide_width - img_width) / 2
            top = Inches(1.5)
            slide.shapes.add_picture(image_path, left, top, width=img_width, height=img_height)
        
        # Ajouter la description si fournie
        if image_description:
            desc_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(12), Inches(0.5))
            desc_frame = desc_box.text_frame
            p = desc_frame.add_paragraph()
            p.text = image_description
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(14)
            p.font.italic = True
        
        return slide

    def _add_footer(self, slide):
        # Logo Olaqin
        logo_small_path = self._find_image('olaqin_logo_small')
        if logo_small_path:
            left = Inches(0.5)
            top = Inches(6.5)
            slide.shapes.add_picture(logo_small_path, left, top, height=Inches(0.3))
        
        # Date
        date_box = slide.shapes.add_textbox(Inches(2), Inches(6.5), Inches(2), Inches(0.3))
        date_frame = date_box.text_frame
        date_para = date_frame.add_paragraph()
        date_para.text = datetime.now().strftime("%d/%m/%Y")
        date_para.font.size = Pt(10)
        
        # Mention confidentielle
        conf_box = slide.shapes.add_textbox(Inches(4), Inches(6.5), Inches(5), Inches(0.3))
        conf_frame = conf_box.text_frame
        conf_para = conf_frame.add_paragraph()
        conf_para.text = "CONFIDENTIEL – NE PAS DIFFUSER ©Olaqin"
        conf_para.font.size = Pt(10)
        conf_para.font.color.rgb = self.COLORS['violet_clair']
        conf_para.alignment = PP_ALIGN.CENTER
        
        # Logo hexagonal
        hex_small_path = self._find_image('olaqin_hex_small')
        if hex_small_path:
            slide.shapes.add_picture(hex_small_path, Inches(12), Inches(6.5), height=Inches(0.3))

    def create_table(self, slide, rows, cols, left=Inches(0.5), top=Inches(1.5), width=Inches(12)):
        height = Inches(0.4 * rows)
        shape = slide.shapes.add_table(rows, cols, left, top, width, height)
        table = shape.table
        
        # Définir les largeurs de colonnes
        col_widths = [2.5, 0.9, 0.9, 0.9, 0.9, 0.9, 0.9, 1.1]  # Total = 9 pouces
        for i, width in enumerate(col_widths):
            if i < len(table.columns):
                table.columns[i].width = Inches(width)
        
        return table

    def format_table_header(self, table, row_idx=0):
        for cell in table.rows[row_idx].cells:
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.COLORS['violet_fonce']
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.color.rgb = self.COLORS['blanc']
            paragraph.font.bold = True
            paragraph.font.size = Pt(9)
            paragraph.alignment = PP_ALIGN.CENTER

    def add_ok_nok(self, cell, is_ok):
        cell.text = "OK" if is_ok else "NOK"
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.color.rgb = self.COLORS['vert'] if is_ok else self.COLORS['rouge']
        paragraph.font.size = Pt(9)
        paragraph.alignment = PP_ALIGN.CENTER

    def format_cell(self, cell, text, bold=False, color=None, alignment=PP_ALIGN.LEFT):
        """Formate une cellule de tableau avec des options personnalisées"""
        cell.text = str(text)
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.size = Pt(9)
        paragraph.font.bold = bold
        if color in self.COLORS:
            paragraph.font.color.rgb = self.COLORS[color]
        paragraph.alignment = alignment 