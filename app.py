import os
# En déploiement (sans hashed_pw.pkl), toujours utiliser config.example pour prendre les secrets à jour
if not os.path.exists('hashed_pw.pkl') and os.path.exists('config.example.py'):
    import shutil
    shutil.copy('config.example.py', 'config.py')

import streamlit as st

# set_page_config DOIT être la première commande Streamlit
st.set_page_config(
    page_title=":bar_chart: Dashboard support",
    layout="wide",
    initial_sidebar_state="collapsed",
)

from datetime import datetime, date, timedelta
from data_processing.aircall_processing import process_aircall_data, def_df_support, agents_all, line_tous, agents_support, line_support, line_armatis, agents_armatis, load_aircall_data
#from data_processing.aircall_processing import get_df_support
from data_processing.hubspot_processing import process_hubspot_data, load_hubspot_data
from data_processing.kpi_generation import generate_kpis, filtrer_par_periode, calculate_ticket_response_time, graph_activite, evo_appels_ticket, graph_charge_affid_stellair, graph_repartition_groupes_stellair, graph_yelda_evaluation, graph_yelda_evaluation_intentions, graph_yelda_interactions_tickets_semaine, graph_yelda_score_llm, graph_yelda_evolution_scores
from utils.streamlit_helpers import load_data
from utils.powerpoint_helpers import create_powerpoint, create_powerpoint_agents, create_powerpoint_stellair_minimal, create_powerpoint_agents_report
import config
from config import CREDENTIALS, AIRCALL_DATA_PATH_V1, AIRCALL_DATA_PATH_V2, AIRCALL_DATA_PATH_V3, HUBSPOT_TICKET_DATA_PATH, EVALUATION_DATA_PATH, YELDA_DATA_PATH
import streamlit_authenticator as stauth
from pptx import Presentation
from pptx.util import Inches, Pt
import io
import tempfile
import pandas as pd

import streamlit as st
#from selenium import webdriver
#from selenium.webdriver.chrome.service import Service
#from webdriver_manager.chrome import ChromeDriverManager


def get_periode_str(periode):
    """Retourne une chaîne pour affichage et noms de fichier (PowerPoint, CSV, etc.)."""
    if isinstance(periode, (tuple, list)) and len(periode) == 2:
        d1, d2 = pd.to_datetime(periode[0]).strftime('%Y-%m-%d'), pd.to_datetime(periode[1]).strftime('%Y-%m-%d')
        return f"Du {d1} au {d2}"
    return str(periode)


def create_single_graph_powerpoint(fig, title, periode):
    """
    Crée une présentation PowerPoint avec un seul graphique.
    
    Args:
        fig: Figure Plotly à exporter
        title: Titre du graphique
        periode: Période analysée
    
    Returns:
        BytesIO: Présentation PowerPoint en mémoire
    """
    from pptx import Presentation
    from pptx.util import Inches
    import io
    import tempfile
    import os
    
    # Créer une nouvelle présentation
    prs = Presentation()
    
    # Slide 1 - Titre
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide_title = slide.shapes.title
    subtitle = slide.placeholders[1]
    slide_title.text = title
    subtitle.text = f"Période : {periode}"
    
    # Slide 2 - Graphique
    img_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(img_slide_layout)
    slide_title = slide.shapes.title
    slide_title.text = title
    
    # Sauvegarder le graphique temporairement
    tmp_file = None
    try:
        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
            tmp_file = tmp.name
            # Utiliser kaleido pour exporter l'image
            fig.write_image(tmp_file, width=1200, height=800, engine='kaleido')
            slide.shapes.add_picture(tmp_file, Inches(0.5), Inches(1.5), width=Inches(9), height=Inches(6))
    finally:
        # Nettoyer le fichier temporaire
        if tmp_file and os.path.exists(tmp_file):
            try:
                os.unlink(tmp_file)
            except:
                pass
    
    # Sauvegarder la présentation en mémoire
    pptx_io = io.BytesIO()
    prs.save(pptx_io)
    pptx_io.seek(0)
    
    return pptx_io


def render_powerpoint_download_button(fig, title, periode, file_name, button_key):
    """
    Affiche un bouton pour exporter un graphique en PowerPoint.
    Le bouton de téléchargement reste visible après génération.
    Régénère automatiquement si la période change.
    
    Args:
        fig: Figure Plotly à exporter
        title: Titre du graphique
        periode: Période analysée
        file_name: Nom du fichier PowerPoint
        button_key: Clé unique pour le bouton Streamlit
    """
    download_key = f"download_{button_key}"
    periode_key = f"periode_{button_key}"
    
    # Vérifier si la période a changé - si oui, invalider le cache
    if periode_key in st.session_state:
        if st.session_state[periode_key] != periode:
            # La période a changé, invalider le cache
            st.session_state[download_key] = None
    
    # Stocker la période actuelle
    st.session_state[periode_key] = periode
    
    # Initialiser session_state si nécessaire
    if download_key not in st.session_state:
        st.session_state[download_key] = None
    
    # Bouton pour générer le PowerPoint
    if st.button("📊 Exporter en PowerPoint", key=button_key):
        try:
            with st.spinner("Génération du PowerPoint en cours..."):
                # Toujours régénérer avec la figure actuelle (qui contient les données filtrées)
                pptx_io = create_single_graph_powerpoint(fig, title, periode)
                st.session_state[download_key] = pptx_io.getvalue()
            st.success("✅ PowerPoint généré avec succès !")
        except Exception as e:
            st.error(f"Erreur : {e}")
            st.session_state[download_key] = None
    
    # Afficher le bouton de téléchargement si le fichier est prêt
    if st.session_state[download_key] is not None:
        st.download_button(
            label="📥 Télécharger",
            data=st.session_state[download_key],
            file_name=file_name,
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            key=f"dl_{button_key}"
        )


# Authentification de l'utilisateur
authenticator = stauth.Authenticate(
    config.CREDENTIALS,
    'dashboard_support',
    config.COOKIE_KEY,  # Assurez-vous que cette clé est unique et sécurisée
    cookie_expiry_days=2
)

# Utilisation du nouvel argument 'fields'
name, authentification_status, username = authenticator.login('main', fields={'Form name': 'custom_form_name'})

if authentification_status == False:
    st.error('Nom d\'utilisateur/mot de passe incorrect')
elif authentification_status == None:
    st.warning('Veuillez entrer votre nom d\'utilisateur et mot de passe')
elif authentification_status:
    st.sidebar.title(f"Bienvenue, {name}")
    authenticator.logout('Logout', 'sidebar')

    # Bouton dans la sidebar pour recharger les données Aircall et HubSpot
    if st.sidebar.button("Recharger les données"):
        import os
        import sqlite3
        
        # Vider tous les caches Streamlit
        cache_keys_to_clear = [
            'df_aircall', 'df_tickets', 'df_evaluation',
            'df_support_processed', 'df_tickets_processed', 'kpis_all',
            'df_aircall_processed'
        ]
        
        # Ajouter les caches des dataframes support
        for option in ["support_stellair", "support_affid", "xmed", "tmaj"]:
            cache_keys_to_clear.append(f"df_support_{option}")
        
        for key in cache_keys_to_clear:
            if key in st.session_state:
                del st.session_state[key]
        
        # Vider le cache Streamlit (@st.cache_data)
        st.cache_data.clear()
        
        # Supprimer les fichiers de cache pour forcer le rechargement complet
        cache_files_to_delete = [
            "data/Cache/data_cache.db",  # Cache SQLite pour HubSpot
            "data/Affid/Cache/df_tickets.parquet",  # Cache Parquet pour HubSpot
            "data/Affid/Cache/df_support.parquet",  # Cache Parquet pour Aircall
            "data/cache.sqlite",  # Cache SQLite pour Aircall
        ]
        
        for cache_file in cache_files_to_delete:
            if os.path.exists(cache_file):
                try:
                    os.remove(cache_file)
                    print(f"✅ Fichier de cache supprimé: {cache_file}")
                except Exception as e:
                    print(f"⚠️ Impossible de supprimer {cache_file}: {e}")
        
        # Recharger les données
        st.session_state['df_aircall'] = load_aircall_data(AIRCALL_DATA_PATH_V1, AIRCALL_DATA_PATH_V2, AIRCALL_DATA_PATH_V3, force_reload=True)
        st.session_state['df_tickets'] = load_hubspot_data(HUBSPOT_TICKET_DATA_PATH)
        st.session_state['df_evaluation'] = pd.read_excel(EVALUATION_DATA_PATH)
        st.success("✅ Données Aircall, HubSpot et Evaluation rechargées ! Les fichiers de cache ont été supprimés.")

    # Chargement initial si pas déjà en mémoire
    if 'df_aircall' not in st.session_state:
        st.session_state['df_aircall'] = load_aircall_data(AIRCALL_DATA_PATH_V1, AIRCALL_DATA_PATH_V2, AIRCALL_DATA_PATH_V3, force_reload=False)
    if 'df_tickets' not in st.session_state:
        st.session_state['df_tickets'] = load_hubspot_data(HUBSPOT_TICKET_DATA_PATH)
    if 'df_evaluation' not in st.session_state:
        st.session_state['df_evaluation'] = pd.read_excel(EVALUATION_DATA_PATH)

    df_aircall = st.session_state['df_aircall']
    df_tickets = st.session_state['df_tickets']
    df_evaluation = st.session_state['df_evaluation']

    # Optimisation : traitement des données une seule fois et mise en cache
    if 'df_support_processed' not in st.session_state:
        st.session_state['df_support_processed'] = def_df_support(
            process_aircall_data(df_aircall), 
            process_aircall_data(df_aircall), 
            line_tous, 
            agents_all
        )
    if 'df_tickets_processed' not in st.session_state:
        st.session_state['df_tickets_processed'] = process_hubspot_data(df_tickets)

    df_support = st.session_state['df_support_processed']
    df_tickets_processed = st.session_state['df_tickets_processed']

    # Génération des KPIs (mise en cache si possible)
    if 'kpis_all' not in st.session_state:
        st.session_state['kpis_all'] = generate_kpis(df_support, df_tickets_processed, 'agents_all')
    
    kpis = st.session_state['kpis_all']

    PAGES = {
        "Support": "support",
        "Agents": "agents",
        "Tickets": "tickets",
        "Analyste IA": "mcp_analyst"
    }


    selection_page = st.sidebar.selectbox("Choix de la page", list(PAGES.keys()), key="unique_page_selection")

    if selection_page == "Support":
        from data_processing.aircall_processing import process_aircall_data, def_df_support, agents_all, line_tous, agents_support, line_support, line_armatis, agents_armatis
        from data_processing.kpi_generation import (
            filtrer_par_periode, generate_kpis, convert_to_sixtieth, 
            graph_activite, graph_taux_jour, graph_taux_heure, 
            graph_activite_xmed, graph_activite_tmaj, calculate_ticket_response_time, 
            graph_charge_affid_stellair, graph_repartition_groupes_stellair,
            categoriser_avec_ia_personnalisee, evo_tickets_par_sujets_mensuel, repartition_lecteurs_par_type
        )


        def support():

            dataframe_option = st.sidebar.selectbox("Choisir le dataframe", ["support_stellair", "support_affid", "xmed", "tmaj"], key="unique_dataframe_selection")

            # Optimisation : traitement des données une seule fois et mise en cache
            if 'df_aircall_processed' not in st.session_state:
                st.session_state['df_aircall_processed'] = process_aircall_data(df_aircall)

            df_aircall_processed = st.session_state['df_aircall_processed']

            # Cache pour les différents dataframes
            cache_key = f"df_support_{dataframe_option}"
            if cache_key not in st.session_state:
                if dataframe_option == "support_stellair":
                    df_stellair = def_df_support(df_aircall_processed, df_aircall_processed, line_tous, agents_all)
                    df_stellair = df_stellair[(df_stellair['line'] == 'armatistechnique') | (df_stellair['Logiciel'] == 'Stellair')]
                    st.session_state[cache_key] = df_stellair
                elif dataframe_option == "support_affid":
                    df_affid = def_df_support(df_aircall_processed, df_aircall_processed, line_tous, agents_all)
                    df_affid = df_affid[(df_affid['Logiciel'] == 'Affid')]    
                    st.session_state[cache_key] = df_affid
                elif dataframe_option == "xmed":
                    df_xmed = def_df_support(df_aircall_processed, df_aircall_processed, line_tous, agents_all)
                    df_xmed = df_xmed[(df_xmed['line'] == 'xmed')]
                    st.session_state[cache_key] = df_xmed
                elif dataframe_option == "tmaj":
                    df_tmaj = def_df_support(df_aircall_processed, df_aircall_processed, line_tous, agents_all)
                    df_tmaj = df_tmaj[(df_tmaj['line'] == 'supporthardware')]
                    st.session_state[cache_key] = df_tmaj

            # Configuration des agents selon l'option
            agents_config = {
                "support_stellair": agents_all,
                "support_affid": agents_support,
                "xmed": agents_support,
                "tmaj": agents_support
            }

            df_support = st.session_state[cache_key]
            agents = agents_config[dataframe_option]

            periode_option = st.selectbox("Sélectionnez une période :", 
                    ["1 an", "6 derniers mois", "3 derniers mois", "Dernier mois", "Personnalisée"],
                    key="periode_support")
            
            if periode_option == "Personnalisée":
                # Bornes min/max des données pour les date pickers
                df_temp = df_support.copy()
                df_temp['Date'] = pd.to_datetime(df_temp['Date'])
                date_min_data = (df_temp['Date'].min() or pd.Timestamp.now() - pd.DateOffset(years=2)).date()
                date_max_data = (df_temp['Date'].max() or pd.Timestamp.now()).date()
                col_d1, col_d2 = st.columns(2)
                with col_d1:
                    date_debut = st.date_input("Date de début", value=date_max_data - timedelta(days=90), min_value=date_min_data, max_value=date_max_data, key="date_debut_support")
                with col_d2:
                    date_fin = st.date_input("Date de fin", value=date_max_data, min_value=date_debut, max_value=date_max_data, key="date_fin_support")
                periode = (date_debut, date_fin)
            else:
                periode = periode_option
            
            periode_str = get_periode_str(periode)
            df_support = filtrer_par_periode(df_support, periode)
            kpis = generate_kpis(filtrer_par_periode(df_support, periode), filtrer_par_periode(df_tickets_processed, periode), agents)

            # Utiliser graph_activite_xmed pour la page xmed, graph_activite_tmaj pour TMAJ, sinon graph_activite normal
            if dataframe_option == "xmed":
                col1, col2, col3 = st.columns(3)
                col1.metric("Taux de service en %", kpis['Taux_de_service'])
                col2.metric("Appels entrant / Jour", kpis['Entrant'])
                col3.metric("Numéros uniques / Jour", kpis['Numero_unique'])
                st.plotly_chart(graph_activite_xmed(df_support), use_container_width=True)
                col1, col2 = st.columns([3, 1])
                with col2:
                    render_powerpoint_download_button(
                        graph_activite_xmed(df_support),
                        "Activité XMED",
                        periode_str,
                        f"activite_xmed_{periode_str.replace(' ', '_').replace('/', '-')}.pptx",
                        "btn_activite_xmed"
                    )
            elif dataframe_option == "tmaj":
                col1, col2, col3 = st.columns(3)
                col1.metric("Taux de service en %", kpis['Taux_de_service'])
                col2.metric("Appels entrant / Jour", kpis['Entrant'])
                col3.metric("Numéros uniques / Jour", kpis['Numero_unique'])
                st.plotly_chart(graph_activite_tmaj(df_support), use_container_width=True)
                col1, col2 = st.columns([3, 1])
                with col2:
                    render_powerpoint_download_button(
                        graph_activite_tmaj(df_support),
                        "Activité TMAJ (Support Hardware)",
                        periode_str,
                        f"activite_tmaj_{periode_str.replace(' ', '_').replace('/', '-')}.pptx",
                        "btn_activite_tmaj"
                    )
            elif dataframe_option == "support_stellair":
                # KPI spécifiques pour support_stellair
                col1, col2, col3 = st.columns(3)
                col1.metric("Taux de service en %", kpis['Taux_de_service'])
                col2.metric("Appels entrant / Jour", kpis['Entrant'])
                col3.metric("Numéros uniques / Jour", kpis['Numero_unique'])

                col4, col5 = st.columns(2)
                col4.metric(
                    "Entrants vs Tickets (%)",
                    f"{round(kpis['activite_appels_pourcentage'] * 100, 2)}% / {round(kpis['activite_tickets_pourcentage'] * 100, 2)}%"
                )
                # Calcul du temps de réponse aux tickets
                from data_processing.kpi_generation import calculate_ticket_response_time
                moyenne_temps_reponse, graph_temps_reponse, df_temps_tickets = calculate_ticket_response_time(
                    filtrer_par_periode(df_tickets_processed, periode), 
                    agents_all
                )

                col5.metric("Temps de réponse moyen aux tickets (h:min)", f"{int(moyenne_temps_reponse)}:{int((moyenne_temps_reponse % 1) * 60):02d}")

                # Graphique principal
                st.plotly_chart(graph_activite(df_support), use_container_width=True)
                col1, col2 = st.columns([3, 1])
                with col2:
                    render_powerpoint_download_button(
                        graph_activite(df_support),
                        "Activité Support Stellair",
                        periode_str,
                        f"activite_stellair_{periode_str.replace(' ', '_').replace('/', '-')}.pptx",
                        "btn_activite"
                    )
                
                # Graphique Activité en % - NXT vs Stellair (nécessite les données complètes Affid + Stellair)
                st.markdown("### 📊 Répartition des appels entre Affid (NXT) et Stellair")
                fig_charge_pct_stellair, fig_charge_vol_stellair, _ = graph_charge_affid_stellair(
                    filtrer_par_periode(st.session_state['df_support_processed'], periode)
                )
                st.plotly_chart(fig_charge_pct_stellair, use_container_width=True)
                col1, col2 = st.columns([3, 1])
                with col2:
                    if st.button("📊 Exporter Activité % NXT/Stellair en PowerPoint", key="btn_charge_pct_stellair"):
                        try:
                            pptx_io = create_single_graph_powerpoint(fig_charge_pct_stellair, "Activité en % - NXT vs Stellair", periode_str)
                            st.download_button(
                                label="📥 Télécharger",
                                data=pptx_io,
                                file_name=f"activite_pct_nxt_stellair_{periode_str.replace(' ', '_').replace('/', '-')}.pptx",
                                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                            )
                        except Exception as e:
                            st.error(f"Erreur : {e}")
                
                # Graphique des temps de réponse aux tickets
                st.plotly_chart(graph_temps_reponse, use_container_width=True)
                col1, col2 = st.columns([3, 1])
                with col2:
                    render_powerpoint_download_button(
                        graph_temps_reponse,
                        "Temps de réponse aux tickets",
                        periode_str,
                        f"temps_reponse_stellair_{periode_str.replace(' ', '_').replace('/', '-')}.pptx",
                        "btn_temps_reponse"
                    )
                
                # Affichage des graphiques de taux
                col_graph1, col_graph2 = st.columns(2)
                col_graph1.plotly_chart(graph_taux_jour(df_support), use_container_width=True)
                col1, col2 = st.columns([3, 1])
                with col2:
                    if st.button("📊 Exporter en PowerPoint", key="btn_taux_jour"):
                        try:
                            pptx_io = create_single_graph_powerpoint(graph_taux_jour(df_support), "Taux par jour", periode_str)
                            st.download_button(
                                label="📥 Télécharger",
                                data=pptx_io,
                                file_name=f"taux_jour_stellair_{periode_str.replace(' ', '_').replace('/', '-')}.pptx",
                                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                            )
                        except Exception as e:
                            st.error(f"Erreur : {e}")
                
                col_graph2.plotly_chart(graph_taux_heure(df_support), use_container_width=True)
                col1, col2 = st.columns([3, 1])
                with col2:
                    if st.button("📊 Exporter en PowerPoint", key="btn_taux_heure"):
                        try:
                            pptx_io = create_single_graph_powerpoint(graph_taux_heure(df_support), "Taux par heure", periode_str)
                            st.download_button(
                                label="📥 Télécharger",
                                data=pptx_io,
                                file_name=f"taux_heure_stellair_{periode_str.replace(' ', '_').replace('/', '-')}.pptx",
                                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                            )
                        except Exception as e:
                            st.error(f"Erreur : {e}")
                
                # ----------- GRAPHIQUE N1 PAR SEMAINE (SSI/SSIA/SPSA/Chatbot Yelda - agents support) -----------
                st.markdown("## 📈 Graphique N1 par semaine (SSI/SSIA/SPSA/Chatbot Yelda - agents support)")
                from data_processing.kpi_generation import graph_tickets_n1_par_semaine_stellair
                df_tickets_stellair_periode = filtrer_par_periode(df_tickets_processed, periode)
                weeks_to_display_n1 = sorted(
                    [w for w in df_tickets_stellair_periode['Semaine'].unique() if isinstance(w, str) and len(w) >= 8],
                    key=lambda x: (int(x[1:5]), int(x[6:]))
                )
                cache_key_n1_semaine_stellair = f'n1_semaine_stellair_{periode_str}_{hash(str(df_tickets_processed.shape))}_{hash(str(weeks_to_display_n1[:5]))}'
                if cache_key_n1_semaine_stellair not in st.session_state:
                    with st.spinner("Génération du graphique N1 par semaine..."):
                        fig_n1_semaine, agents_disponibles_n1, tickets_n1_en_cours = graph_tickets_n1_par_semaine_stellair(df_tickets_processed, weeks_to_display=weeks_to_display_n1)
                        st.session_state[cache_key_n1_semaine_stellair] = {'fig_n1_semaine': fig_n1_semaine, 'agents_disponibles': agents_disponibles_n1, 'tickets_n1_en_cours': tickets_n1_en_cours}
                else:
                    fig_n1_semaine = st.session_state[cache_key_n1_semaine_stellair]['fig_n1_semaine']
                    agents_disponibles_n1 = st.session_state[cache_key_n1_semaine_stellair]['agents_disponibles']
                    tickets_n1_en_cours = st.session_state[cache_key_n1_semaine_stellair]['tickets_n1_en_cours']
                agent_selection_n1 = st.selectbox("Filtrer par agent :", ["Tous les agents"] + (agents_disponibles_n1 or []), key="agent_filter_n1_semaine_stellair")
                if agent_selection_n1 != "Tous les agents":
                    cache_key_n1_filtered = f'n1_semaine_stellair_{agent_selection_n1}_{periode_str}_{hash(str(weeks_to_display_n1[:5]))}'
                    if cache_key_n1_filtered not in st.session_state:
                        fig_n1_semaine, _, tickets_n1_en_cours = graph_tickets_n1_par_semaine_stellair(df_tickets_processed, agent_selection_n1, weeks_to_display=weeks_to_display_n1)
                        st.session_state[cache_key_n1_filtered] = {'fig_n1_semaine': fig_n1_semaine, 'tickets_n1_en_cours': tickets_n1_en_cours}
                    else:
                        fig_n1_semaine = st.session_state[cache_key_n1_filtered]['fig_n1_semaine']
                        tickets_n1_en_cours = st.session_state[cache_key_n1_filtered]['tickets_n1_en_cours']
                st.plotly_chart(fig_n1_semaine, use_container_width=True)
                col1, col2 = st.columns([3, 1])
                with col2:
                    if st.button("📊 Exporter N1 par semaine en PowerPoint", key="btn_n1_semaine_stellair"):
                        try:
                            title = f"N1 par semaine - {agent_selection_n1}" if agent_selection_n1 != "Tous les agents" else "N1 par semaine"
                            pptx_io = create_single_graph_powerpoint(fig_n1_semaine, title, periode_str)
                            st.download_button(
                                label="📥 Télécharger",
                                data=pptx_io,
                                file_name=f"n1_par_semaine_{agent_selection_n1.replace(' ', '_')}_{periode_str.replace(' ', '_').replace('/', '-')}.pptx",
                                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                            )
                        except Exception as e:
                            st.error(f"Erreur : {e}")
                if not tickets_n1_en_cours.empty:
                    st.markdown(f"**🔴 {len(tickets_n1_en_cours)} ticket(s) N1 en cours**")
                    if st.button("📥 Télécharger la liste des tickets N1 en cours", key="btn_dl_n1_stellair"):
                        csv_n1 = tickets_n1_en_cours[['Ticket ID', 'Date', 'Statut du ticket', 'Propriétaire du ticket']].to_csv(index=False, encoding='utf-8-sig')
                        st.download_button("📄 Télécharger CSV", data=csv_n1, file_name=f"tickets_n1_en_cours_{periode_str.replace(' ', '_').replace('/', '-')}.csv", mime="text/csv")
                
                # ----------- INDICATEURS YELDA (Chatbot Stellair - fse.stellair.fr) -----------
                st.markdown("## 🤖 Indicateurs Yelda (Chatbot Stellair)")
                try:
                    from data_processing.yelda_processing import load_yelda_data, filter_yelda_stellair, filter_yelda_evaluated, filtrer_yelda_par_periode, compute_yelda_kpis
                    df_yelda_raw = load_yelda_data(YELDA_DATA_PATH)
                    if df_yelda_raw is not None and not df_yelda_raw.empty:
                        df_yelda_fse = filter_yelda_stellair(df_yelda_raw)
                        df_yelda_periode = filtrer_yelda_par_periode(df_yelda_fse, periode)
                        df_yelda_eval = filter_yelda_evaluated(df_yelda_periode)
                        kpis_yelda = compute_yelda_kpis(df_yelda_eval)
                        nb_satisfait = kpis_yelda['evaluation_counts']['Satisfait']
                        nb_insatisfait = kpis_yelda['evaluation_counts']['Insatisfait']
                        nb_revoir = kpis_yelda['evaluation_counts']['À revoir']
                        total_eval = nb_satisfait + nb_insatisfait + nb_revoir
                        taux_satisfaction = round(100 * nb_satisfait / total_eval, 1) if total_eval > 0 else 0
                        col1, col2, col3, col4, col5 = st.columns(5)
                        col1.metric("Conversations évaluées (fse.stellair.fr)", kpis_yelda['nb_interactions'])
                        col2.metric("Tickets créés", kpis_yelda['nb_tickets_crees'])
                        col3.metric("✅ Satisfait (éval. LLM)", nb_satisfait)
                        col4.metric("❌ Insatisfait (éval. LLM)", nb_insatisfait)
                        col5.metric("Taux satisfaction*", f"{taux_satisfaction}%")
                        nb_intent_sat = kpis_yelda.get('intentions_satisfaisant', 0)
                        nb_intent_non = kpis_yelda.get('intentions_non_satisfaisant', 0)
                        tot_intent = nb_intent_sat + nb_intent_non
                        taux_intent = round(100 * nb_intent_sat / tot_intent, 1) if tot_intent > 0 else 0
                        st.markdown("**Évaluation utilisateurs (Intentions)** — reponse_agent_satisfaisante / non_satisfaisante")
                        col_a, col_b, col_c = st.columns(3)
                        col_a.metric("👍 Satisfaisant", nb_intent_sat)
                        col_b.metric("👎 Non satisfaisant", nb_intent_non)
                        col_c.metric("Taux satisfaction utilisateurs", f"{taux_intent}%")
                        st.caption("*Taux satisfaction = Satisfait / (Satisfait + Insatisfait + À revoir) — Seules les conversations évaluées sont comptabilisées (ouvertures sans question exclues)")
                        st.plotly_chart(graph_yelda_evaluation(kpis_yelda['evaluation_counts']), use_container_width=True)
                        st.plotly_chart(graph_yelda_evaluation_intentions(nb_intent_sat, nb_intent_non), use_container_width=True)
                        st.plotly_chart(graph_yelda_interactions_tickets_semaine(df_yelda_eval), use_container_width=True)
                        st.plotly_chart(graph_yelda_evolution_scores(df_yelda_eval), use_container_width=True)
                        if kpis_yelda['score_llm_moyen'] > 0:
                            st.plotly_chart(graph_yelda_score_llm(df_yelda_eval), use_container_width=True)
                    else:
                        st.info("📁 Fichier Yelda non trouvé. Placez `yelda.xlsx` dans `data/Affid/yelda/` pour afficher les indicateurs.")
                except Exception as e:
                    st.warning(f"Indicateurs Yelda non disponibles : {e}")
                
                # Nouveau graphique de répartition des groupes d'agents
                st.plotly_chart(graph_repartition_groupes_stellair(df_support), use_container_width=True)
                col1, col2 = st.columns([3, 1])
                with col2:
                    if st.button("📊 Exporter en PowerPoint", key="btn_repartition_groupes"):
                        try:
                            pptx_io = create_single_graph_powerpoint(graph_repartition_groupes_stellair(df_support), "Répartition des groupes d'agents", periode_str)
                            st.download_button(
                                label="📥 Télécharger",
                                data=pptx_io,
                                file_name=f"repartition_groupes_stellair_{periode_str.replace(' ', '_').replace('/', '-')}.pptx",
                                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                            )
                        except Exception as e:
                            st.error(f"Erreur : {e}")
                
                st.plotly_chart(kpis['evo_appels_tickets'])
                col1, col2 = st.columns([3, 1])
                with col2:
                    if st.button("📊 Exporter évolution hebdomadaire en PowerPoint", key="btn_evo_appels_tickets_hebdo"):
                        try:
                            pptx_io = create_single_graph_powerpoint(kpis['evo_appels_tickets'], "Évolution hebdomadaire : appels entrants + tickets", periode_str)
                            st.download_button(
                                label="📥 Télécharger",
                                data=pptx_io,
                                file_name=f"evolution_hebdomadaire_appels_tickets_{periode_str.replace(' ', '_').replace('/', '-')}.pptx",
                                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                            )
                        except Exception as e:
                            st.error(f"Erreur : {e}")
                
                st.plotly_chart(kpis.get('evo_appels_tickets_mensuel'))
                col1, col2 = st.columns([3, 1])
                with col2:
                    if st.button("📊 Exporter évolution mensuelle en PowerPoint", key="btn_evo_appels_tickets_mensuel"):
                        try:
                            pptx_io = create_single_graph_powerpoint(kpis.get('evo_appels_tickets_mensuel'), "Évolution mensuelle : appels entrants + tickets", periode_str)
                            st.download_button(
                                label="📥 Télécharger",
                                data=pptx_io,
                                file_name=f"evolution_mensuelle_appels_tickets_{periode_str.replace(' ', '_').replace('/', '-')}.pptx",
                                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                            )
                        except Exception as e:
                            st.error(f"Erreur : {e}")
                
                st.plotly_chart(kpis.get('sunburst_categories'))
                col1, col2 = st.columns([3, 1])
                with col2:
                    if st.button("📊 Exporter Sunburst en PowerPoint", key="btn_sunburst_categories"):
                        try:
                            pptx_io = create_single_graph_powerpoint(kpis.get('sunburst_categories'), "Répartition des tickets SSI par sujets et catégories", periode_str)
                            st.download_button(
                                label="📥 Télécharger",
                                data=pptx_io,
                                file_name=f"sunburst_categories_ssi_{periode_str.replace(' ', '_').replace('/', '-')}.pptx",
                                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                            )
                        except Exception as e:
                            st.error(f"Erreur : {e}")
                
                # Nouveaux graphiques d'évolution mensuelle des tickets
                fig_total, fig_sujets = evo_tickets_par_sujets_mensuel(filtrer_par_periode(df_tickets_processed, periode))
                
                # Graphique du total mensuel
                st.plotly_chart(fig_total, use_container_width=True)
                col1, col2 = st.columns([3, 1])
                with col2:
                    if st.button("📊 Exporter total mensuel en PowerPoint", key="btn_total_mensuel"):
                        try:
                            pptx_io = create_single_graph_powerpoint(fig_total, "Évolution mensuelle du nombre total de tickets SSI", periode_str)
                            st.download_button(
                                label="📥 Télécharger",
                                data=pptx_io,
                                file_name=f"total_mensuel_{periode_str.replace(' ', '_').replace('/', '-')}.pptx",
                                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                            )
                        except Exception as e:
                            st.error(f"Erreur : {e}")
                
                # Graphique des sujets par mois
                st.plotly_chart(fig_sujets, use_container_width=True)
                col1, col2 = st.columns([3, 1])
                with col2:
                    if st.button("📊 Exporter évolution sujets en PowerPoint", key="btn_evo_sujets_mensuel"):
                        try:
                            pptx_io = create_single_graph_powerpoint(fig_sujets, "Évolution mensuelle des tickets SSI par sujets", periode_str)
                            st.download_button(
                                label="📥 Télécharger",
                                data=pptx_io,
                                file_name=f"evolution_sujets_mensuel_{periode_str.replace(' ', '_').replace('/', '-')}.pptx",
                                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                            )
                        except Exception as e:
                            st.error(f"Erreur : {e}")
                
                # Graphique camembert des types de lecteur
                st.subheader("📊 Répartition des tickets Lecteur par type de lecteur")
                try:
                    fig_lecteurs = repartition_lecteurs_par_type(filtrer_par_periode(df_tickets_processed, periode))
                    st.plotly_chart(fig_lecteurs, use_container_width=True)
                    
                    col1, col2 = st.columns([3, 1])
                    with col2:
                        if st.button("📊 Exporter répartition lecteurs en PowerPoint", key="btn_repartition_lecteurs"):
                            try:
                                pptx_io = create_single_graph_powerpoint(fig_lecteurs, "Répartition des tickets Lecteur par type de lecteur", periode_str)
                                st.download_button(
                                    label="📥 Télécharger",
                                    data=pptx_io,
                                    file_name=f"repartition_lecteurs_{periode_str.replace(' ', '_').replace('/', '-')}.pptx",
                                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                                )
                            except Exception as e:
                                st.error(f"Erreur : {e}")
                except Exception as e:
                    st.error(f"Erreur lors de la génération du graphique des lecteurs : {e}")
            elif dataframe_option == "support_affid":
                col1, col2, col3 = st.columns(3)
                col1.metric("Taux de service en %", kpis['Taux_de_service'])
                col2.metric("Appels entrant / Jour", kpis['Entrant'])
                col3.metric("Numéros uniques / Jour", kpis['Numero_unique'])

                col4, col5 = st.columns(2)
                col4.metric(
                    "Entrants vs Tickets (%)",
                    f"{round(kpis['activite_appels_pourcentage'] * 100, 2)}% / {round(kpis['activite_tickets_pourcentage'] * 100, 2)}%"
                )
                col5.empty()

                # Affichage du graphique d'activité
                st.plotly_chart(graph_activite(df_support), use_container_width=True)
                col1, col2 = st.columns([3, 1])
                with col2:
                    if st.button("📊 Exporter en PowerPoint", key="btn_activite_affid"):
                        try:
                            pptx_io = create_single_graph_powerpoint(graph_activite(df_support), "Activité Support Affid", periode_str)
                            st.download_button(
                                label="📥 Télécharger",
                                data=pptx_io,
                                file_name=f"activite_affid_{periode_str.replace(' ', '_').replace('/', '-')}.pptx",
                                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                            )
                        except Exception as e:
                            st.error(f"Erreur : {e}")
                
                # Affichage des graphiques de taux
                col_graph1, col_graph2 = st.columns(2)
                col_graph1.plotly_chart(graph_taux_jour(df_support), use_container_width=True)
                col1, col2 = st.columns([3, 1])
                with col2:
                    if st.button("📊 Exporter en PowerPoint", key="btn_taux_jour_affid"):
                        try:
                            pptx_io = create_single_graph_powerpoint(graph_taux_jour(df_support), "Taux par jour Affid", periode_str)
                            st.download_button(
                                label="📥 Télécharger",
                                data=pptx_io,
                                file_name=f"taux_jour_affid_{periode_str.replace(' ', '_').replace('/', '-')}.pptx",
                                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                            )
                        except Exception as e:
                            st.error(f"Erreur : {e}")
                
                col_graph2.plotly_chart(graph_taux_heure(df_support), use_container_width=True)
                col1, col2 = st.columns([3, 1])
                with col2:
                    if st.button("📊 Exporter en PowerPoint", key="btn_taux_heure_affid"):
                        try:
                            pptx_io = create_single_graph_powerpoint(graph_taux_heure(df_support), "Taux par heure Affid", periode_str)
                            st.download_button(
                                label="📥 Télécharger",
                                data=pptx_io,
                                file_name=f"taux_heure_affid_{periode_str.replace(' ', '_').replace('/', '-')}.pptx",
                                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                            )
                        except Exception as e:
                            st.error(f"Erreur : {e}")
                
                st.plotly_chart(kpis['evo_appels_tickets'])
                col1, col2 = st.columns([3, 1])
                with col2:
                    if st.button("📊 Exporter évolution hebdomadaire en PowerPoint", key="btn_evo_appels_tickets_affid_hebdo"):
                        try:
                            pptx_io = create_single_graph_powerpoint(kpis['evo_appels_tickets'], "Évolution hebdomadaire : appels entrants + tickets Affid", periode_str)
                            st.download_button(
                                label="📥 Télécharger",
                                data=pptx_io,
                                file_name=f"evolution_hebdomadaire_appels_tickets_affid_{periode_str.replace(' ', '_').replace('/', '-')}.pptx",
                                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                            )
                        except Exception as e:
                            st.error(f"Erreur : {e}")
                
                st.plotly_chart(kpis.get('evo_appels_tickets_mensuel'))
                col1, col2 = st.columns([3, 1])
                with col2:
                    if st.button("📊 Exporter évolution mensuelle en PowerPoint", key="btn_evo_appels_tickets_affid_mensuel"):
                        try:
                            pptx_io = create_single_graph_powerpoint(kpis.get('evo_appels_tickets_mensuel'), "Évolution mensuelle : appels entrants + tickets Affid", periode_str)
                            st.download_button(
                                label="📥 Télécharger",
                                data=pptx_io,
                                file_name=f"evolution_mensuelle_appels_tickets_affid_{periode_str.replace(' ', '_').replace('/', '-')}.pptx",
                                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                            )
                        except Exception as e:
                            st.error(f"Erreur : {e}")
                
                col1, col2 = st.columns([3, 1])
                with col2:
                    if st.button("📊 Exporter en PowerPoint", key="btn_evo_appels_tickets_affid"):
                        try:
                            pptx_io = create_single_graph_powerpoint(kpis['evo_appels_tickets'], "Évolution appels/tickets Affid", periode_str)
                            st.download_button(
                                label="📥 Télécharger",
                                data=pptx_io,
                                file_name=f"evolution_appels_tickets_affid_{periode_str.replace(' ', '_').replace('/', '-')}.pptx",
                                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                            )
                        except Exception as e:
                            st.error(f"Erreur : {e}")

            import os


            if st.sidebar.button("Générer Rapport PowerPoint support 6mois"):
                try:
                    # Construction du DataFrame support_stellair comme dans la page
                    df_stellair = def_df_support(df_aircall_processed, df_aircall_processed, line_tous, agents_all)
                    df_stellair = df_stellair[(df_stellair['line'] == 'armatistechnique') | (df_stellair['Logiciel'] == 'Stellair')]
                    # Calculer les KPIs sur 6 mois et 3 mois pour support_stellair
                    df_support_6m = filtrer_par_periode(df_stellair, "6 derniers mois")
                    df_tickets_6m = filtrer_par_periode(df_tickets_processed, "6 derniers mois")
                    kpis_6m = generate_kpis(df_support_6m, df_tickets_6m, agents_all)
                    # Générer le graphique temps de réponse
                    from data_processing.kpi_generation import calculate_ticket_response_time
                    pptx_io_6m = create_powerpoint(kpis_6m, df_support_6m, df_tickets_6m , '6m')
                    st.sidebar.success("Rapport PowerPoint support 6mois !")
                    st.sidebar.download_button(
                        label="📥 Télécharger le rapport minimaliste",
                        data=pptx_io_6m,  
                        file_name="rapport_minimaliste_6m.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                    )
                except Exception as e:
                    st.sidebar.error(f"Erreur lors de la génération du PowerPoint : {e}")

            if st.sidebar.button("Générer Rapport PowerPoint support 3mois"):
                try:
                    # Construction du DataFrame support_stellair comme dans la page
                    df_stellair = def_df_support(df_aircall_processed, df_aircall_processed, line_tous, agents_all)
                    df_stellair = df_stellair[(df_stellair['line'] == 'armatistechnique') | (df_stellair['Logiciel'] == 'Stellair')]
                    df_support_3m = filtrer_par_periode(df_stellair, "3 derniers mois")
                    df_tickets_3m = filtrer_par_periode(df_tickets_processed, "3 derniers mois")
                    kpis_3m = generate_kpis(df_support_3m, df_tickets_3m, agents_all)
                    # Générer le graphique temps de réponse
                    moyenne_temps_reponse_3m, fig_temps_reponse_3m, _ = calculate_ticket_response_time(df_tickets_3m, agents_all)
                    pptx_io_3m = create_powerpoint(kpis_3m, df_support_3m, '3m')

                    st.sidebar.success("Rapport PowerPoint support 3mois !")
                    st.sidebar.download_button(
                        label="📥 Télécharger le rapport minimaliste",
                        data=pptx_io_3m,  
                        file_name="rapport_minimaliste_6m.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                    )
                except Exception as e:
                    st.sidebar.error(f"Erreur lors de la génération du PowerPoint : {e}")
            
            # Bouton pour télécharger les données Aircall traitées
            st.sidebar.markdown("---")
            st.sidebar.markdown("### 📊 Export des données")
            if st.sidebar.button("📥 Télécharger données Support (def_df_support)"):
                try:
                    # Récupérer le DataFrame final actuellement utilisé selon l'option sélectionnée
                    df_final_support = st.session_state.get(cache_key)
                    if df_final_support is None:
                        # Secours: recalculer selon l'option
                        if dataframe_option == "support_stellair":
                            df_final_support = def_df_support(df_aircall_processed, df_aircall_processed, line_tous, agents_all)
                            df_final_support = df_final_support[(df_final_support['line'] == 'armatistechnique') | (df_final_support['Logiciel'] == 'Stellair')]
                        elif dataframe_option == "support_affid":
                            df_final_support = def_df_support(df_aircall_processed, df_aircall_processed, line_tous, agents_all)
                            df_final_support = df_final_support[(df_final_support['Logiciel'] == 'Affid')]
                        elif dataframe_option == "xmed":
                            df_final_support = def_df_support(df_aircall_processed, df_aircall_processed, line_tous, agents_all)
                            df_final_support = df_final_support[(df_final_support['line'] == 'xmed')]
                        elif dataframe_option == "tmaj":
                            df_final_support = def_df_support(df_aircall_processed, df_aircall_processed, line_tous, agents_all)
                            df_final_support = df_final_support[(df_final_support['line'] == 'supporthardware')]

                    # Convertir en CSV
                    csv_data = df_final_support.to_csv(index=False, encoding='utf-8-sig')

                    st.sidebar.success("Données Support (def_df_support) prêtes !")
                    st.sidebar.download_button(
                        label="📥 Télécharger CSV",
                        data=csv_data,
                        file_name=f"donnees_support_{dataframe_option}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.csv",
                        mime="text/csv"
                    )
                except Exception as e:
                    st.sidebar.error(f"Erreur lors de l'export des données : {e}")

        support()



    elif selection_page == "Agents":
        import pandas as pd
        from data_processing.kpi_generation import (
            filtrer_par_periode, generate_kpis, convert_to_sixtieth, 
            filtrer_par_agent, charge_entrant_sortant, df_compute_ticket_appels_metrics, 
            historique_scores_total, charge_ticket, 
            filter_evaluation, calculate_performance_score,
            graph_tickets_n2_par_semaine, graph_tickets_n2_resolus_par_agent, graph_tickets_pierre_goupillon,
            graph_tickets_n1_par_semaine, get_n1_agents_list
        )
        from data_processing.hubspot_processing import process_hubspot_data
        from data_processing.aircall_processing import process_aircall_data, def_df_support, agents_all, line_tous


        

        def agents():
            # Chargement des données
            periode_option = st.selectbox("Période KPI automatique :", 
                ["1 an", "6 derniers mois", "3 derniers mois", "Dernier mois", "Personnalisée"],
                index=0, key="periode_agents")
            
            if periode_option == "Personnalisée":
                df_temp = def_df_support(process_aircall_data(df_aircall), process_aircall_data(df_aircall), line_tous, agents_all)
                df_temp['Date'] = pd.to_datetime(df_temp['Date'])
                date_min_data = (df_temp['Date'].min() or pd.Timestamp.now() - pd.DateOffset(years=2)).date()
                date_max_data = (df_temp['Date'].max() or pd.Timestamp.now()).date()
                col_d1, col_d2 = st.columns(2)
                with col_d1:
                    date_debut = st.date_input("Date de début", value=date_max_data - timedelta(days=90), min_value=date_min_data, max_value=date_max_data, key="date_debut_agents")
                with col_d2:
                    date_fin = st.date_input("Date de fin", value=date_max_data, min_value=date_debut, max_value=date_max_data, key="date_fin_agents")
                periode_selectbox = (date_debut, date_fin)
                periode_str = get_periode_str(periode_selectbox)
            else:
                periode_selectbox = periode_option
                periode_str = periode_selectbox
            
            # Paramètres de scoring
            st.sidebar.markdown("## ⚙️ Paramètres de scoring")
            with st.sidebar.expander("🎯 Objectifs", expanded=True):
                objectif_total = st.number_input("Objectif total de demandes par jour", min_value=1, max_value=50, value=30)
                ratio_appels = st.slider("Ratio d'appels (%)", min_value=0, max_value=100, value=70) / 100
                ratio_tickets = 1 - ratio_appels  # Calculé automatiquement
                st.write(f"Ratio de tickets : {ratio_tickets:.0%}")
                objectif_taux_service = st.slider("Objectif taux de service (%)", min_value=0, max_value=100, value=70) / 100
            
            # Bouton pour vider le cache de la page agents
            if st.sidebar.button("🔄 Vider le cache de la page agents", help="Force le recalcul de tous les graphiques et métriques"):
                # Supprimer tous les caches liés à la page agents
                cache_keys_to_clear = [
                    'agents_data_', 'agents_metrics_', 'agents_graphs_', 'agents_n1_list_',
                    'n1_graph_', 'historique_scores_', 'n2_graph_', 'cumul_n1_', 'cumul_n2_'
                ]
                
                for key in list(st.session_state.keys()):
                    if any(prefix in key for prefix in cache_keys_to_clear):
                        del st.session_state[key]
                
                st.sidebar.success("✅ Cache de la page agents vidé !")
                st.rerun()
            
            # Cache pour les données filtrées par période
            cache_key_data = f'agents_data_{periode_str}_{hash(str(df_tickets_processed.shape))}_{hash(str(df_aircall.shape))}'
            
            if cache_key_data not in st.session_state:
                with st.spinner(f"Chargement des données pour la période {periode_str}..."):
                    df_support = def_df_support(process_aircall_data(df_aircall), process_aircall_data(df_aircall), line_tous, agents_all)
                    df_support_filtered = filtrer_par_periode(df_support, periode_selectbox)
                    df_tickets_filtered = filtrer_par_periode(df_tickets_processed, periode_selectbox)
                    
                    st.session_state[cache_key_data] = {
                        'df_support_filtered': df_support_filtered,
                        'df_tickets_filtered': df_tickets_filtered
                    }
                    st.success(f"✅ Données mises en cache pour {periode_str}")
            else:
                df_support_filtered = st.session_state[cache_key_data]['df_support_filtered']
                df_tickets_filtered = st.session_state[cache_key_data]['df_tickets_filtered']
                st.info(f"📋 Utilisation des données en cache pour {periode_str}")
            
            agents_n1 = ['Mourad HUMBLOT', 'Archimede KESSI', 'Celine Crendal', 'Melinda Marmin', 'Emilie GEST', 'Sandrine Sauvage', 'Cédeline DUVAL']
            agents_n1_tickets = agents_n1 + ['Frederic SAUVAN']
            agents_scores = ['Mourad HUMBLOT', 'Archimede KESSI', 'Celine Crendal', 'Emilie GEST', 'Cédeline DUVAL']

            agents_score_suresnes = ['Mourad HUMBLOT', 'Archimede KESSI', 'Celine Crendal', 'Emilie GEST', 'Cédeline DUVAL']
            agents_score_calais = ['Celine Crendal', 'Melinda Marmin', 'Emilie GEST', 'Sandrine Sauvage', 'Cédeline DUVAL']

            # Cache pour les métriques de scoring
            cache_key_metrics = f'agents_metrics_{periode_str}_{objectif_total}_{ratio_appels}_{objectif_taux_service}_{hash(str(df_tickets_filtered.shape))}_{hash(str(agents_score_suresnes))}'
            
            if cache_key_metrics not in st.session_state:
                with st.spinner("Calcul des métriques de scoring..."):
                    df_conforme = df_compute_ticket_appels_metrics(agents_score_suresnes, df_tickets_filtered, df_support_filtered)
                    df_conforme_calais = df_compute_ticket_appels_metrics(agents_score_calais, df_tickets_filtered, df_support_filtered)
                    
                    st.session_state[cache_key_metrics] = {
                        'df_conforme': df_conforme,
                        'df_conforme_calais': df_conforme_calais
                    }
                    st.success("✅ Métriques de scoring calculées et mises en cache")
            else:
                df_conforme = st.session_state[cache_key_metrics]['df_conforme']
                df_conforme_calais = st.session_state[cache_key_metrics]['df_conforme_calais']
                st.info("📋 Utilisation des métriques en cache")

            # Vérifier que df_conforme est un DataFrame avant d'appliquer les calculs
            if isinstance(df_conforme, pd.DataFrame):
                # Appliquer le score de performance avec les paramètres personnalisés
                df_conforme['score_performance'] = df_conforme.apply(
                    lambda row: calculate_performance_score(
                        row,
                        objectif_total=objectif_total,
                        ratio_appels=ratio_appels,
                        ratio_tickets=ratio_tickets,
                        objectif_taux_service=objectif_taux_service
                    ),
                    axis=1
                )

                df_conforme = df_conforme[['Agent', 'score_performance', "Nombre d'appel traité", 
                                        'Nombre de ticket traité', '% appel entrant agent', 
                                        '% tickets', '% appels']]
            else:
                st.error("Erreur: Les données ne sont pas dans le format attendu.")
                return

            # Fonction de style
            def style_scores(df):
                def apply_style(row):
                    styles = {}
                    score = row['score_performance']
                    if score >= 70:
                        styles['score_performance'] = 'background-color: #a7dba7'  # vert
                    elif score >= 60:
                        styles['score_performance'] = 'background-color: #f7c97f'  # orange
                    else:
                        styles['score_performance'] = 'background-color: #f28e8e'  # rouge
                    return pd.Series(styles)
                return df.style.apply(apply_style, axis=1).format({
                    'score_performance': '{:.2f}',
                    "Nombre d'appel traité": '{:.2f}',
                    'Nombre de ticket traité': '{:.2f}',
                    '% appel entrant agent': '{:.2f}',
                    '% tickets': '{:.2f}',
                    '% appels': '{:.2f}'
                })

            # ----------- SIDEBAR - FILTRES -------------------
            with st.sidebar:
                st.header("🔍 Filtres évaluations")
                with st.form("filter_form"):
                    agents = st.multiselect("Agent(s)", options=df_evaluation["agent"].unique(), default=df_evaluation["agent"].unique())
                    periodes_eval = st.multiselect("Période", options=df_evaluation["quarter"].unique(), default=df_evaluation["quarter"].unique())
                    submit = st.form_submit_button(label="🎯 Appliquer les filtres")

            # ----------- TITRE -------------------
            st.title("📊 Tableau de bord des agents N1")

            st.markdown(f"**Période KPI automatique sélectionnée** : `{periode_str}`")

            # ----------- SCORING AUTOMATIQUE -------------------
            st.markdown("## ⚖️ Scores automatiques")

            st.markdown("""
            ### 🧠 Méthodologie de scoring

            Le score de performance (0 à 100) est calculé selon quatre critères :
            - Volume total (55%) : Nombre total de demandes traitées par rapport à l'objectif
            - Répartition (15%) : Équilibre entre appels et tickets selon le ratio défini
            - Comparaison à la moyenne (15%) : Performance par rapport à la moyenne du service
            - Taux d'appels entrants (15%) : Pourcentage d'appels entrants par rapport au total

            Codes couleur :
            - 🟢 **Vert** : score ≥ 70%
            - 🟠 **Orange** : 60% ≤ score < 70%
            - 🔴 **Rouge** : score < 60%
            """)

            # Calcul du nombre total de tickets traités par Céline Crendal (toutes orthographes)


            styled_df = style_scores(df_conforme)
            st.dataframe(styled_df, use_container_width=True)

            st.dataframe(df_conforme_calais, use_container_width=True)


            # ----------- EXPORT EXCEL PAR AGENT ET PAR JOUR -------------------
            st.markdown("### 📥 Export des données détaillées")
            
            def generate_agent_daily_data(df_support, df_tickets, agents_list):
                """
                Génère un DataFrame avec le nombre d'appels et de tickets traités par agent et par jour.
                
                Args:
                    df_support: DataFrame contenant les données d'appels
                    df_tickets: DataFrame contenant les données de tickets
                    agents_list: Liste des agents à inclure
                    
                Returns:
                    DataFrame avec colonnes: Agent, Date, Nombre d'appels traités, Nombre de tickets traités
                """
                import pandas as pd
                
                # Dictionnaire de correspondance des noms
                correspondance_noms = {
                    'Celine Crendal': ['Celine Crendal', 'Céline Crendal', 'CELINE CRENDAL', 'CÉLINE CRENDAL'],
                    'Cédeline DUVAL': ['Cédeline DUVAL', 'Cedeline DUVAL', 'CÉDELINE DUVAL', 'CEDELINE DUVAL', 'CÃ©deline DUVAL', 'cedeline duval'],
                }
                
                # Liste pour stocker les résultats
                results = []
                
                for agent in agents_list:
                    # Gestion des variantes de nom
                    if agent in correspondance_noms:
                        selected_agents = correspondance_noms[agent]
                    else:
                        selected_agents = [agent]
                    
                    # Filtrer les appels pour cet agent (appels connectés uniquement)
                    df_agent_calls = df_support[
                        (df_support['UserName'].isin(selected_agents)) & 
                        (df_support['LastState'] == 'yes')
                    ]
                    
                    # Compter les appels par jour
                    if not df_agent_calls.empty:
                        calls_per_day = df_agent_calls.groupby('Date').size().reset_index(name='Nombre_appels')
                    else:
                        calls_per_day = pd.DataFrame(columns=['Date', 'Nombre_appels'])
                    
                    # Filtrer les tickets pour cet agent
                    df_agent_tickets = df_tickets[
                        (df_tickets['Propriétaire du ticket'].isin(selected_agents)) & 
                        (df_tickets['Pipeline'].isin(['SSIA', 'SSI', 'SPSA'])) &
                        (df_tickets['Source'].isin(['Chat', 'E-mail', 'Formulaire']) | pd.isna(df_tickets['Source']))
                    ]
                    
                    # Compter les tickets par jour
                    if not df_agent_tickets.empty:
                        tickets_per_day = df_agent_tickets.groupby('Date').size().reset_index(name='Nombre_tickets')
                    else:
                        tickets_per_day = pd.DataFrame(columns=['Date', 'Nombre_tickets'])
                    
                    # Obtenir toutes les dates de la période
                    all_dates = pd.DataFrame({'Date': df_support['Date'].unique()})
                    
                    # Fusionner les données
                    agent_data = all_dates.merge(calls_per_day, on='Date', how='left')
                    agent_data = agent_data.merge(tickets_per_day, on='Date', how='left')
                    
                    # Remplir les valeurs manquantes avec 0
                    agent_data['Nombre_appels'] = agent_data['Nombre_appels'].fillna(0).astype(int)
                    agent_data['Nombre_tickets'] = agent_data['Nombre_tickets'].fillna(0).astype(int)
                    
                    # Ajouter le nom de l'agent
                    agent_data['Agent'] = agent
                    
                    # Réorganiser les colonnes
                    agent_data = agent_data[['Agent', 'Date', 'Nombre_appels', 'Nombre_tickets']]
                    
                    results.append(agent_data)
                
                # Combiner tous les résultats
                final_df = pd.concat(results, ignore_index=True)
                
                # Trier par agent puis par date
                final_df = final_df.sort_values(['Agent', 'Date']).reset_index(drop=True)
                
                return final_df
            
            # Clé pour le cache de l'export
            export_cache_key = f'agent_export_{periode_str}'
            
            # Bouton pour générer l'export Excel
            if st.button("📊 Générer l'export Excel par agent et par jour", key="btn_export_agent_daily"):
                try:
                    with st.spinner("Génération de l'export Excel..."):
                        # Générer les données
                        df_export = generate_agent_daily_data(df_support_filtered, df_tickets_filtered, agents_n1)
                        
                        # Créer le fichier Excel en mémoire
                        output = io.BytesIO()
                        with pd.ExcelWriter(output, engine='openpyxl') as writer:
                            df_export.to_excel(writer, sheet_name='Données par agent et par jour', index=False)
                            
                            # Obtenir le workbook et la feuille pour la mise en forme
                            workbook = writer.book
                            worksheet = writer.sheets['Données par agent et par jour']
                            
                            # Ajuster la largeur des colonnes
                            worksheet.column_dimensions['A'].width = 25  # Agent
                            worksheet.column_dimensions['B'].width = 15  # Date
                            worksheet.column_dimensions['C'].width = 20  # Nombre_appels
                            worksheet.column_dimensions['D'].width = 20  # Nombre_tickets
                            
                            # Formater la ligne d'en-tête
                            from openpyxl.styles import Font, PatternFill, Alignment
                            header_fill = PatternFill(start_color="366092", end_color="366092", fill_type="solid")
                            header_font = Font(bold=True, color="FFFFFF")
                            
                            for cell in worksheet[1]:
                                cell.fill = header_fill
                                cell.font = header_font
                                cell.alignment = Alignment(horizontal='center', vertical='center')
                            
                            # Créer un résumé par agent
                            df_summary = df_export.groupby('Agent').agg({
                                'Nombre_appels': 'sum',
                                'Nombre_tickets': 'sum'
                            }).reset_index()
                            
                            # Calculer les moyennes par jour
                            nb_jours = df_export['Date'].nunique()
                            df_summary['Moyenne_appels_par_jour'] = (df_summary['Nombre_appels'] / nb_jours).round(2)
                            df_summary['Moyenne_tickets_par_jour'] = (df_summary['Nombre_tickets'] / nb_jours).round(2)
                            
                            # Ajouter la feuille de résumé
                            df_summary.to_excel(writer, sheet_name='Résumé par agent', index=False)
                            
                            # Formater la feuille de résumé
                            summary_sheet = writer.sheets['Résumé par agent']
                            summary_sheet.column_dimensions['A'].width = 25
                            summary_sheet.column_dimensions['B'].width = 20
                            summary_sheet.column_dimensions['C'].width = 20
                            summary_sheet.column_dimensions['D'].width = 25
                            summary_sheet.column_dimensions['E'].width = 25
                            
                            for cell in summary_sheet[1]:
                                cell.fill = header_fill
                                cell.font = header_font
                                cell.alignment = Alignment(horizontal='center', vertical='center')
                        
                        output.seek(0)
                        
                        # Stocker dans session_state
                        st.session_state[export_cache_key] = {
                            'excel_data': output.getvalue(),
                            'df_preview': df_export,
                            'nb_lignes': len(df_export)
                        }
                        
                        st.success(f"✅ Export généré avec succès ! {len(df_export)} lignes de données.")
                        
                except Exception as e:
                    st.error(f"❌ Erreur lors de la génération de l'export : {str(e)}")
                    import traceback
                    st.error(traceback.format_exc())
            
            # Afficher l'aperçu et le bouton de téléchargement si l'export existe
            if export_cache_key in st.session_state:
                export_data = st.session_state[export_cache_key]
                
                st.success(f"✅ Export prêt ! {export_data['nb_lignes']} lignes de données.")
                st.markdown("**Aperçu des données :**")
                st.dataframe(export_data['df_preview'].head(20), use_container_width=True)
                
                # Bouton de téléchargement persistant
                st.download_button(
                    label="📥 Télécharger le fichier Excel",
                    data=export_data['excel_data'],
                    file_name=f"export_agents_appels_tickets_{periode_str.replace(' ', '_').replace('/', '-')}.xlsx",
                    mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                    key="download_agent_daily_excel"
                )

            # ----------- GRAPHIQUES KPI -------------------
            st.markdown("## 📈 Graphiques d'activité")

            # Cache pour les graphiques d'activité
            cache_key_graphs = f'agents_graphs_{periode_str}_{hash(str(df_support_filtered.shape))}_{hash(str(df_tickets_filtered.shape))}'
            
            if cache_key_graphs not in st.session_state:
                with st.spinner("Génération des graphiques d'activité..."):
                    # Créer les graphiques avec template sombre et couleurs vives
                    fig_line_entrant, fig_pie_entrant = charge_entrant_sortant(df_support_filtered, agents_n1)
                    fig_line_ticket, fig_pie_ticket = charge_ticket(df_tickets_filtered, agents_n1)
                    
                    st.session_state[cache_key_graphs] = {
                        'fig_line_entrant': fig_line_entrant,
                        'fig_pie_entrant': fig_pie_entrant,
                        'fig_line_ticket': fig_line_ticket,
                        'fig_pie_ticket': fig_pie_ticket
                    }
                    st.success("✅ Graphiques d'activité générés et mis en cache")
            else:
                fig_line_entrant = st.session_state[cache_key_graphs]['fig_line_entrant']
                fig_pie_entrant = st.session_state[cache_key_graphs]['fig_pie_entrant']
                fig_line_ticket = st.session_state[cache_key_graphs]['fig_line_ticket']
                fig_pie_ticket = st.session_state[cache_key_graphs]['fig_pie_ticket']
                st.info("📋 Utilisation des graphiques d'activité en cache")

            st.markdown("### 📞 Activité téléphonique")
            col1, col2 = st.columns(2)
            col1.plotly_chart(fig_line_entrant, use_container_width=True)
            col2.plotly_chart(fig_pie_entrant, use_container_width=True)
            
            # Boutons d'export pour les graphiques téléphoniques
            col1, col2 = st.columns(2)
            with col1:
                if st.button("📊 Exporter en PowerPoint", key="btn_activite_telephone"):
                    try:
                        pptx_io = create_single_graph_powerpoint(fig_line_entrant, "Activité téléphonique - Ligne", periode_str)
                        st.download_button(
                            label="📥 Télécharger",
                            data=pptx_io,
                            file_name=f"activite_telephone_ligne_{periode_str.replace(' ', '_').replace('/', '-')}.pptx",
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                        )
                    except Exception as e:
                        st.error(f"Erreur : {e}")
            
            with col2:
                if st.button("📊 Exporter en PowerPoint", key="btn_repartition_telephone"):
                    try:
                        pptx_io = create_single_graph_powerpoint(fig_pie_entrant, "Répartition téléphonique", periode_str)
                        st.download_button(
                            label="📥 Télécharger",
                            data=pptx_io,
                            file_name=f"repartition_telephone_{periode_str.replace(' ', '_').replace('/', '-')}.pptx",
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                        )
                    except Exception as e:
                        st.error(f"Erreur : {e}")

            st.markdown("### 🧾 Activité tickets")
            col1, col2 = st.columns(2)
            col1.plotly_chart(fig_line_ticket, use_container_width=True)
            col2.plotly_chart(fig_pie_ticket, use_container_width=True)
            
            # Boutons d'export pour les graphiques tickets
            col1, col2 = st.columns(2)
            with col1:
                if st.button("📊 Exporter en PowerPoint", key="btn_activite_tickets"):
                    try:
                        pptx_io = create_single_graph_powerpoint(fig_line_ticket, "Activité tickets - Ligne", periode_str)
                        st.download_button(
                            label="📥 Télécharger",
                            data=pptx_io,
                            file_name=f"activite_tickets_ligne_{periode_str.replace(' ', '_').replace('/', '-')}.pptx",
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                        )
                    except Exception as e:
                        st.error(f"Erreur : {e}")
            
            with col2:
                if st.button("📊 Exporter en PowerPoint", key="btn_repartition_tickets"):
                    try:
                        pptx_io = create_single_graph_powerpoint(fig_pie_ticket, "Répartition tickets", periode_str)
                        st.download_button(
                            label="📥 Télécharger",
                            data=pptx_io,
                            file_name=f"repartition_tickets_{periode_str.replace(' ', '_').replace('/', '-')}.pptx",
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                        )
                    except Exception as e:
                        st.error(f"Erreur : {e}")

            # ----------- GRAPHIQUE N1 PAR SEMAINE -------------------
            st.markdown("## 📊 Graphique N1 par semaine")
            
            # Cache pour les agents N1 disponibles
            cache_key_agents_n1 = f'agents_n1_list_{periode_str}_{hash(str(df_tickets_filtered.shape))}'
            
            if cache_key_agents_n1 not in st.session_state:
                agents_n1_disponibles = get_n1_agents_list(df_tickets_filtered)
                st.session_state[cache_key_agents_n1] = agents_n1_disponibles
            else:
                agents_n1_disponibles = st.session_state[cache_key_agents_n1]
            
            # Sélecteur d'agent pour le graphique N1
            if agents_n1_disponibles:
                agent_n1_selection = st.selectbox(
                    "Sélectionner un agent N1 pour afficher ses tickets ouverts :",
                    ["Aucun"] + agents_n1_disponibles,
                    key="agent_n1_selection"
                )
                
                # Cache pour le graphique N1
                cache_key_n1_graph = f'n1_graph_{periode_str}_{agent_n1_selection}_{hash(str(df_tickets_filtered.shape))}'
                
                if cache_key_n1_graph not in st.session_state:
                    with st.spinner("Génération du graphique N1..."):
                        if agent_n1_selection != "Aucun":
                            fig_n1 = graph_tickets_n1_par_semaine(df_tickets_filtered, selected_agent=agent_n1_selection)
                        else:
                            fig_n1 = graph_tickets_n1_par_semaine(df_tickets_filtered)
                        st.session_state[cache_key_n1_graph] = fig_n1
                        st.success("✅ Graphique N1 généré et mis en cache")
                else:
                    fig_n1 = st.session_state[cache_key_n1_graph]
                    st.info("📋 Utilisation du graphique N1 en cache")
                
                st.plotly_chart(fig_n1, use_container_width=True)
                
                # Bouton d'export pour le graphique N1
                col1, col2 = st.columns([3, 1])
                with col2:
                    if st.button("📊 Exporter en PowerPoint", key="btn_graphique_n1"):
                        try:
                            pptx_io = create_single_graph_powerpoint(fig_n1, f"Tickets N1 par semaine - {agent_n1_selection}", periode_str)
                            st.download_button(
                                label="📥 Télécharger",
                                data=pptx_io,
                                file_name=f"tickets_n1_semaine_{agent_n1_selection.replace(' ', '_')}_{periode_str.replace(' ', '_').replace('/', '-')}.pptx",
                                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                            )
                        except Exception as e:
                            st.error(f"Erreur : {e}")
            else:
                st.info("Aucun agent N1 trouvé pour la période sélectionnée.")

            # ----------- HISTORIQUE SCORES -------------------
            st.markdown("## 📊 Historique des scores")
            
            # Cache pour l'historique des scores
            cache_key_historique = f'historique_scores_{periode_str}_{objectif_total}_{ratio_appels}_{ratio_tickets}_{objectif_taux_service}_{hash(str(df_tickets_filtered.shape))}_{hash(str(df_support_filtered.shape))}_{hash(str(agents_scores))}'
            
        
            
            if cache_key_historique not in st.session_state:
                with st.spinner("Génération de l'historique des scores..."):
                    fig_historique = historique_scores_total(
                        agents_scores, 
                        df_tickets_filtered, 
                        df_support_filtered,
                        objectif_total=objectif_total,
                        ratio_appels=ratio_appels,
                        ratio_tickets=ratio_tickets,
                        objectif_taux_service=objectif_taux_service
                    )
                    st.session_state[cache_key_historique] = fig_historique
                    st.success("✅ Historique des scores généré et mis en cache")
            else:
                fig_historique = st.session_state[cache_key_historique]
                st.info("📋 Utilisation de l'historique des scores en cache")
            
            st.plotly_chart(fig_historique, use_container_width=True)
            
            # Bouton d'export pour l'historique des scores
            col1, col2 = st.columns([3, 1])
            with col2:
                if st.button("📊 Exporter en PowerPoint", key="btn_historique_scores"):
                    try:
                        pptx_io = create_single_graph_powerpoint(fig_historique, "Historique des scores", periode_str)
                        st.download_button(
                            label="📥 Télécharger",
                            data=pptx_io,
                            file_name=f"historique_scores_{periode_str.replace(' ', '_').replace('/', '-')}.pptx",
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                        )
                    except Exception as e:
                        st.error(f"Erreur : {e}")

            # ----------- TABLEAU DE BORD N1 -------------------
            st.markdown("## Tableau de bord des agents N2")

            # Cache pour le graphique N2
            cache_key_n2_graph = f'n2_graph_{periode_str}_{hash(str(df_tickets_filtered.shape))}'
            
            if cache_key_n2_graph not in st.session_state:
                with st.spinner("Génération du graphique N2..."):
                    fig1 = graph_tickets_n2_par_semaine(df_tickets_filtered)
                    st.session_state[cache_key_n2_graph] = fig1
                    st.success("✅ Graphique N2 généré et mis en cache")
            else:
                fig1 = st.session_state[cache_key_n2_graph]
                st.info("📋 Utilisation du graphique N2 en cache")
            
            st.plotly_chart(fig1, use_container_width=True)
            
            # Bouton d'export pour le graphique N2
            col1, col2 = st.columns([3, 1])
            with col2:
                if st.button("📊 Exporter en PowerPoint", key="btn_graphique_n2"):
                    try:
                        pptx_io = create_single_graph_powerpoint(fig1, "Tickets N2 par semaine", periode_str)
                        st.download_button(
                            label="📥 Télécharger",
                            data=pptx_io,
                            file_name=f"tickets_n2_semaine_{periode_str.replace(' ', '_').replace('/', '-')}.pptx",
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                        )
                    except Exception as e:
                        st.error(f"Erreur : {e}")

            # Graphique 2
            #fig2 = graph_tickets_n2_resolus_par_agent(filtrer_par_periode(df_tickets_processed, periode_selectbox))
            #st.plotly_chart(fig2, use_container_width=True)

            # Graphique 3
            #fig3 = graph_tickets_pierre_goupillon(filtrer_par_periode(df_tickets_processed, periode_selectbox))
            #st.plotly_chart(fig3, use_container_width=True)

            # Graphique des tickets ouverts de Pierre Goupillon (comparatif)
            from data_processing.kpi_generation import graph_tickets_ouverts_pierre_goupillon
            #fig3_ouverts = graph_tickets_ouverts_pierre_goupillon(filtrer_par_periode(df_tickets_processed, periode_selectbox))
           # st.plotly_chart(fig3_ouverts, use_container_width=True)

            # ----------- GRAPHIQUE CUMULATIF N2 -------------------
            st.markdown("## 📈 Cumulatif tickets N2 (passés vs résolus)")
            from data_processing.kpi_generation import graph_tickets_n2_cumulatif
            
            # Cache pour le graphique cumulatif N2
            cache_key_cumul_n2 = f'cumul_n2_{periode_str}_{hash(str(df_tickets_filtered.shape))}'
            
            if cache_key_cumul_n2 not in st.session_state:
                with st.spinner("Génération du graphique cumulatif N2..."):
                    fig_cumul_n2, pipelines_disponibles, tickets_en_cours = graph_tickets_n2_cumulatif(df_tickets_filtered)
                    st.session_state[cache_key_cumul_n2] = {
                        'fig_cumul_n2': fig_cumul_n2,
                        'pipelines_disponibles': pipelines_disponibles,
                        'tickets_en_cours': tickets_en_cours
                    }
                    st.success("✅ Graphique cumulatif N2 généré et mis en cache")
            else:
                fig_cumul_n2 = st.session_state[cache_key_cumul_n2]['fig_cumul_n2']
                pipelines_disponibles = st.session_state[cache_key_cumul_n2]['pipelines_disponibles']
                tickets_en_cours = st.session_state[cache_key_cumul_n2]['tickets_en_cours']
                st.info("📋 Utilisation du graphique cumulatif N2 en cache")
            
            # Filtre par pipeline
            if pipelines_disponibles:
                pipeline_selection = st.selectbox(
                    "Filtrer par pipeline :",
                    ["Tous les pipelines"] + pipelines_disponibles,
                    key="pipeline_filter_cumul"
                )
                
                # Cache pour le graphique filtré par pipeline
                cache_key_cumul_n2_filtered = f'cumul_n2_filtered_{periode_str}_{pipeline_selection}_{hash(str(df_tickets_filtered.shape))}'
                
                if pipeline_selection != "Tous les pipelines":
                    if cache_key_cumul_n2_filtered not in st.session_state:
                        with st.spinner(f"Génération du graphique filtré pour {pipeline_selection}..."):
                            fig_cumul_n2, _, tickets_en_cours = graph_tickets_n2_cumulatif(df_tickets_filtered, pipeline_selection)
                            st.session_state[cache_key_cumul_n2_filtered] = {
                                'fig_cumul_n2': fig_cumul_n2,
                                'tickets_en_cours': tickets_en_cours
                            }
                            st.success(f"✅ Graphique filtré pour {pipeline_selection} généré et mis en cache")
                    else:
                        fig_cumul_n2 = st.session_state[cache_key_cumul_n2_filtered]['fig_cumul_n2']
                        tickets_en_cours = st.session_state[cache_key_cumul_n2_filtered]['tickets_en_cours']
                        st.info(f"📋 Utilisation du graphique filtré pour {pipeline_selection} en cache")
            
            st.plotly_chart(fig_cumul_n2, use_container_width=True)
            
            # Bouton d'export pour le graphique cumulatif N2
            col1, col2 = st.columns([3, 1])
            with col2:
                if st.button("📊 Exporter en PowerPoint", key="btn_cumul_n2"):
                    try:
                        title = f"Cumulatif tickets N2 - {pipeline_selection}" if pipeline_selection != "Tous les pipelines" else "Cumulatif tickets N2"
                        pptx_io = create_single_graph_powerpoint(fig_cumul_n2, title, periode_str)
                        st.download_button(
                            label="📥 Télécharger",
                            data=pptx_io,
                            file_name=f"cumulatif_n2_{pipeline_selection.replace(' ', '_')}_{periode_str.replace(' ', '_').replace('/', '-')}.pptx",
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                        )
                    except Exception as e:
                        st.error(f"Erreur : {e}")
            
            # Bouton pour afficher les tickets en cours
            if not tickets_en_cours.empty:
                nb_tickets_en_cours = len(tickets_en_cours)
                st.markdown(f"**🔴 {nb_tickets_en_cours} ticket(s) N2 en cours**")
                
                # Bouton pour télécharger le DataFrame des tickets en cours
                if st.button(f"📥 Télécharger la liste des {nb_tickets_en_cours} tickets en cours", key="btn_download_tickets_en_cours"):
                    # Préparer les données pour le téléchargement
                    tickets_download = tickets_en_cours[['Ticket ID', 'Date', 'Statut du ticket']].copy()
                    tickets_download = tickets_download.rename(columns={
                        'Ticket ID': 'ID_Ticket',
                        'Date': 'Date_Derniere_Mise_A_Jour',
                        'Statut du ticket': 'Statut'
                    })
                    tickets_download = tickets_download.sort_values('Date_Derniere_Mise_A_Jour', ascending=False)
                    
                    # Convertir en CSV
                    csv_data = tickets_download.to_csv(index=False, encoding='utf-8-sig')
                    
                    # Bouton de téléchargement
                    st.download_button(
                        label="📄 Télécharger CSV des tickets en cours",
                        data=csv_data,
                        file_name=f"tickets_n2_en_cours_{periode_str.replace(' ', '_').replace('/', '-')}.csv",
                        mime="text/csv"
                    )
                    
                    # Afficher un aperçu du DataFrame
                    st.markdown("### 📋 Aperçu des tickets en cours")
                    st.dataframe(
                        tickets_download,
                        use_container_width=True,
                        column_config={
                            "ID_Ticket": st.column_config.TextColumn("ID Ticket", width="medium"),
                            "Date_Derniere_Mise_A_Jour": st.column_config.DateColumn("Date dernière mise à jour", width="medium"),
                            "Statut": st.column_config.TextColumn("Statut", width="small")
                        }
                    )
                
                if st.button(f"📋 Voir la liste complète des {nb_tickets_en_cours} tickets en cours", key="btn_tickets_en_cours"):
                    st.markdown("### 📋 Liste complète des tickets N2 en cours")
                    
                    # Préparer les données pour l'affichage complet
                    tickets_display = tickets_en_cours[['Ticket ID', 'Propriétaire du ticket', 'Pipeline', 'Statut du ticket', 'Date', 'Semaine']].copy()
                    tickets_display = tickets_display.sort_values('Date', ascending=False)
                    
                    # Afficher le tableau
                    st.dataframe(
                        tickets_display,
                        use_container_width=True,
                        column_config={
                            "Ticket ID": st.column_config.TextColumn("ID Ticket", width="medium"),
                            "Propriétaire du ticket": st.column_config.TextColumn("Propriétaire", width="medium"),
                            "Pipeline": st.column_config.TextColumn("Pipeline", width="small"),
                            "Statut du ticket": st.column_config.TextColumn("Statut", width="small"),
                            "Date": st.column_config.DateColumn("Date", width="small"),
                            "Semaine": st.column_config.TextColumn("Semaine", width="small")
                        }
                    )
                    
                    # Option pour télécharger la liste complète
                    csv_data_complete = tickets_display.to_csv(index=False, encoding='utf-8-sig')
                    st.download_button(
                        label="📥 Télécharger la liste complète en CSV",
                        data=csv_data_complete,
                        file_name=f"tickets_n2_en_cours_complet_{periode_str.replace(' ', '_').replace('/', '-')}.csv",
                        mime="text/csv"
                    )
            else:
                st.success("✅ Aucun ticket N2 en cours !")

            # ----------- ÉVALUATIONS MANAGERIALES ------------------
            
            if submit:
                df_evaluation_filre = filter_evaluation(df_evaluation, agents, periodes_eval)
                
                st.markdown("## 📝 Évaluations managériales")
                if not df_evaluation_filre.empty:
                    moyenne_score = round(df_evaluation_filre["average_score"].mean(), 2)
                    st.markdown(f"**🎯 Score moyen des évaluations : `{moyenne_score}`**")
                else:
                    st.warning("Aucune donnée disponible pour la sélection.")
                
                st.dataframe(df_evaluation_filre, use_container_width=True)

            def render_evaluation_table(df):
                """Affiche les évaluations en tableau Plotly (sans matplotlib pour déploiement léger)."""
                if df.empty:
                    st.warning("⚠️ Aucune donnée disponible pour la sélection.")
                    return
                import plotly.graph_objects as go
                moyenne_score = round(df["average_score"].mean(), 2)
                st.markdown(f"📝 **Évaluations managériales** — Score moyen : **{moyenne_score}**")
                df_display = df.reset_index()
                fig = go.Figure(data=[go.Table(
                    header=dict(values=list(df_display.columns), fill_color='#1f77b4', font=dict(color='white'), align='left'),
                    cells=dict(values=[df_display[c].astype(str).tolist() for c in df_display.columns], align='left')
                )])
                fig.update_layout(margin=dict(l=20, r=20, t=20, b=20), height=min(400, 50 + len(df) * 25))
                st.plotly_chart(fig, use_container_width=True)

            df_eval_filtered = filter_evaluation(df_evaluation, agents, periodes_eval)
            render_evaluation_table(df_eval_filtered)


            
        agents()

    elif selection_page == "Tickets":
        import pandas as pd
        from data_processing.kpi_generation import filtrer_par_periode, sla, metrics_nombre_ticket_categorie
        import plotly.express as px
        import plotly.graph_objects as go

        def tickets():
            

            periode = st.selectbox("Sélectionnez une période :", 
                    ["1 an", "Toute la période", "6 derniers mois", "3 derniers mois", "Dernier mois"],
                    index=0)
            
            # --- SYSTÈME DE CACHE POUR LES KPI ---
            # Utiliser les données déjà traitées en cache
            df_support = st.session_state['df_support_processed']
            df_tickets_processed = st.session_state['df_tickets_processed']
            
            # Cache pour les KPI par période
            cache_key_kpi = f'kpi_tickets_{periode}_{hash(str(df_tickets_processed.shape))}'
            
            if cache_key_kpi not in st.session_state:
                with st.spinner("Calcul des KPI en cours..."):
                    kpis = generate_kpis(filtrer_par_periode(df_support, periode), filtrer_par_periode(df_tickets_processed, periode), 'agents_all', None)
                    st.session_state[cache_key_kpi] = kpis
                    st.success("✅ KPI calculés et mis en cache")
            else:
                kpis = st.session_state[cache_key_kpi]
                st.info("📋 Utilisation des KPI en cache")

            # --- NOUVELLE SECTION : CATÉGORISATION ENRICHIE AVEC DÉFINITIONS ---
            st.markdown("## 🤖 Catégorisation Enrichie avec Définitions (SSI/Chat uniquement)")
            st.markdown("Cette section utilise le nouveau système de catégorisation enrichi avec des définitions de contexte pour améliorer la précision. **Uniquement pour les tickets avec Pipeline = 'SSI' et Source = 'Chat'.**")
            
            # Informations sur le modèle
            model_info = get_model_info()
            
            col1, col2, col3 = st.columns(3)
            with col1:
                if model_info['enhanced_model_available']:
                    st.success("✅ Modèle enrichi disponible")
                else:
                    st.warning("⚠️ Modèle enrichi non disponible")
            
            with col2:
                if model_info['definitions_available']:
                    st.success("✅ Définitions disponibles")
                else:
                    st.warning("⚠️ Définitions manquantes")
            
            with col3:
                if model_info['enhanced_model_available'] and 'enhanced_categories_count' in model_info:
                    st.info(f"📊 {model_info['enhanced_categories_count']} catégories")
                else:
                    st.info("📊 Modèle classique")
            
            # Paramètres de catégorisation
            col1, col2 = st.columns(2)
            with col1:
                seuil_confiance = st.slider(
                    "Seuil de confiance minimum",
                    min_value=0.0,
                    max_value=1.0,
                    value=0.15,
                    step=0.05,
                    help="Score minimum pour accepter une catégorisation"
                )
            
            with col2:
                afficher_scores = st.checkbox(
                    "Afficher les scores détaillés",
                    value=False,
                    help="Afficher les scores pour chaque catégorie"
                )
            
            # Cache pour la catégorisation enrichie
            cache_key_enhanced = f'categorisation_enhanced_{periode}_{seuil_confiance}_{hash(str(df_tickets_processed.shape))}'
            
            if cache_key_enhanced not in st.session_state:
                with st.spinner("Catégorisation enrichie en cours..."):
                    try:
                        # Trouver la colonne de description
                        colonnes_description_possibles = ['Description', 'Description du ticket', 'Contenu', 'Message', 'Sujet', 'Titre', 'Résumé', 'Détails', 'Commentaires', 'Subject']
                        colonne_description = None
                        for col in colonnes_description_possibles:
                            if col in df_tickets_processed.columns:
                                colonne_description = col
                                break
                        
                        if colonne_description is None:
                            st.error("❌ Aucune colonne de description trouvée dans les données")
                            return
                        
                        # Filtrer les tickets par période
                        df_tickets_periode = filtrer_par_periode(df_tickets_processed, periode)
                        
                        # Filtrer uniquement les tickets SSI/Chat
                        df_tickets_ssi_chat = df_tickets_periode[
                            (df_tickets_periode['Pipeline'].str.lower().str.contains('ssi', na=False)) & 
                            (df_tickets_periode['Source'].str.lower().str.contains('chat', na=False))
                        ].copy()
                        
                        if df_tickets_ssi_chat.empty:
                            st.warning("⚠️ Aucun ticket SSI/Chat trouvé pour cette période")
                            return
                        
                        st.info(f"📊 {len(df_tickets_ssi_chat)} tickets SSI/Chat trouvés sur {len(df_tickets_periode)} tickets totaux")
                        
                        # Prédictions avec le modèle enrichi (uniquement pour SSI/Chat)
                        descriptions = df_tickets_ssi_chat[colonne_description].fillna('').astype(str).tolist()
                        predictions = predict_with_enhanced_model(descriptions)
                        
                        # Créer le DataFrame de résultats (uniquement pour SSI/Chat)
                        df_categorise = df_tickets_ssi_chat.copy()
                        
                        # Extraire les prédictions avec gestion d'erreur
                        categories_predites = []
                        confiances = []
                        scores_par_categorie = []
                        
                        for pred in predictions:
                            # Catégorie prédite
                            if 'categorie_predite' in pred:
                                categories_predites.append(pred['categorie_predite'])
                            else:
                                categories_predites.append('Non catégorisé')
                            
                            # Confiance
                            if 'confiance' in pred:
                                confiances.append(pred['confiance'])
                            else:
                                confiances.append(0.0)
                            
                            # Scores par catégorie
                            if 'scores_par_categorie' in pred:
                                scores_par_categorie.append(pred['scores_par_categorie'])
                            else:
                                scores_par_categorie.append({})
                        
                        df_categorise['Categorie_Enrichie'] = categories_predites
                        df_categorise['Confiance_Enrichie'] = confiances
                        df_categorise['Scores_par_categorie'] = scores_par_categorie
                        
                        # Appliquer le seuil de confiance
                        df_categorise['Categorie_Final'] = df_categorise.apply(
                            lambda row: row['Categorie_Enrichie'] if row['Confiance_Enrichie'] >= seuil_confiance else 'Non catégorisé',
                            axis=1
                        )
                        
                        st.session_state[cache_key_enhanced] = df_categorise
                        st.success("✅ Catégorisation enrichie terminée et mise en cache")
                        
                    except Exception as e:
                        st.error(f"❌ Erreur lors de la catégorisation : {str(e)}")
                        st.exception(e)
            else:
                df_categorise = st.session_state[cache_key_enhanced]
                st.info("📋 Utilisation de la catégorisation enrichie en cache")
            
            # Statistiques de catégorisation
            if 'df_categorise' in locals() and df_categorise is not None:
                try:
                    total_tickets_ssi_chat = len(df_categorise)
                    
                    # Vérifier que la colonne Categorie_Final existe
                    if 'Categorie_Final' in df_categorise.columns:
                        tickets_categorises = len(df_categorise[df_categorise['Categorie_Final'] != 'Non catégorisé'])
                    else:
                        tickets_categorises = 0
                        st.warning("⚠️ Colonne 'Categorie_Final' manquante")
                    
                    taux_categorisation = (tickets_categorises / total_tickets_ssi_chat) * 100 if total_tickets_ssi_chat > 0 else 0
                    
                    # Vérifier que la colonne Confiance_Enrichie existe
                    if 'Confiance_Enrichie' in df_categorise.columns:
                        confiance_moyenne = df_categorise['Confiance_Enrichie'].mean()
                    else:
                        confiance_moyenne = 0.0
                        st.warning("⚠️ Colonne 'Confiance_Enrichie' manquante")
                except Exception as e:
                    st.error(f"❌ Erreur lors du calcul des statistiques: {str(e)}")
                    total_tickets_ssi_chat = 0
                    tickets_categorises = 0
                    taux_categorisation = 0.0
                    confiance_moyenne = 0.0
                
                col1, col2, col3, col4 = st.columns(4)
                col1.metric("Tickets SSI/Chat", total_tickets_ssi_chat)
                col2.metric("Tickets catégorisés", tickets_categorises)
                col3.metric("Taux de catégorisation", f"{taux_categorisation:.1f}%")
                col4.metric("Confiance moyenne", f"{confiance_moyenne:.3f}")
                
                # Graphique des catégories
                if 'df_categorise' in locals() and df_categorise is not None and 'Categorie_Final' in df_categorise.columns:
                    import plotly.express as px
                    stats_categories = df_categorise['Categorie_Final'].value_counts()
                    
                    # Retirer la catégorie "Non catégorisé" du graphique
                    if 'Non catégorisé' in stats_categories.index:
                        stats_categories = stats_categories.drop('Non catégorisé')
                    
                    if not stats_categories.empty:
                        fig_categories = px.pie(
                            values=stats_categories.values,
                            names=stats_categories.index,
                            title=f"Répartition des tickets SSI/Chat par catégorie (Seuil: {seuil_confiance})"
                        )
                        fig_categories.update_layout(template="plotly_dark")
                        st.plotly_chart(fig_categories, use_container_width=True)
                    else:
                        st.warning("⚠️ Aucune catégorie valide à afficher (tous les tickets sont 'Non catégorisé')")
                else:
                    st.warning("⚠️ Impossible d'afficher le graphique : données de catégorisation manquantes")
                
                # Bouton d'export PowerPoint pour le graphique
                if st.button("📊 Exporter ce graphique en PowerPoint", key="btn_categories_enhanced"):
                    pptx_io = create_single_graph_powerpoint(fig_categories, "Catégories de tickets enrichies", periode)
                    st.download_button(
                        label="📥 Télécharger le graphique Catégories Enrichies (PowerPoint)",
                        data=pptx_io,
                        file_name=f"categories_enhanced_{periode.replace(' ', '_')}.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                    )
                
                # Tableau des résultats
                if 'df_categorise' in locals() and df_categorise is not None:
                    st.markdown("### 📊 Résultats de la catégorisation enrichie")
                    
                    # Vérifier que les colonnes nécessaires existent
                    colonnes_requises = ['Ticket ID', 'Categorie_Final', 'Confiance_Enrichie']
                    colonnes_disponibles = [col for col in colonnes_requises if col in df_categorise.columns]
                    
                    if len(colonnes_disponibles) == len(colonnes_requises):
                        # Préparer le tableau de résultats
                        df_resultats = df_categorise[colonnes_disponibles].copy()
                        df_resultats = df_resultats.sort_values('Confiance_Enrichie', ascending=False)
                    else:
                        st.error(f"❌ Colonnes manquantes : {[col for col in colonnes_requises if col not in df_categorise.columns]}")
                        df_resultats = None
                else:
                    st.warning("⚠️ Aucune donnée de catégorisation disponible")
                    df_resultats = None
                
                # Afficher les scores détaillés si demandé
                if afficher_scores and 'Scores_par_categorie' in df_categorise.columns:
                    st.markdown("#### 🔍 Scores détaillés par catégorie")
                    # Prendre un échantillon pour l'affichage
                    sample_size = min(10, len(df_categorise))
                    df_sample = df_categorise.head(sample_size)
                    
                    for idx, row in df_sample.iterrows():
                        with st.expander(f"Ticket {row['Ticket ID']} - {row['Categorie_Final']} (Confiance: {row['Confiance_Enrichie']:.3f})"):
                            scores = row['Scores_par_categorie']
                            if scores and isinstance(scores, dict):
                                # Trier les scores par ordre décroissant
                                scores_tries = sorted(scores.items(), key=lambda x: x[1], reverse=True)
                                
                                for categorie, score in scores_tries[:5]:  # Top 5
                                    st.write(f"**{categorie}**: {score:.3f}")
                            else:
                                st.write("Aucun score détaillé disponible")
                
                # Afficher le tableau principal
                if df_resultats is not None and not df_resultats.empty:
                    st.dataframe(
                        df_resultats,
                        use_container_width=True,
                        hide_index=True,
                        column_config={
                            "Ticket ID": st.column_config.TextColumn("Ticket ID", width="medium"),
                            "Categorie_Final": st.column_config.TextColumn("Catégorie", width="medium"),
                            "Confiance_Enrichie": st.column_config.NumberColumn("Confiance", format="%.3f")
                        }
                    )
                else:
                    st.info("📋 Aucun résultat à afficher")
                
                # Bouton d'export Excel pour les résultats
                if st.button("📥 Exporter les résultats en Excel", key="btn_export_enhanced"):
                    try:
                        # Créer un fichier Excel avec les résultats
                        output_path = f"categorisation_enhanced_{pd.Timestamp.now().strftime('%Y%m%d_%H%M%S')}.xlsx"
                        
                        with pd.ExcelWriter(output_path, engine='openpyxl') as writer:
                            # Onglet principal
                            df_categorise.to_excel(writer, sheet_name='Résultats', index=False)
                            
                            # Onglet statistiques - recalculer les statistiques
                            if 'df_categorise' in locals() and df_categorise is not None and 'Categorie_Final' in df_categorise.columns:
                                stats_categories_export = df_categorise['Categorie_Final'].value_counts()
                                # Retirer "Non catégorisé" pour les statistiques d'export
                                if 'Non catégorisé' in stats_categories_export.index:
                                    stats_categories_export = stats_categories_export.drop('Non catégorisé')
                                
                                stats_df = pd.DataFrame({
                                    'Catégorie': stats_categories_export.index,
                                    'Nombre de tickets': stats_categories_export.values,
                                    'Pourcentage': (stats_categories_export.values / total_tickets_ssi_chat) * 100
                                })
                                stats_df.to_excel(writer, sheet_name='Statistiques', index=False)
                            else:
                                # Créer un DataFrame vide si pas de données
                                stats_df = pd.DataFrame({
                                    'Catégorie': [],
                                    'Nombre de tickets': [],
                                    'Pourcentage': []
                                })
                                stats_df.to_excel(writer, sheet_name='Statistiques', index=False)
                        
                        # Lire le fichier pour le téléchargement
                        with open(output_path, 'rb') as f:
                            excel_data = f.read()
                        
                        st.download_button(
                            label="📥 Télécharger les résultats (Excel)",
                            data=excel_data,
                            file_name=output_path,
                            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
                        )
                        
                        # Nettoyer le fichier temporaire
                        import os
                        if os.path.exists(output_path):
                            os.remove(output_path)
                            
                    except Exception as e:
                        st.error(f"❌ Erreur lors de l'export : {str(e)}")
            
            # Bouton pour vider le cache
            if st.button("🔄 Recalculer la catégorisation", key="btn_recalcul_enhanced"):
                if cache_key_enhanced in st.session_state:
                    del st.session_state[cache_key_enhanced]
                st.rerun()
            
            st.markdown("---")

            # --- SECTION : CATÉGORISATION CLASSIQUE ---
            st.markdown("## 🔧 Catégorisation Classique (SSI/Chat uniquement)")
            st.markdown("Cette section utilise le modèle de catégorisation classique (sans définitions enrichies) pour comparaison. **Uniquement pour les tickets avec Pipeline = 'SSI' et Source = 'Chat'.**")
            
            # Informations sur le modèle classique
            col1, col2, col3 = st.columns(3)
            with col1:
                if model_info['classic_model_available']:
                    st.success("✅ Modèle classique disponible")
                else:
                    st.warning("⚠️ Modèle classique non disponible")
            
            with col2:
                st.info("🔧 Modèle standard")
            
            with col3:
                st.info("📊 Comparaison")
            
            # Cache pour la catégorisation classique
            cache_key_classic = f'categorisation_classic_{periode}_{seuil_confiance}_{hash(str(df_tickets_processed.shape))}'
            
            if cache_key_classic not in st.session_state:
                with st.spinner("Catégorisation classique en cours..."):
                    try:
                        # Utiliser les mêmes données filtrées SSI/Chat
                        if 'df_tickets_ssi_chat' in locals() and df_tickets_ssi_chat is not None:
                            # Prédictions avec le modèle classique
                            from data_processing.ticket_classifier import load_ticket_classifier
                            from data_processing.enhanced_ticket_classifier import load_enhanced_ticket_classifier
                            classifier_classic = load_ticket_classifier()
                            predictions_classic = classifier_classic.predict(descriptions)
                            
                            # Créer le DataFrame de résultats
                            df_categorise_classic = df_tickets_ssi_chat.copy()
                            
                            # Traitement sécurisé des prédictions
                            categories_predites_classic = []
                            confiances_classic = []
                            scores_par_categorie_classic = []
                            
                            for pred in predictions_classic:
                                # Catégorie prédite
                                if 'categorie_predite' in pred:
                                    categories_predites_classic.append(pred['categorie_predite'])
                                else:
                                    categories_predites_classic.append('Non catégorisé')
                                
                                # Confiance
                                if 'confiance' in pred:
                                    confiances_classic.append(pred['confiance'])
                                else:
                                    confiances_classic.append(0.0)
                                
                                # Scores par catégorie
                                if 'scores_par_categorie' in pred:
                                    scores_par_categorie_classic.append(pred['scores_par_categorie'])
                                else:
                                    scores_par_categorie_classic.append({})
                            
                            df_categorise_classic['Categorie_Classique'] = categories_predites_classic
                            df_categorise_classic['Confiance_Classique'] = confiances_classic
                            df_categorise_classic['Scores_par_categorie_Classique'] = scores_par_categorie_classic
                            
                            # Appliquer le seuil de confiance
                            df_categorise_classic['Categorie_Final_Classique'] = df_categorise_classic.apply(
                                lambda row: row['Categorie_Classique'] if row['Confiance_Classique'] >= seuil_confiance else 'Non catégorisé',
                                axis=1
                            )
                            
                            st.session_state[cache_key_classic] = df_categorise_classic
                            st.success("✅ Catégorisation classique terminée et mise en cache")
                        else:
                            st.warning("⚠️ Aucune donnée SSI/Chat disponible pour la catégorisation classique")
                            
                    except Exception as e:
                        st.error(f"❌ Erreur lors de la catégorisation classique : {str(e)}")
                        st.exception(e)
            else:
                df_categorise_classic = st.session_state[cache_key_classic]
                st.info("📋 Utilisation de la catégorisation classique en cache")
            
            # Statistiques de catégorisation classique
            if 'df_categorise_classic' in locals() and df_categorise_classic is not None:
                try:
                    total_tickets_ssi_chat_classic = len(df_categorise_classic)
                    
                    if 'Categorie_Final_Classique' in df_categorise_classic.columns:
                        tickets_categorises_classic = len(df_categorise_classic[df_categorise_classic['Categorie_Final_Classique'] != 'Non catégorisé'])
                    else:
                        tickets_categorises_classic = 0
                    
                    taux_categorisation_classic = (tickets_categorises_classic / total_tickets_ssi_chat_classic) * 100 if total_tickets_ssi_chat_classic > 0 else 0
                    
                    if 'Confiance_Classique' in df_categorise_classic.columns:
                        confiance_moyenne_classic = df_categorise_classic['Confiance_Classique'].mean()
                    else:
                        confiance_moyenne_classic = 0.0
                        
                except Exception as e:
                    st.error(f"❌ Erreur lors du calcul des statistiques classiques: {str(e)}")
                    total_tickets_ssi_chat_classic = 0
                    tickets_categorises_classic = 0
                    taux_categorisation_classic = 0.0
                    confiance_moyenne_classic = 0.0
                
                col1, col2, col3, col4 = st.columns(4)
                col1.metric("Tickets SSI/Chat", total_tickets_ssi_chat_classic)
                col2.metric("Tickets catégorisés", tickets_categorises_classic)
                col3.metric("Taux de catégorisation", f"{taux_categorisation_classic:.1f}%")
                col4.metric("Confiance moyenne", f"{confiance_moyenne_classic:.3f}")
                
                # Graphique des catégories classiques
                if 'Categorie_Final_Classique' in df_categorise_classic.columns:
                    stats_categories_classic = df_categorise_classic['Categorie_Final_Classique'].value_counts()
                    
                    if 'Non catégorisé' in stats_categories_classic.index:
                        stats_categories_classic = stats_categories_classic.drop('Non catégorisé')
                    
                    if not stats_categories_classic.empty:
                        fig_categories_classic = px.pie(
                            values=stats_categories_classic.values,
                            names=stats_categories_classic.index,
                            title=f"Répartition des tickets SSI/Chat par catégorie - Modèle Classique (Seuil: {seuil_confiance})"
                        )
                        fig_categories_classic.update_layout(template="plotly_dark")
                        st.plotly_chart(fig_categories_classic, use_container_width=True)
                    else:
                        st.warning("⚠️ Aucune catégorie valide à afficher (tous les tickets sont 'Non catégorisé')")
            
            # Bouton pour vider le cache classique
            if st.button("🔄 Recalculer la catégorisation classique", key="btn_recalcul_classic"):
                if cache_key_classic in st.session_state:
                    del st.session_state[cache_key_classic]
                st.rerun()
            
            st.markdown("---")

            # --- SECTION : CATÉGORISATION HUBSPOT (SANS ENTRAÎNEMENT) ---
            st.markdown("## 📊 Catégorisation HubSpot (SSI/Chat uniquement)")
            st.markdown("Cette section utilise directement le champ catégorie de HubSpot sans entraînement de modèle. **Uniquement pour les tickets avec Pipeline = 'SSI' et Source = 'Chat'.**")
            
            # Informations sur la catégorisation HubSpot
            col1, col2, col3 = st.columns(3)
            with col1:
                st.success("✅ Données HubSpot disponibles")
            
            with col2:
                st.info("📊 Catégories natives")
            
            with col3:
                st.info("🔍 Analyse directe")
            
            # Cache pour la catégorisation HubSpot
            cache_key_hubspot = f'categorisation_hubspot_{periode}_{hash(str(df_tickets_processed.shape))}'
            
            if cache_key_hubspot not in st.session_state:
                with st.spinner("Analyse des catégories HubSpot en cours..."):
                    try:
                        # Utiliser les mêmes données filtrées SSI/Chat
                        if 'df_tickets_ssi_chat' in locals() and df_tickets_ssi_chat is not None:
                            # Créer le DataFrame de résultats HubSpot
                            df_categorise_hubspot = df_tickets_ssi_chat.copy()
                            
                            # Chercher la colonne catégorie dans les données HubSpot
                            colonnes_categorie_possibles = ['Catégorie', 'Category', 'Categorie', 'Type', 'Type de ticket', 'Ticket Type', 'Classification', 'Classe']
                            colonne_categorie_hubspot = None
                            for col in colonnes_categorie_possibles:
                                if col in df_categorise_hubspot.columns:
                                    colonne_categorie_hubspot = col
                                    break
                            
                            if colonne_categorie_hubspot is not None:
                                # Utiliser directement la catégorie HubSpot
                                df_categorise_hubspot['Categorie_HubSpot'] = df_categorise_hubspot[colonne_categorie_hubspot].fillna('Non catégorisé')
                                
                                # Calculer le top 20 des catégories les plus fréquentes
                                categories_counts = df_categorise_hubspot['Categorie_HubSpot'].value_counts()
                                top_20_categories = categories_counts.head(20).index.tolist()
                                
                                # Appliquer le filtrage : garder seulement les top 20, les autres deviennent "Autres"
                                df_categorise_hubspot['Categorie_Final_HubSpot'] = df_categorise_hubspot['Categorie_HubSpot'].apply(
                                    lambda x: x if x in top_20_categories else 'Autres'
                                )
                                df_categorise_hubspot['Confiance_HubSpot'] = 1.0  # Confiance maximale car c'est la vraie catégorie
                                
                                st.success(f"✅ Catégorisation HubSpot terminée (colonne: {colonne_categorie_hubspot})")
                                st.info(f"📊 Top 20 des catégories HubSpot appliqué ({len(top_20_categories)} catégories principales)")
                            else:
                                # Si pas de colonne catégorie trouvée
                                df_categorise_hubspot['Categorie_HubSpot'] = 'Non catégorisé'
                                df_categorise_hubspot['Categorie_Final_HubSpot'] = 'Non catégorisé'
                                df_categorise_hubspot['Confiance_HubSpot'] = 0.0
                                
                                st.warning("⚠️ Aucune colonne catégorie trouvée dans les données HubSpot")
                            
                            st.session_state[cache_key_hubspot] = df_categorise_hubspot
                            
                    except Exception as e:
                        st.error(f"❌ Erreur lors de l'analyse HubSpot : {str(e)}")
                        st.exception(e)
            else:
                df_categorise_hubspot = st.session_state[cache_key_hubspot]
                st.info("📋 Utilisation de la catégorisation HubSpot en cache")
            
            # Statistiques de catégorisation HubSpot
            if 'df_categorise_hubspot' in locals() and df_categorise_hubspot is not None:
                try:
                    total_tickets_ssi_chat_hubspot = len(df_categorise_hubspot)
                    
                    if 'Categorie_Final_HubSpot' in df_categorise_hubspot.columns:
                        tickets_categorises_hubspot = len(df_categorise_hubspot[df_categorise_hubspot['Categorie_Final_HubSpot'] != 'Non catégorisé'])
                    else:
                        tickets_categorises_hubspot = 0
                    
                    taux_categorisation_hubspot = (tickets_categorises_hubspot / total_tickets_ssi_chat_hubspot) * 100 if total_tickets_ssi_chat_hubspot > 0 else 0
                    
                    if 'Confiance_HubSpot' in df_categorise_hubspot.columns:
                        confiance_moyenne_hubspot = df_categorise_hubspot['Confiance_HubSpot'].mean()
                    else:
                        confiance_moyenne_hubspot = 0.0
                        
                except Exception as e:
                    st.error(f"❌ Erreur lors du calcul des statistiques HubSpot: {str(e)}")
                    total_tickets_ssi_chat_hubspot = 0
                    tickets_categorises_hubspot = 0
                    taux_categorisation_hubspot = 0.0
                    confiance_moyenne_hubspot = 0.0
                
                col1, col2, col3, col4 = st.columns(4)
                col1.metric("Tickets SSI/Chat", total_tickets_ssi_chat_hubspot)
                col2.metric("Tickets catégorisés", tickets_categorises_hubspot)
                col3.metric("Taux de catégorisation", f"{taux_categorisation_hubspot:.1f}%")
                col4.metric("Confiance moyenne", f"{confiance_moyenne_hubspot:.3f}")
                
                # Graphique des catégories HubSpot
                if 'Categorie_Final_HubSpot' in df_categorise_hubspot.columns:
                    stats_categories_hubspot = df_categorise_hubspot['Categorie_Final_HubSpot'].value_counts()
                    
                    if 'Non catégorisé' in stats_categories_hubspot.index:
                        stats_categories_hubspot = stats_categories_hubspot.drop('Non catégorisé')
                    
                    # Exclure "Autres" du graphique si elle représente moins de 5% des tickets
                    if 'Autres' in stats_categories_hubspot.index:
                        total_tickets_graph = stats_categories_hubspot.sum()
                        autres_percentage = (stats_categories_hubspot['Autres'] / total_tickets_graph) * 100
                        if autres_percentage < 5.0:
                            stats_categories_hubspot = stats_categories_hubspot.drop('Autres')
                            st.info(f"ℹ️ Catégorie 'Autres' exclue du graphique (représente {autres_percentage:.1f}% des tickets)")
                    
                    if not stats_categories_hubspot.empty:
                        fig_categories_hubspot = px.pie(
                            values=stats_categories_hubspot.values,
                            names=stats_categories_hubspot.index,
                            title=f"Répartition des tickets SSI/Chat par catégorie - HubSpot (Top 20)"
                        )
                        fig_categories_hubspot.update_layout(template="plotly_dark")
                        st.plotly_chart(fig_categories_hubspot, use_container_width=True)
                    else:
                        st.warning("⚠️ Aucune catégorie valide à afficher (tous les tickets sont 'Non catégorisé')")
            
            # Bouton pour vider le cache HubSpot
            if st.button("🔄 Recalculer la catégorisation HubSpot", key="btn_recalcul_hubspot"):
                if cache_key_hubspot in st.session_state:
                    del st.session_state[cache_key_hubspot]
                st.rerun()
            
            st.markdown("---")

            # --- SECTION : COMPARAISON DES MÉTHODES ---
            st.markdown("## 📈 Comparaison des Méthodes de Catégorisation")
            st.markdown("Comparaison des performances des trois méthodes de catégorisation pour les tickets SSI/Chat.")
            
            # Créer un DataFrame de comparaison
            comparison_data = []
            
            # Méthode 1 : Enrichie
            if 'df_categorise' in locals() and df_categorise is not None:
                comparison_data.append({
                    'Méthode': 'Enrichie avec Définitions',
                    'Tickets SSI/Chat': total_tickets_ssi_chat,
                    'Tickets catégorisés': tickets_categorises,
                    'Taux de catégorisation': f"{taux_categorisation:.1f}%",
                    'Confiance moyenne': f"{confiance_moyenne:.3f}",
                    'Seuil de confiance': seuil_confiance
                })
            
            # Méthode 2 : Classique
            if 'df_categorise_classic' in locals() and df_categorise_classic is not None:
                comparison_data.append({
                    'Méthode': 'Modèle Classique',
                    'Tickets SSI/Chat': total_tickets_ssi_chat_classic,
                    'Tickets catégorisés': tickets_categorises_classic,
                    'Taux de catégorisation': f"{taux_categorisation_classic:.1f}%",
                    'Confiance moyenne': f"{confiance_moyenne_classic:.3f}",
                    'Seuil de confiance': seuil_confiance
                })
            
            # Méthode 3 : HubSpot
            if 'df_categorise_hubspot' in locals() and df_categorise_hubspot is not None:
                comparison_data.append({
                    'Méthode': 'HubSpot (Référence)',
                    'Tickets SSI/Chat': total_tickets_ssi_chat_hubspot,
                    'Tickets catégorisés': tickets_categorises_hubspot,
                    'Taux de catégorisation': f"{taux_categorisation_hubspot:.1f}%",
                    'Confiance moyenne': f"{confiance_moyenne_hubspot:.3f}",
                    'Seuil de confiance': 'N/A'
                })
            
            if comparison_data:
                df_comparison = pd.DataFrame(comparison_data)
                st.dataframe(
                    df_comparison,
                    use_container_width=True,
                    hide_index=True,
                    column_config={
                        "Méthode": st.column_config.TextColumn("Méthode", width="medium"),
                        "Tickets SSI/Chat": st.column_config.NumberColumn("Tickets SSI/Chat", format="%d"),
                        "Tickets catégorisés": st.column_config.NumberColumn("Tickets catégorisés", format="%d"),
                        "Taux de catégorisation": st.column_config.TextColumn("Taux de catégorisation", width="medium"),
                        "Confiance moyenne": st.column_config.TextColumn("Confiance moyenne", width="medium"),
                        "Seuil de confiance": st.column_config.TextColumn("Seuil de confiance", width="medium")
                    }
                )
                
                # Graphique de comparaison des taux de catégorisation
                if len(comparison_data) > 1:
                    fig_comparison = px.bar(
                        df_comparison,
                        x='Méthode',
                        y='Tickets catégorisés',
                        title="Comparaison du nombre de tickets catégorisés par méthode",
                        color='Méthode',
                        text='Tickets catégorisés'
                    )
                    fig_comparison.update_layout(template="plotly_dark")
                    fig_comparison.update_traces(textposition='outside')
                    st.plotly_chart(fig_comparison, use_container_width=True)
            else:
                st.warning("⚠️ Aucune donnée de comparaison disponible")
            
            st.markdown("---")

            st.plotly_chart(kpis['fig_activite_ticket'])

            # --- SECTION KPI STANDARD ---
            st.markdown("## 📊 KPI Standard")
            
            st.plotly_chart(kpis['activite_ticket_semaine'])
            st.plotly_chart(kpis['activite_categorie'])

            partenaire = st.selectbox("Sélectionnez un partenaire :", ["FOLLOW", "Odaiji", "Help Info", "Oppysoft"], key="unique_partenaire_selection")

            # Cache pour les KPI partenaire
            cache_key_kpi_partenaire = f'kpi_partenaire_{partenaire}_{periode}_{hash(str(df_tickets_processed.shape))}'
            
            if cache_key_kpi_partenaire not in st.session_state:
                with st.spinner(f"Calcul des KPI pour {partenaire} en cours..."):
                    kpis_partenaire = generate_kpis(filtrer_par_periode(df_support, periode), filtrer_par_periode(df_tickets_processed, periode), 'agents_all', partenaire)
                    st.session_state[cache_key_kpi_partenaire] = kpis_partenaire
                    st.success(f"✅ KPI {partenaire} calculés et mis en cache")
            else:
                kpis_partenaire = st.session_state[cache_key_kpi_partenaire]
                st.info(f"📋 Utilisation des KPI {partenaire} en cache")

            col1, col2, col3= st.columns(3) 
            #col1.metric(f'% sla < 2h',int(kpis_partenaire['%_sla_2']))
            col2.metric(f'Délai moyen de réponse (heure)',round(kpis_partenaire['delai_moyen_reponse'], 2))

            st.plotly_chart(kpis_partenaire['sla_fig'])
            st.plotly_chart(kpis_partenaire['activite_categorie'])

        # --- Bloc Tickets simplifié (partenaires uniquement) ---
        def tickets():
            periode = st.selectbox(
                "Sélectionnez une période :",
                ["1 an", "Toute la période", "6 derniers mois", "3 derniers mois", "Dernier mois"],
                index=0
            )

            df_tickets_processed = st.session_state['df_tickets_processed']
            df_tickets_periode = filtrer_par_periode(df_tickets_processed, periode)

            st.markdown("### 🎯 KPI Partenaires (FOLLOW, Odaiji, Help Info, Oppysoft)")

            partenaire = st.selectbox(
                "Sélectionnez un partenaire :", ["FOLLOW", "Odaiji", "Help Info", "Oppysoft"],
                key="unique_partenaire_selection"
            )

            cache_key_kpi_partenaire = f'kpi_partenaire_{partenaire}_{periode}_{hash(str(df_tickets_periode.shape))}'
            if cache_key_kpi_partenaire not in st.session_state:
                with st.spinner(f"Calcul des KPI pour {partenaire} en cours..."):
                    sla_fig, sla_2h, delai_moyen, df_partenaire, _ = sla(df_tickets_periode, partenaire, canal="partenaire")
                    fig_categories = metrics_nombre_ticket_categorie(df_tickets_periode, partenaire)
                    st.session_state[cache_key_kpi_partenaire] = {
                        "sla_fig": sla_fig,
                        "delai_moyen": delai_moyen,
                        "df_partenaire": df_partenaire,
                        "fig_categories": fig_categories,
                    }
                    st.success(f"✅ KPI {partenaire} calculés et mis en cache")

            kpi_part = st.session_state[cache_key_kpi_partenaire]
            df_partenaire = kpi_part["df_partenaire"]

            if df_partenaire.empty:
                st.warning(f"Aucun ticket pour {partenaire} sur la période sélectionnée.")
                return

            def pick_date_column(df):
                candidates = ['Date', 'Date de création', 'Created', 'Create date', 'created_at', 'Date entrée N2']
                for col in candidates:
                    if col in df.columns:
                        return col
                return None

            date_col = pick_date_column(df_partenaire)
            if date_col is None:
                st.error("Impossible de trouver une colonne de date pour calculer les métriques mensuelles.")
                return

            df_partenaire = df_partenaire.copy()
            df_partenaire[date_col] = pd.to_datetime(df_partenaire[date_col], errors='coerce')
            df_partenaire = df_partenaire.dropna(subset=[date_col])

            df_partenaire['Mois'] = df_partenaire[date_col].dt.to_period('M').astype(str)
            df_partenaire['Heure'] = df_partenaire[date_col].dt.hour

            id_col = 'Ticket ID' if 'Ticket ID' in df_partenaire.columns else df_partenaire.columns[0]
            tickets_par_mois = df_partenaire.groupby('Mois')[id_col].count().reset_index(name='Tickets')
            tickets_par_mois = tickets_par_mois.sort_values('Mois')
            moyenne_tickets_mois = tickets_par_mois['Tickets'].mean() if not tickets_par_mois.empty else 0

            resp_par_mois = None
            df_resp = None
            moyenne_reponse_mensuelle = None
            mediane_reponse_mensuelle = None
            if 'working_hours' in df_partenaire.columns:
                df_partenaire['working_hours'] = pd.to_numeric(df_partenaire['working_hours'], errors='coerce')
                df_resp = df_partenaire.dropna(subset=['working_hours'])
                if not df_resp.empty:
                    # Limiter les valeurs extrêmes pour éviter d'aplatir le graphique
                    cap = df_resp['working_hours'].quantile(0.995)
                    df_resp_viz = df_resp.copy()
                    if pd.notna(cap):
                        df_resp_viz['working_hours_capped'] = df_resp_viz['working_hours'].clip(upper=cap)
                    else:
                        df_resp_viz['working_hours_capped'] = df_resp_viz['working_hours']

                    resp_par_mois = df_resp_viz.groupby(['Mois'])['working_hours_capped'].mean().reset_index()
                    moyenne_reponse_mensuelle = resp_par_mois['working_hours_capped'].mean()
                    mediane_reponse_mensuelle = df_resp_viz['working_hours_capped'].median()

            # Graphique combiné : tickets (barres) + temps moyen de réponse (courbe) avec axe secondaire
            merged = tickets_par_mois.copy()
            if resp_par_mois is not None:
                merged = merged.merge(resp_par_mois, on='Mois', how='left')
            merged = merged.sort_values('Mois')

            fig_tickets_mois = go.Figure()
            fig_tickets_mois.add_bar(
                x=merged['Mois'],
                y=merged['Tickets'],
                name="Tickets / mois",
                marker_color="#1f77b4",
            )
            if resp_par_mois is not None:
                fig_tickets_mois.add_scatter(
                    x=merged['Mois'],
                    y=merged['working_hours_capped'],
                    name="Temps moyen de réponse (h)",
                    mode="lines+markers",
                    yaxis="y2",
                    line=dict(color="#ff7f0e"),
                )

            fig_tickets_mois.update_layout(
                template="plotly_dark",
                title="Tickets reçus par mois & temps moyen de réponse",
                yaxis=dict(title="Tickets par mois"),
                yaxis2=dict(
                    title="Temps moyen de réponse (h)",
                    overlaying="y",
                    side="right",
                ),
                legend=dict(orientation="h", yanchor="bottom", y=1.02, xanchor="right", x=1),
            )

            fig_box_resp = None
            if df_resp is not None and not df_resp.empty:
                df_box = df_resp_viz if 'df_resp_viz' in locals() else df_resp
                y_box = 'working_hours_capped' if 'working_hours_capped' in df_box.columns else 'working_hours'
                fig_box_resp = px.box(
                    df_box,
                    x='Mois',
                    y=y_box,
                    title="Distribution des temps de réponse par mois",
                    labels={y_box: 'Temps de réponse (h)', 'Mois': 'Mois'},
                    points="outliers"
                )
                fig_box_resp.update_layout(template="plotly_dark")

            col1, col2, col3 = st.columns(3)
            col1.metric("Délai moyen de réponse (heures)", f"{kpi_part['delai_moyen']:.2f}")
            col2.metric("Médiane temps de réponse (heures)", f"{(mediane_reponse_mensuelle or 0):.2f}")
            col3.metric("Tickets / mois (moyenne)", f"{moyenne_tickets_mois:.1f}")

            st.plotly_chart(fig_tickets_mois, use_container_width=True)
            if fig_box_resp is not None:
                st.plotly_chart(fig_box_resp, use_container_width=True)
            st.plotly_chart(kpi_part["fig_categories"], use_container_width=True)

        tickets()

    elif selection_page == "Analyste IA":
        from utils.mcp_ui import render_mcp_analyst_page
        render_mcp_analyst_page()

    # --- RAPPORT POWERPOINT AGENTS ---
    from utils.powerpoint_helpers import create_powerpoint_agents_report
    from data_processing.kpi_generation import charge_entrant_sortant, charge_ticket, df_compute_ticket_appels_metrics
    st.sidebar.markdown("---")
    if st.sidebar.button("Générer rapport PowerPoint Agents"):
        # Texte méthodologie scoring (markdown)
        markdown_methodo = """## Paramètres de scoring\n\n- Objectif total de demandes par jour : 25\n- Ratio d'appels : 70%\n- Ratio de tickets : 30%\n- Objectif taux de service : 70%\n\nLe score est calculé selon le volume total traité, la répartition appels/tickets, la comparaison à la moyenne et le taux d'appels entrants."""
        # Préparation des données 6 mois
        periode_6m = "6 derniers mois"
        periode_3m = "3 derniers mois"
        agents_n1 = ['Mourad HUMBLOT', 'Archimede KESSI', 'Céline Crendal', 'Melinda Marmin', 'Emilie GEST', 'Sandrine Sauvage']
        df_support_6m = filtrer_par_periode(def_df_support(process_aircall_data(df_aircall), process_aircall_data(df_aircall), line_tous, agents_all), periode_6m)
        df_tickets_6m = filtrer_par_periode(df_tickets_processed, periode_6m)
        df_scores_6m = df_compute_ticket_appels_metrics(agents_n1, df_tickets_6m, df_support_6m)
        # Préparation des données 3 mois
        df_support_3m = filtrer_par_periode(def_df_support(process_aircall_data(df_aircall), process_aircall_data(df_aircall), line_tous, agents_all), periode_3m)
        df_tickets_3m = filtrer_par_periode(df_tickets_processed, periode_3m)
        df_scores_3m = df_compute_ticket_appels_metrics(agents_n1, df_tickets_3m, df_support_3m)
        # Graphiques activité téléphone et tickets (6 mois)
        fig_line_entrant, fig_pie_entrant = charge_entrant_sortant(df_support_6m, agents_n1)
        fig_ticket_6m, _ = charge_ticket(df_tickets_6m, agents_n1)
        # Graphique temps de réponse (6 mois)
        from data_processing.kpi_generation import calculate_ticket_response_time
        moyenne_temps_reponse_6m, fig_temps_reponse_6m, _ = calculate_ticket_response_time(df_tickets_6m, agents_n1)
        # Génération du rapport
        pptx_io = create_powerpoint_agents_report(df_scores_6m, df_scores_3m, fig_pie_entrant, fig_ticket_6m, markdown_methodo, fig_temps_reponse_6m)
        st.sidebar.success("Rapport PowerPoint Agents généré !")
        st.sidebar.download_button(
            label="📥 Télécharger le rapport Agents",
            data=pptx_io,
            file_name="rapport_agents.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
    # --- FIN RAPPORT POWERPOINT AGENTS ---


def create_powerpoint(kpis, df_support, periode):
    # Créer une nouvelle présentation
    prs = Presentation()
    
    # Slide 1 - Titre
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Dashboard Support"
    subtitle.text = f"Période : {periode}"
    
    # Slide 2 - KPIs principaux
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = 'KPIs Principaux'
    
    tf = body_shape.text_frame
    tf.text = f"Taux de service : {kpis['Taux_de_service']}%"
    p = tf.add_paragraph()
    p.text = f"Appels entrants par jour : {kpis['Entrant']}"
    p = tf.add_paragraph()
    p.text = f"Numéros uniques par jour : {kpis['Numero_unique']}"
    p = tf.add_paragraph()
    p.text = f"Temps moyen par appel : {kpis['temps_moy_appel']}"
    
    # Slide 3 - Graphiques existants
    for fig in [kpis['charge_affid_stellair_%'], kpis['charge_affid_stellair_v'], kpis['evo_appels_tickets']]:
        img_slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(img_slide_layout)
        
        # Sauvegarder le graphique temporairement
        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
            fig.write_image(tmp.name)
            slide.shapes.add_picture(tmp.name, Inches(1), Inches(1), width=Inches(8))
    
    # Slide 4 - Graphique répartition des groupes d'agents
    from data_processing.kpi_generation import graph_repartition_groupes_stellair
    repartition_groupes = graph_repartition_groupes_stellair(df_support)
    
    img_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(img_slide_layout)
    title = slide.shapes.title
    title.text = "Répartition des groupes d'agents"
    
    # Sauvegarder le graphique temporairement
    with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
        repartition_groupes.write_image(tmp.name)
        slide.shapes.add_picture(tmp.name, Inches(1), Inches(1), width=Inches(8))
    
    # Sauvegarder la présentation en mémoire
    pptx_io = io.BytesIO()
    prs.save(pptx_io)
    pptx_io.seek(0)
    
    return pptx_io

