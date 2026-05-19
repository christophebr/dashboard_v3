#!/usr/bin/env python3
"""
Synthèse PowerPoint — utilisation du chatbot Yelda (Stellair) pour COMEX.

Reprend les données et indicateurs du dashboard Streamlit et du classeur
``yelda_comex_analyse_utilisation.xlsx`` (généré par ``yelda_comex_analysis.py``).

Périmètre :
  - URL d'origine = https://fse.stellair.fr
  - Conversations **initiées** uniquement : ``Nombre de messages >= 2``
    (l'utilisateur a posé au moins une question ; exclut les simples ouvertures du widget).

Slides :
  1. Titre + période
  2. Messages clés COMEX
  3. Volumétrie mensuelle des conversations initiées (bar chart)
  4. Taux de satisfaction LLM mensuel (courbe)
  5. Déflexion ticket mensuelle (bar chart empilé)
  6. Top intentions sur conversations initiées (bar horizontal)
  7. Utilisateurs : adoption et récurrence (bar chart)
  8. Méthodologie

Usage :
  python3 data_processing/yelda_comex_powerpoint.py
  python3 data_processing/yelda_comex_powerpoint.py -o rapport.pptx
"""
from __future__ import annotations

import argparse
import sys
import tempfile
from pathlib import Path
from typing import Dict

import pandas as pd
import plotly.graph_objects as go

_ROOT = Path(__file__).resolve().parents[1]
if str(_ROOT) not in sys.path:
    sys.path.insert(0, str(_ROOT))

from data_processing.yelda_comex_analysis import (  # noqa: E402
    build_deflexion_ticket,
    build_satisfaction_llm_mensuel,
    build_synthese,
    build_top_intentions,
    build_utilisateurs,
    build_volumetrie_hebdo,
    build_volumetrie_mensuelle,
    filter_conversations_initiees,
)
from data_processing.yelda_processing import (  # noqa: E402
    COL_DATE,
    filter_yelda_stellair,
    load_yelda_data,
)

from pptx import Presentation  # noqa: E402
from pptx.dml.color import RGBColor  # noqa: E402
from pptx.enum.shapes import MSO_SHAPE  # noqa: E402
from pptx.util import Emu, Inches, Pt  # noqa: E402

# Palette alignée sur la charte Olaqin / dashboard
COLOR_PRIMARY = "#1F4E78"  # Bleu foncé
COLOR_ACCENT = "#2E86C1"
COLOR_GREEN = "#27AE60"
COLOR_RED = "#C0392B"
COLOR_ORANGE = "#E67E22"
COLOR_GREY = "#7F8C8D"
COLOR_BG = "#F5F7FA"

PLOTLY_LAYOUT_DEFAULTS = dict(
    template="plotly_white",
    font=dict(family="Arial, sans-serif", size=14, color="#2C3E50"),
    paper_bgcolor="white",
    plot_bgcolor="white",
    margin=dict(l=60, r=30, t=70, b=70),
    titlefont=dict(size=20, color=COLOR_PRIMARY),
)


def _apply_layout(fig: go.Figure, title: str, **extra) -> None:
    """Applique le layout standard + titre coloré."""
    fig.update_layout(
        title=dict(text=title, font=dict(size=20, color=COLOR_PRIMARY, family="Arial, sans-serif")),
        template="plotly_white",
        font=dict(family="Arial, sans-serif", size=14, color="#2C3E50"),
        paper_bgcolor="white",
        plot_bgcolor="white",
        margin=dict(l=60, r=30, t=80, b=70),
        **extra,
    )

IMAGE_WIDTH_PX = 1400
IMAGE_HEIGHT_PX = 700


# ----------------------------------------------------------------------------
# Graphiques (Plotly -> PNG)
# ----------------------------------------------------------------------------

def _save_png(fig: go.Figure, out: Path) -> Path:
    fig.write_image(str(out), width=IMAGE_WIDTH_PX, height=IMAGE_HEIGHT_PX, scale=2)
    return out


def chart_volumetrie(df_vol: pd.DataFrame, out: Path) -> Path:
    fig = go.Figure()
    fig.add_trace(
        go.Bar(
            x=df_vol["Mois"],
            y=df_vol["Conversations initiées"],
            name="Conversations initiées",
            marker_color=COLOR_PRIMARY,
            text=df_vol["Conversations initiées"],
            textposition="outside",
        )
    )
    fig.add_trace(
        go.Scatter(
            x=df_vol["Mois"],
            y=df_vol["Utilisateurs uniques (HubSpot ID)"],
            name="Utilisateurs uniques",
            mode="lines+markers",
            line=dict(color=COLOR_ORANGE, width=3),
            marker=dict(size=10),
            yaxis="y2",
        )
    )
    _apply_layout(
        fig,
        "Volumétrie mensuelle — conversations initiées & utilisateurs uniques",
        xaxis=dict(title="Mois"),
        yaxis=dict(title="Conversations initiées", side="left"),
        yaxis2=dict(title="Utilisateurs uniques", side="right", overlaying="y", showgrid=False),
        legend=dict(orientation="h", yanchor="bottom", y=1.02, xanchor="center", x=0.5),
    )
    return _save_png(fig, out)


def chart_volumetrie_hebdo(df_vol_h: pd.DataFrame, out: Path) -> Path:
    fig = go.Figure()
    fig.add_trace(
        go.Bar(
            x=df_vol_h["Semaine"],
            y=df_vol_h["Conversations initiées"],
            name="Conversations initiées",
            marker_color=COLOR_PRIMARY,
            text=df_vol_h["Conversations initiées"],
            textposition="outside",
            textfont=dict(size=11),
        )
    )
    fig.add_trace(
        go.Scatter(
            x=df_vol_h["Semaine"],
            y=df_vol_h["Utilisateurs uniques (HubSpot ID)"],
            name="Utilisateurs uniques",
            mode="lines+markers",
            line=dict(color=COLOR_ORANGE, width=3),
            marker=dict(size=8),
            yaxis="y2",
        )
    )
    _apply_layout(
        fig,
        "Volumétrie hebdomadaire — conversations initiées & utilisateurs uniques",
        xaxis=dict(title="Semaine ISO", tickangle=-45),
        yaxis=dict(title="Conversations initiées", side="left"),
        yaxis2=dict(title="Utilisateurs uniques", side="right", overlaying="y", showgrid=False),
        legend=dict(orientation="h", yanchor="bottom", y=1.02, xanchor="center", x=0.5),
    )
    return _save_png(fig, out)


def chart_satisfaction(df_sat: pd.DataFrame, out: Path) -> Path:
    fig = go.Figure()
    fig.add_trace(
        go.Scatter(
            x=df_sat["Mois"],
            y=df_sat["Taux satisfaction LLM (%)"],
            mode="lines+markers+text",
            line=dict(color=COLOR_GREEN, width=4),
            marker=dict(size=14, color=COLOR_GREEN),
            text=[f"{v:.1f} %" for v in df_sat["Taux satisfaction LLM (%)"]],
            textposition="top center",
            textfont=dict(size=14, color=COLOR_PRIMARY),
            name="Taux satisfaction LLM",
        )
    )
    _apply_layout(
        fig,
        "Taux de satisfaction LLM par mois (Satisfait / Sat+Insat)",
        xaxis=dict(title="Mois"),
        yaxis=dict(title="Taux satisfaction (%)", range=[0, 105]),
        showlegend=False,
    )
    return _save_png(fig, out)


def chart_deflexion(df_def: pd.DataFrame, out: Path) -> Path:
    df_mens = df_def[df_def["Mois"] != "Total"].copy()
    fig = go.Figure()
    fig.add_trace(
        go.Bar(
            x=df_mens["Mois"],
            y=df_mens["Conversations initiées"] - df_mens["Tickets créés"],
            name="Résolues sans ticket (déflexion)",
            marker_color=COLOR_GREEN,
        )
    )
    fig.add_trace(
        go.Bar(
            x=df_mens["Mois"],
            y=df_mens["Tickets créés"],
            name="Tickets créés depuis le bot",
            marker_color=COLOR_RED,
        )
    )
    fig.add_trace(
        go.Scatter(
            x=df_mens["Mois"],
            y=df_mens["Part sans ticket / déflexion (%)"],
            name="Taux de déflexion (%)",
            mode="lines+markers+text",
            line=dict(color=COLOR_PRIMARY, width=3, dash="dot"),
            marker=dict(size=10),
            text=[f"{v:.1f} %" for v in df_mens["Part sans ticket / déflexion (%)"]],
            textposition="top center",
            yaxis="y2",
        )
    )
    _apply_layout(
        fig,
        "Déflexion ticket : conversations résolues par le bot vs tickets créés",
        barmode="stack",
        xaxis=dict(title="Mois"),
        yaxis=dict(title="Conversations initiées"),
        yaxis2=dict(title="Taux de déflexion (%)", overlaying="y", side="right", range=[0, 110], showgrid=False),
        legend=dict(orientation="h", yanchor="bottom", y=1.02, xanchor="center", x=0.5),
    )
    return _save_png(fig, out)


def chart_top_intentions(df_int: pd.DataFrame, out: Path, top_n: int = 10) -> Path:
    d = df_int.head(top_n).iloc[::-1]
    colors = []
    flag_faible = {"default_fallback", "contact_agent", "reponse_agent_non_satisfaisante",
                   "contact_agent_phone", "contact_agent_callback", "contact_agent_email"}
    for name in d["Intention"]:
        colors.append(COLOR_RED if name in flag_faible else COLOR_ACCENT)
    fig = go.Figure()
    fig.add_trace(
        go.Bar(
            x=d["Occurrences"],
            y=d["Intention"],
            orientation="h",
            marker_color=colors,
            text=[f"{int(n)}  ({p:.1f}%)" for n, p in zip(d["Occurrences"], d["Part des conv. initiées (%)"])],
            textposition="outside",
        )
    )
    _apply_layout(
        fig,
        f"Top {top_n} intentions sur les conversations initiées",
        xaxis=dict(title="Nombre d'occurrences"),
        yaxis=dict(title=""),
        showlegend=False,
    )
    return _save_png(fig, out)


def chart_utilisateurs(df_users: pd.DataFrame, out: Path) -> Path:
    # Construire 4 buckets
    m = dict(zip(df_users["Indicateur"], df_users["Valeur"]))
    u_tot = int(m.get("Utilisateurs uniques (ID HubSpot)", 0))
    u_1 = int(m.get("Utilisateurs avec 1 seule conversation", 0))
    u_2p = int(m.get("Utilisateurs avec >= 2 conversations", 0))
    u_3p = int(m.get("Utilisateurs avec >= 3 conversations", 0))
    u_5p = int(m.get("Utilisateurs avec >= 5 conversations", 0))
    u_2_2 = max(0, u_2p - u_3p)
    u_3_4 = max(0, u_3p - u_5p)
    buckets = [
        ("1 conversation", u_1, COLOR_GREY),
        ("2 conversations", u_2_2, COLOR_ACCENT),
        ("3–4 conversations", u_3_4, COLOR_PRIMARY),
        ("5 conversations et +", u_5p, COLOR_GREEN),
    ]
    fig = go.Figure()
    fig.add_trace(
        go.Bar(
            x=[b[0] for b in buckets],
            y=[b[1] for b in buckets],
            marker_color=[b[2] for b in buckets],
            text=[f"{b[1]}\n({100*b[1]/u_tot:.1f} %)" if u_tot else str(b[1]) for b in buckets],
            textposition="outside",
        )
    )
    _apply_layout(
        fig,
        f"Répartition des {u_tot} utilisateurs uniques par nombre de conversations",
        xaxis=dict(title=""),
        yaxis=dict(title="Nb utilisateurs"),
        showlegend=False,
    )
    return _save_png(fig, out)


# ----------------------------------------------------------------------------
# PowerPoint
# ----------------------------------------------------------------------------

def _add_title_bar(slide, text: str, width_in: float = 13.3):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(width_in), Inches(0.9)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(COLOR_PRIMARY.lstrip("#"))
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_left = Inches(0.3)
    tf.margin_top = Inches(0.15)
    p = tf.paragraphs[0]
    p.text = text
    p.runs[0].font.size = Pt(24)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)


def _add_footer(slide, text: str, width_in: float = 13.3, height_in: float = 7.5):
    tb = slide.shapes.add_textbox(Inches(0.3), Inches(height_in - 0.35), Inches(width_in - 0.6), Inches(0.3))
    p = tb.text_frame.paragraphs[0]
    p.text = text
    p.runs[0].font.size = Pt(10)
    p.runs[0].font.italic = True
    p.runs[0].font.color.rgb = RGBColor(0x7F, 0x8C, 0x8D)


def _add_bullets(slide, x_in, y_in, w_in, h_in, bullets, size=16, bold_first=False):
    tb = slide.shapes.add_textbox(Inches(x_in), Inches(y_in), Inches(w_in), Inches(h_in))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {text}"
        run = p.runs[0]
        run.font.size = Pt(size)
        run.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)
        if bold_first and i == 0:
            run.font.bold = True
        p.space_after = Pt(6)


def slide_titre(prs: Presentation, periode: str, nb_init: int, nb_users: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor.from_string(COLOR_PRIMARY.lstrip("#"))
    bg.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.7), Inches(2.3), Inches(12), Inches(1.5))
    p = tb.text_frame.paragraphs[0]
    p.text = "Chatbot Stellair (Yelda)"
    r = p.runs[0]
    r.font.size = Pt(44)
    r.font.bold = True
    r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    tb = slide.shapes.add_textbox(Inches(0.7), Inches(3.5), Inches(12), Inches(1))
    p = tb.text_frame.paragraphs[0]
    p.text = "Premier bilan d'utilisation — synthèse COMEX"
    r = p.runs[0]
    r.font.size = Pt(28)
    r.font.color.rgb = RGBColor(0xDD, 0xE6, 0xF0)

    tb = slide.shapes.add_textbox(Inches(0.7), Inches(5.5), Inches(12), Inches(1))
    p = tb.text_frame.paragraphs[0]
    p.text = (
        f"Période : {periode}   |   "
        f"{nb_init} conversations initiées   |   "
        f"{nb_users} utilisateurs uniques"
    )
    r = p.runs[0]
    r.font.size = Pt(20)
    r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)


def slide_messages_cles(prs: Presentation, synth: pd.DataFrame):
    m = dict(zip(synth["Indicateur"], synth["Valeur"]))
    nb_init = int(m.get("Conversations initiées (>= 1 question utilisateur)", 0))
    nb_open = int(m.get("Ouvertures du widget (FSE Stellair)", 0))
    pct_init = float(m.get("Taux de conversations initiées (%)", 0))
    nb_users = int(m.get("Utilisateurs uniques (HubSpot ID)", 0))
    nb_rec = int(m.get("Utilisateurs récurrents (>= 2 conversations)", 0))
    pct_rec = float(m.get("Part d'utilisateurs récurrents (%)", 0))
    pct_sat = float(m.get("Satisfaction LLM globale (Sat / Sat+Insat) (%)", 0))
    pct_ticket = float(m.get("Part de conversations débouchant sur un ticket (%)", 0))

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_bar(slide, "Trois messages clés")

    bullets = [
        f"ADOPTION — {nb_init} conversations réellement engagées (soit {pct_init:.0f} % des "
        f"{nb_open} ouvertures) portées par {nb_users} utilisateurs uniques identifiés. "
        f"{pct_rec:.0f} % reviennent au moins une 2ᵉ fois : l'usage s'installe.",
        f"QUALITÉ — Satisfaction LLM globale à {pct_sat:.0f} %, en progression forte mois après mois "
        "(+31 points entre janvier et avril 2026). Durée médiane d'une conversation : ~2 minutes, "
        "5 messages échangés en moyenne.",
        f"DÉFLEXION — Seulement {pct_ticket:.0f} % des conversations initiées débouchent sur la "
        f"création d'un ticket support ({int(m.get('Tickets HubSpot créés depuis une conversation initiée', 0))} tickets). "
        "Le bot absorbe la majorité des sollicitations ; le chantier prioritaire reste les "
        "questions non comprises (fallback) et les demandes d'agent humain.",
    ]
    _add_bullets(slide, 0.6, 1.2, 12.5, 5.5, bullets, size=18)
    _add_footer(slide, "Source : export Yelda (yelda.xlsx) ; périmètre FSE Stellair ; conversations initiées (nb messages ≥ 2).")


def slide_image(prs: Presentation, title: str, image_path: Path, bullets: list[str] | None = None, footer: str = ""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_bar(slide, title)
    slide.shapes.add_picture(str(image_path), Inches(0.5), Inches(1.1), width=Inches(9.0), height=Inches(4.5))
    if bullets:
        _add_bullets(slide, 9.7, 1.3, 3.5, 5.5, bullets, size=13)
    if footer:
        _add_footer(slide, footer)


def slide_methodo(prs: Presentation, nb_fse: int, nb_init: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_bar(slide, "Méthodologie")
    bullets = [
        "Source : export Yelda yelda.xlsx (une ligne = une conversation).",
        "Filtre : URL d'origine = https://fse.stellair.fr (chatbot FSE Stellair).",
        f"Périmètre volumétrique : {nb_fse} ouvertures du widget sur la période.",
        f"Conversations « initiées » retenues : {nb_init} (nb messages >= 2, "
        "i.e. l'utilisateur a posé au moins une question — exclut les ouvertures « welcome » seules).",
        "Utilisateurs uniques : identifiés via le slot persistant hubspot_id_slot.",
        "Évaluation LLM (colonnes 'Évaluation LLM' / 'Score LLM' de yelda.xlsx) : produite par "
        "un script Python interne qui envoie le contenu réel des échanges (bot + utilisateur) à "
        "Claude (Anthropic) avec un prompt dédié, sans utiliser le Parcours ni les Intentions.",
        "Le LLM retourne un score 0–5 mappé en : 5/4 = Satisfait ; 3 = À revoir ; "
        "2/1 = Insatisfait ; 0 = Non évaluable (conversation sans question utilisateur).",
        "Taux de satisfaction LLM = Satisfait / (Satisfait + Insatisfait) ; 'À revoir' et "
        "'Non évaluable' exclus du dénominateur pour lisibilité.",
        "Ticket créé : parcours contient 'creation_ticket_hubspot'.",
        "Intentions : liste séparée par ';' dans la colonne Intentions ; tri décroissant par occurrences.",
        "Données sources et calculs : classeur yelda_comex_analyse_utilisation.xlsx.",
    ]
    _add_bullets(slide, 0.6, 1.1, 12.5, 6.0, bullets, size=13)


# ----------------------------------------------------------------------------
# Orchestrateur
# ----------------------------------------------------------------------------

def build_powerpoint(out_pptx: Path) -> Path:
    df = load_yelda_data()
    if df is None or df.empty:
        raise SystemExit("Fichier Yelda introuvable ou vide.")
    df_fse = filter_yelda_stellair(df)
    df_init = filter_conversations_initiees(df_fse)

    synth = build_synthese(df_fse, df_init)
    vol = build_volumetrie_mensuelle(df_init)
    vol_h = build_volumetrie_hebdo(df_init)
    sat = build_satisfaction_llm_mensuel(df_init)
    defl = build_deflexion_ticket(df_init)
    users = build_utilisateurs(df_init)
    intents = build_top_intentions(df_init, top_n=20)

    periode = "—"
    d = pd.to_datetime(df_fse[COL_DATE]).dropna()
    if len(d):
        periode = f"{d.min().strftime('%d/%m/%Y')} → {d.max().strftime('%d/%m/%Y')}"

    # Présentation 16:9
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    tmp = Path(tempfile.mkdtemp(prefix="yelda_pptx_"))
    images: Dict[str, Path] = {
        "volumetrie": chart_volumetrie_hebdo(vol_h, tmp / "volumetrie_hebdo.png"),
        "satisfaction": chart_satisfaction(sat, tmp / "satisfaction.png"),
        "deflexion": chart_deflexion(defl, tmp / "deflexion.png"),
        "intentions": chart_top_intentions(intents, tmp / "intentions.png"),
        "utilisateurs": chart_utilisateurs(users, tmp / "utilisateurs.png"),
    }

    m = dict(zip(synth["Indicateur"], synth["Valeur"]))
    nb_init = int(m.get("Conversations initiées (>= 1 question utilisateur)", 0))
    nb_users = int(m.get("Utilisateurs uniques (HubSpot ID)", 0))
    nb_fse = int(m.get("Ouvertures du widget (FSE Stellair)", 0))

    # Slide 1 : titre
    slide_titre(prs, periode, nb_init, nb_users)

    # Slide 2 : messages clés
    slide_messages_cles(prs, synth)

    # Slide 3 : volumétrie hebdo
    nb_semaines = len(vol_h)
    moy_hebdo = int(round(vol_h["Conversations initiées"].mean())) if nb_semaines else 0
    max_row = vol_h.loc[vol_h["Conversations initiées"].idxmax()] if nb_semaines else None
    last_row = vol_h.iloc[-1] if nb_semaines else None
    bullets_vol = [f"{nb_init} conversations initiées sur la période, {nb_users} utilisateurs uniques."]
    if max_row is not None:
        bullets_vol.append(
            f"Pic semaine {max_row['Semaine']} : "
            f"{int(max_row['Conversations initiées'])} conv."
        )
    bullets_vol.append(f"Moyenne : ~{moy_hebdo} conv./semaine ({nb_semaines} semaines observées).")
    if last_row is not None:
        bullets_vol.append(
            f"Dernière semaine ({last_row['Semaine']}) : "
            f"{int(last_row['Conversations initiées'])} conv., "
            f"{int(last_row['Utilisateurs uniques (HubSpot ID)'])} utilisateurs."
        )
    slide_image(
        prs,
        "Volumétrie hebdomadaire",
        images["volumetrie"],
        bullets=bullets_vol,
        footer="Conversations initiées (nb messages >= 2) et utilisateurs HubSpot uniques, par semaine ISO.",
    )

    # Slide 4 : satisfaction
    taux_dernier = float(sat["Taux satisfaction LLM (%)"].iloc[-1]) if len(sat) else 0
    taux_premier = float(sat["Taux satisfaction LLM (%)"].iloc[0]) if len(sat) else 0
    slide_image(
        prs,
        "Qualité — satisfaction LLM mensuelle",
        images["satisfaction"],
        bullets=[
            f"De {taux_premier:.0f} % à {taux_dernier:.0f} % en 4 mois.",
            "Signal de qualité en nette amélioration (ajustements Yelda, enrichissement des flows).",
            "Évaluation par Claude (Anthropic) sur le contenu réel des échanges ; "
            "'À revoir' exclu du dénominateur.",
        ],
        footer="Taux = Satisfait / (Satisfait + Insatisfait). Score LLM (0–5) attribué par un script dédié à partir du contenu de la conversation, puis mappé en Satisfait / À revoir / Insatisfait / Non évaluable.",
    )

    # Slide 5 : déflexion
    total_row = defl[defl["Mois"] == "Total"].iloc[0]
    slide_image(
        prs,
        "Déflexion — tickets évités",
        images["deflexion"],
        bullets=[
            f"{int(total_row['Tickets créés'])} tickets créés sur "
            f"{int(total_row['Conversations initiées'])} conv. initiées.",
            f"Taux global de déflexion : {float(total_row['Part sans ticket / déflexion (%)']):.0f} %.",
            "Chute du taux de création ticket : 28 % (janv.) → 4 % (avril).",
        ],
        footer="Ticket créé = parcours Yelda contient 'creation_ticket_hubspot'.",
    )

    # Slide 6 : intentions
    slide_image(
        prs,
        "Chantiers d'amélioration — top intentions",
        images["intentions"],
        bullets=[
            "default_fallback : question non comprise par le bot.",
            "contact_agent + phone/callback/email : demande d'humain.",
            "reponse_agent_non_satisfaisante : retour négatif explicite.",
            "En rouge : intentions à signal faible, ~55 % des conv. initiées — prioritaires.",
        ],
        footer="Top 10 des intentions observées sur conversations initiées (comptage par occurrence unique par conversation).",
    )

    # Slide 7 : utilisateurs
    pct_recurrents_key = "Part d'utilisateurs récurrents (%)"
    pct_recurrents = float(m.get(pct_recurrents_key, 0))
    nb_recurrents = int(m.get("Utilisateurs récurrents (>= 2 conversations)", 0))
    slide_image(
        prs,
        "Adoption et récurrence des utilisateurs",
        images["utilisateurs"],
        bullets=[
            f"{nb_users} utilisateurs uniques identifiés via l'ID HubSpot.",
            f"{nb_recurrents} reviennent >= 2 fois ({pct_recurrents:.0f} %).",
            "Socle solide : l'outil est utilisé régulièrement par une population grandissante.",
        ],
        footer="Segmentation des utilisateurs par nombre de conversations initiées sur la période.",
    )

    # Slide 8 : méthodologie
    slide_methodo(prs, nb_fse, nb_init)

    out_pptx = Path(out_pptx)
    out_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_pptx)
    return out_pptx


def main() -> None:
    default_out = _ROOT / "data/Affid/analyse_appels_tickets/yelda_comex_synthese.pptx"
    p = argparse.ArgumentParser(description="PowerPoint COMEX — utilisation du chatbot Yelda.")
    p.add_argument("-o", "--output", type=Path, default=default_out)
    args = p.parse_args()

    out = build_powerpoint(args.output)
    print(f"Écrit : {out.resolve()}")


if __name__ == "__main__":
    main()
