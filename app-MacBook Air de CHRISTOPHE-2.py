import streamlit as st
from utils.authentification import authenticate_user
from data_processing.aircall_processing import process_aircall_data, def_df_support, agents_all, line_tous, load_aircall_data
#from data_processing.aircall_processing import get_df_support
from data_processing.hubspot_processing import process_hubspot_data, load_hubspot_data
from data_processing.kpi_generation import generate_kpis, filtrer_par_periode, calculate_ticket_response_time, graph_activite, evo_appels_ticket
from utils.streamlit_helpers import load_data
from utils.powerpoint_helpers import create_powerpoint, create_powerpoint_agents, create_powerpoint_stellair_minimal, create_powerpoint_agents_report
import config
from config import CREDENTIALS, AIRCALL_DATA_PATH_V1, AIRCALL_DATA_PATH_V2, HUBSPOT_TICKET_DATA_PATH, EVALUATION_DATA_PATH
import streamlit_authenticator as stauth
from pptx import Presentation
from pptx.util import Inches, Pt
import io
import tempfile
import pandas as pd

import streamlit as st
#from selenium import webdriver
#from selenium.webdriver.chrome.service import Service
#from webdriver_manager.chrome import ChromeDriverManager


st.set_page_config(
    page_title=":bar_chart: Dashboard support",
    layout="wide",
    initial_sidebar_state="collapsed",
)

# Authentification de l'utilisateur
authenticator = stauth.Authenticate(
    config.CREDENTIALS,
    'dashboard_support',
    config.COOKIE_KEY,  # Assurez-vous que cette clé est unique et sécurisée
    cookie_expiry_days=2
)

# Utilisation du nouvel argument 'fields'
name, authentification_status, username = authenticator.login('main', fields={'Form name': 'custom_form_name'})

if authentification_status == False:
    st.error('Nom d\'utilisateur/mot de passe incorrect')
elif authentification_status == None:
    st.warning('Veuillez entrer votre nom d\'utilisateur et mot de passe')
elif authentification_status:
    st.sidebar.title(f"Bienvenue, {name}")
    authenticator.logout('Logout', 'sidebar')

    # Bouton dans la sidebar pour recharger les données Aircall et HubSpot
    if st.sidebar.button("Recharger les données"):
        # Vider tous les caches
        cache_keys_to_clear = [
            'df_aircall', 'df_tickets', 'df_evaluation',
            'df_support_processed', 'df_tickets_processed', 'kpis_all',
            'df_aircall_processed'
        ]
        
        # Ajouter les caches des dataframes support
        for option in ["support_suresnes", "support_armatis", "support_stellair", "support_affid", "xmed"]:
            cache_keys_to_clear.append(f"df_support_{option}")
        
        for key in cache_keys_to_clear:
            if key in st.session_state:
                del st.session_state[key]
        
        # Recharger les données
        st.session_state['df_aircall'] = load_aircall_data(AIRCALL_DATA_PATH_V1, AIRCALL_DATA_PATH_V2, force_reload=True)
        st.session_state['df_tickets'] = load_hubspot_data(HUBSPOT_TICKET_DATA_PATH)
        st.session_state['df_evaluation'] = pd.read_excel(EVALUATION_DATA_PATH)
        st.success("Données Aircall, HubSpot et Evaluation rechargées !")

    # Chargement initial si pas déjà en mémoire
    if 'df_aircall' not in st.session_state:
        st.session_state['df_aircall'] = load_aircall_data(AIRCALL_DATA_PATH_V1, AIRCALL_DATA_PATH_V2, force_reload=False)
    if 'df_tickets' not in st.session_state:
        st.session_state['df_tickets'] = load_hubspot_data(HUBSPOT_TICKET_DATA_PATH)
    if 'df_evaluation' not in st.session_state:
        st.session_state['df_evaluation'] = pd.read_excel(EVALUATION_DATA_PATH)

    df_aircall = st.session_state['df_aircall']
    df_tickets = st.session_state['df_tickets']
    df_evaluation = st.session_state['df_evaluation']

    # Optimisation : traitement des données une seule fois et mise en cache
    if 'df_support_processed' not in st.session_state:
        st.session_state['df_support_processed'] = def_df_support(
            process_aircall_data(df_aircall), 
            process_aircall_data(df_aircall), 
            line_tous, 
            agents_all
        )
    if 'df_tickets_processed' not in st.session_state:
        st.session_state['df_tickets_processed'] = process_hubspot_data(df_tickets)

    df_support = st.session_state['df_support_processed']
    df_tickets_processed = st.session_state['df_tickets_processed']

    # Génération des KPIs (mise en cache si possible)
    if 'kpis_all' not in st.session_state:
        st.session_state['kpis_all'] = generate_kpis(df_support, df_tickets_processed, 'agents_all')
    
    kpis = st.session_state['kpis_all']

    PAGES = {
        "Support": "support",
        "Agents": "agents",
        "Tickets": "tickets"
    }


    selection_page = st.sidebar.selectbox("Choix de la page", list(PAGES.keys()), key="unique_page_selection")

    if selection_page == "Support":
        from data_processing.aircall_processing import process_aircall_data, def_df_support, agents_all, line_tous, agents_support, line_support, line_armatis, agents_armatis
        from data_processing.kpi_generation import (
            filtrer_par_periode, generate_kpis, convert_to_sixtieth, 
            graph_activite, graph_taux_jour, graph_taux_heure, 
            graph_activite_xmed, calculate_ticket_response_time
        )


        def support():
            st.write('DEBUG: PAGE SUPPORT')

            dataframe_option = st.sidebar.selectbox("Choisir le dataframe", ["support_suresnes", "support_armatis", "support_stellair", "support_affid", "xmed"], key="unique_dataframe_selection")

            # Optimisation : traitement des données une seule fois et mise en cache
            if 'df_aircall_processed' not in st.session_state:
                st.session_state['df_aircall_processed'] = process_aircall_data(df_aircall)

            df_aircall_processed = st.session_state['df_aircall_processed']

            # Cache pour les différents dataframes
            cache_key = f"df_support_{dataframe_option}"
            if cache_key not in st.session_state:
                if dataframe_option == "support_stellair":
                    df_stellair = def_df_support(df_aircall_processed, df_aircall_processed, line_tous, agents_all)
                    df_stellair = df_stellair[(df_stellair['line'] == 'armatistechnique') | (df_stellair['IVR Branch'] == 'Stellair')]
                    st.session_state[cache_key] = df_stellair
                elif dataframe_option == "support_affid":
                    df_affid = def_df_support(df_aircall_processed, df_aircall_processed, line_tous, agents_support)
                    df_affid = df_affid[df_affid['IVR Branch'] == 'Affid']
                    st.session_state[cache_key] = df_affid
                elif dataframe_option == "xmed":
                    df_xmed = def_df_support(df_aircall_processed, df_aircall_processed, line_tous, agents_all)
                    df_xmed = df_xmed[(df_xmed['line'] == 'xmed')]
                    st.session_state[cache_key] = df_xmed
                else:
                    # Pour les autres options, utiliser la configuration existante
                    dataframe_config = {
                        "support_suresnes": {
                            "df": def_df_support(df_aircall_processed, df_aircall_processed, line_support, agents_support),
                            "agents": agents_support
                        },
                        "support_armatis": {
                            "df": def_df_support(df_aircall_processed, df_aircall_processed, line_armatis, agents_armatis),
                            "agents": agents_armatis
                        }
                    }
                    st.session_state[cache_key] = dataframe_config[dataframe_option]["df"]

            # Configuration des agents selon l'option
            agents_config = {
                "support_suresnes": agents_support,
                "support_armatis": agents_armatis,
                "support_stellair": agents_all,
                "support_affid": agents_support,
                "xmed": agents_support
            }

            df_support = st.session_state[cache_key]
            agents = agents_config[dataframe_option]

            periode = st.selectbox("Sélectionnez une période :", 
                    ["1 an", "6 derniers mois", "3 derniers mois", "Dernier mois"],
                    index=0)
            
            df_support = filtrer_par_periode(df_support, periode)
            kpis = generate_kpis(filtrer_par_periode(df_support, periode), filtrer_par_periode(df_tickets_processed, periode), agents)

            # Utiliser graph_activite_xmed pour la page xmed, sinon graph_activite normal
            if dataframe_option == "xmed":
                col1, col2, col3 = st.columns(3)
                col1.metric("Taux de service en %", kpis['Taux_de_service'])
                col2.metric("Appels entrant / Jour", kpis['Entrant'])
                col3.metric("Numéros uniques / Jour", kpis['Numero_unique'])
                st.plotly_chart(graph_activite_xmed(df_support), use_container_width=True)
            elif dataframe_option == "support_stellair":
                # KPI spécifiques pour support_stellair
                col1, col2, col3 = st.columns(3)
                col1.metric("Taux de service en %", kpis['Taux_de_service'])
                col2.metric("Appels entrant / Jour", kpis['Entrant'])
                col3.metric("Numéros uniques / Jour", kpis['Numero_unique'])

                col4, col5 = st.columns(2)
                col4.metric(
                    "Entrants vs Tickets (%)",
                    f"{round(kpis['activite_appels_pourcentage'] * 100, 2)}% / {round(kpis['activite_tickets_pourcentage'] * 100, 2)}%"
                )
                # Calcul du temps de réponse aux tickets
                moyenne_temps_reponse, graph_temps_reponse, df_temps_tickets = calculate_ticket_response_time(
                    filtrer_par_periode(df_tickets_processed, periode), 
                    agents_all
                )

                col5.metric("Temps de réponse moyen aux tickets (h:min)", f"{int(moyenne_temps_reponse)}:{int((moyenne_temps_reponse % 1) * 60):02d}")

                # Graphique principal
                st.plotly_chart(graph_activite(df_support), use_container_width=True)
                # Graphique des temps de réponse aux tickets
                st.plotly_chart(graph_temps_reponse, use_container_width=True)
                # Affichage des graphiques de taux
                col_graph1, col_graph2 = st.columns(2)
                col_graph1.plotly_chart(graph_taux_jour(df_support), use_container_width=True)
                col_graph2.plotly_chart(graph_taux_heure(df_support), use_container_width=True)
                st.plotly_chart(kpis['evo_appels_tickets'])
            elif dataframe_option == "support_armatis" or dataframe_option == "support_affid":
                col1, col2, col3 = st.columns(3)
                col1.metric("Taux de service en %", kpis['Taux_de_service'])
                col2.metric("Appels entrant / Jour", kpis['Entrant'])
                col3.metric("Numéros uniques / Jour", kpis['Numero_unique'])

                col4, col5 = st.columns(2)
                col4.metric(
                    "Entrants vs Tickets (%)",
                    f"{round(kpis['activite_appels_pourcentage'] * 100, 2)}% / {round(kpis['activite_tickets_pourcentage'] * 100, 2)}%"
                )
                col5.empty() # Pas de KPI supplémentaire ici
                
                # Condition spécifique pour Affid
                if dataframe_option == "support_affid":
                    st.info("Les graphiques ci-dessous montrent la répartition des appels entre les branches IVR 'Affid' et 'Stellair'.")
                    
                    fig_charge_pourcentage, fig_charge_volume, _ = graph_charge_affid_stellair(
                        filtrer_par_periode(
                            st.session_state['df_support_processed'],
                            periode
                        )
                    )
                    st.plotly_chart(fig_charge_pourcentage, use_container_width=True)
                    st.plotly_chart(fig_charge_volume, use_container_width=True)
                else:
                    # Graphiques pour Armatis
                    st.plotly_chart(graph_activite(df_support), use_container_width=True)
                    col_graph1, col_graph2 = st.columns(2)
                    col_graph1.plotly_chart(graph_taux_jour(df_support), use_container_width=True)
                    col_graph2.plotly_chart(graph_taux_heure(df_support), use_container_width=True)
                    st.plotly_chart(kpis['evo_appels_tickets'])
            else:
                col1, col2, col3 = st.columns(3)
                col1.metric("Taux de service en %", kpis['Taux_de_service'])
                col2.metric("Appels entrant / Jour", kpis['Entrant'])
                col3.metric("Numéros uniques / Jour", kpis['Numero_unique'])

                col4, col5 = st.columns(2)
                col4.metric(
                    "Entrants vs Tickets (%)",
                    f"{round(kpis['activite_appels_pourcentage'] * 100, 2)}% / {round(kpis['activite_tickets_pourcentage'] * 100, 2)}%"
                )
                col5.empty()

                st.plotly_chart(graph_activite(df_support), use_container_width=True)
                # Affichage des graphiques de taux
                col_graph1, col_graph2 = st.columns(2)
                col_graph1.plotly_chart(graph_taux_jour(df_support), use_container_width=True)
                col_graph2.plotly_chart(graph_taux_heure(df_support), use_container_width=True)
                st.plotly_chart(kpis['evo_appels_tickets'])

            import os

            if st.button("Exporter les KPI en HTML et PDF"):
                try:
                    html_file = generate_html_report(kpis, df_support, periode)
                    st.success(f"Les KPI ont été exportés en HTML sous '{html_file}'.")

                    # Affiche bouton de téléchargement HTML
                    with open(html_file, "r", encoding="utf-8") as file:
                        html_content = file.read()

                    st.download_button(
                        label="📥 Télécharger le HTML",
                        data=html_content,
                        file_name=os.path.basename(html_file),
                        mime="text/html"
                    )

                    # Bouton de téléchargement PDF
                    pdf_file = html_file.replace(".html", ".pdf")
                    if os.path.exists(pdf_file):
                        with open(pdf_file, "rb") as f:
                            pdf_bytes = f.read()
                        st.download_button(
                            label="📥 Télécharger le PDF",
                            data=pdf_bytes,
                            file_name=os.path.basename(pdf_file),
                            mime="application/pdf"
                        )
                    else:
                        st.warning("Le fichier PDF n'a pas été généré.")
                except Exception as e:
                    st.error(f"Erreur lors de la génération du rapport : {e}")

            if st.button("Exporter en PowerPoint"):
                try:
                    st.info("Début génération PowerPoint (debug)")
                    st.write("KPIs:", kpis)
                    st.write("df_support shape:", df_support.shape)
                    st.write("Période:", periode)
                    pptx_io = create_powerpoint(kpis, df_support, periode)
                    st.info("PowerPoint généré (debug)")
                    st.download_button(
                        label="📥 Télécharger la présentation PowerPoint",
                        data=pptx_io,
                        file_name="dashboard_support.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                    )
                    st.success("La présentation PowerPoint a été générée avec succès!")
                except Exception as e:
                    st.error(f"Erreur lors de la génération de la présentation : {e}")

            # --- RAPPORT POWERPOINT STELLAIR ---
            if st.sidebar.button("Générer rapport PowerPoint Stellair"):
                # Préparation des données 6 mois
                df_support_stellair_6m = def_df_support(process_aircall_data(df_aircall), process_aircall_data(df_aircall), line_tous, agents_all)
                df_support_stellair_6m = df_support_stellair_6m[(df_support_stellair_6m['line'] == 'armatistechnique') | (df_support_stellair_6m['IVR Branch'] == 'Stellair')]
                df_support_stellair_6m = filtrer_par_periode(df_support_stellair_6m, "6 derniers mois")
                df_tickets_6m = filtrer_par_periode(df_tickets_processed, "6 derniers mois")
                kpis_6m = generate_kpis(df_support_stellair_6m, df_tickets_6m, agents_all)
                moyenne_temps_reponse_6m, graph_temps_reponse_6m = calculate_ticket_response_time(df_tickets_6m, agents_all)
                kpis_6m['moyenne_temps_reponse'] = moyenne_temps_reponse_6m
                graph_activite_6m = graph_activite(df_support_stellair_6m)
                evo_appels_tickets_6m, _, _, _ = evo_appels_ticket(df_tickets_6m, df_support_stellair_6m)

                # Préparation des données 3 mois
                df_support_stellair_3m = def_df_support(process_aircall_data(df_aircall), process_aircall_data(df_aircall), line_tous, agents_all)
                df_support_stellair_3m = df_support_stellair_3m[(df_support_stellair_3m['line'] == 'armatistechnique') | (df_support_stellair_3m['IVR Branch'] == 'Stellair')]
                df_support_stellair_3m = filtrer_par_periode(df_support_stellair_3m, "3 derniers mois")
                df_tickets_3m = filtrer_par_periode(df_tickets_processed, "3 derniers mois")
                kpis_3m = generate_kpis(df_support_stellair_3m, df_tickets_3m, agents_all)
                moyenne_temps_reponse_3m, _ = calculate_ticket_response_time(df_tickets_3m, agents_all)

                # Génération du rapport
                st.write('DEBUG: bouton rapport Stellair cliqué')
                pptx_io = create_powerpoint_stellair_minimal(
                    kpis_6m, "6 derniers mois", graph_activite_6m
                )
                st.sidebar.success("Rapport PowerPoint généré !")
                st.sidebar.download_button(
                    label="📥 Télécharger le rapport Stellair",
                    data=pptx_io,
                    file_name="rapport_stellair.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )
            # --- FIN RAPPORT POWERPOINT STELLAIR ---

        support()



    elif selection_page == "Agents":
        import pandas as pd
        from data_processing.kpi_generation import (
            filtrer_par_periode, generate_kpis, convert_to_sixtieth, 
            filtrer_par_agent, charge_entrant_sortant, df_compute_ticket_appels_metrics, 
            historique_scores_total, charge_ticket, 
            filter_evaluation, calculate_performance_score,
            graph_tickets_n2_par_semaine, graph_tickets_n2_resolus_par_agent, graph_tickets_pierre_goupillon,
            graph_tickets_n1_par_semaine, get_n1_agents_list
        )
        from data_processing.hubspot_processing import process_hubspot_data
        from data_processing.aircall_processing import process_aircall_data, def_df_support, agents_all, line_tous


        

        def agents():
            # Chargement des données
            periode_selectbox = st.selectbox("Période KPI automatique :", 
                ["1 an", "6 derniers mois", "3 derniers mois", "Dernier mois"],
                index=0)
            
            # Paramètres de scoring
            st.sidebar.markdown("## ⚙️ Paramètres de scoring")
            with st.sidebar.expander("🎯 Objectifs", expanded=True):
                objectif_total = st.number_input("Objectif total de demandes par jour", min_value=1, max_value=50, value=25)
                ratio_appels = st.slider("Ratio d'appels (%)", min_value=0, max_value=100, value=70) / 100
                ratio_tickets = 1 - ratio_appels  # Calculé automatiquement
                st.write(f"Ratio de tickets : {ratio_tickets:.0%}")
                objectif_taux_service = st.slider("Objectif taux de service (%)", min_value=0, max_value=100, value=70) / 100
            
            df_support = def_df_support(process_aircall_data(df_aircall), process_aircall_data(df_aircall), line_tous, agents_all)
            
            agents_n1 = ['Olivier Sainte-Rose','Mourad HUMBLOT', 'Archimede KESSI', 'Celine Crendal']
            agents_n1_tickets = agents_n1 + ['Frederic SAUVAN']

            # Récupérer le DataFrame et les métriques
            df_conforme = df_compute_ticket_appels_metrics(agents_n1, filtrer_par_periode(df_tickets_processed, periode_selectbox), filtrer_par_periode(df_support, periode_selectbox))
            
            # Vérifier que df_conforme est un DataFrame avant d'appliquer les calculs
            if isinstance(df_conforme, pd.DataFrame):
                # Appliquer le score de performance avec les paramètres personnalisés
                df_conforme['score_performance'] = df_conforme.apply(
                    lambda row: calculate_performance_score(
                        row,
                        objectif_total=objectif_total,
                        ratio_appels=ratio_appels,
                        ratio_tickets=ratio_tickets,
                        objectif_taux_service=objectif_taux_service
                    ),
                    axis=1
                )

                df_conforme = df_conforme[['Agent', 'score_performance', "Nombre d'appel traité", 
                                        'Nombre de ticket traité', '% appel entrant agent', 
                                        '% tickets', '% appels']]
            else:
                st.error("Erreur: Les données ne sont pas dans le format attendu.")
                return

            # Fonction de style
            def style_scores(df):
                def apply_style(row):
                    styles = {}
                    score = row['score_performance']
                    taux_service = row['% appel entrant agent'] / 100

                    if score >= 70:
                        styles['score_performance'] = 'background-color: #a7dba7'  # vert
                    elif score >= 60:
                        styles['score_performance'] = 'background-color: #f7c97f'  # orange
                    else:
                        styles['score_performance'] = 'background-color: #f28e8e'  # rouge

                    return pd.Series(styles)

                # Formater toutes les colonnes numériques avec 2 décimales
                return df.style.apply(apply_style, axis=1).format({
                    'score_performance': '{:.2f}',
                    "Nombre d'appel traité": '{:.2f}',
                    'Nombre de ticket traité': '{:.2f}',
                    '% appel entrant agent': '{:.2f}',
                    '% tickets': '{:.2f}',
                    '% appels': '{:.2f}'
                })

            # ----------- SIDEBAR - FILTRES -------------------
            with st.sidebar:
                st.header("🔍 Filtres évaluations")
                with st.form("filter_form"):
                    agents = st.multiselect("Agent(s)", options=df_evaluation["agent"].unique(), default=df_evaluation["agent"].unique())
                    periodes_eval = st.multiselect("Période", options=df_evaluation["quarter"].unique(), default=df_evaluation["quarter"].unique())
                    submit = st.form_submit_button(label="🎯 Appliquer les filtres")

            # ----------- TITRE -------------------
            st.title("📊 Tableau de bord des agents N1")

            st.markdown(f"**Période KPI automatique sélectionnée** : `{periode_selectbox}`")

            # ----------- SCORING AUTOMATIQUE -------------------
            st.markdown("## ⚖️ Scores automatiques")

            st.markdown("""
            ### 🧠 Méthodologie de scoring

            Le score de performance (0 à 100) est calculé selon quatre critères :
            - Volume total (45%) : Nombre total de demandes traitées par rapport à l'objectif
            - Répartition (25%) : Équilibre entre appels et tickets selon le ratio défini
            - Comparaison à la moyenne (15%) : Performance par rapport à la moyenne du service
            - Taux d'appels entrants (15%) : Pourcentage d'appels entrants par rapport au total

            Codes couleur :
            - 🟢 **Vert** : score ≥ 70%
            - 🟠 **Orange** : 60% ≤ score < 70%
            - 🔴 **Rouge** : score < 60%
            """)

            # Calcul du nombre total de tickets traités par Céline Crendal (toutes orthographes)


            styled_df = style_scores(df_conforme)
            st.dataframe(styled_df, use_container_width=True)

            # ----------- GRAPHIQUES KPI -------------------
            st.markdown("## 📈 Graphiques d'activité")

            fig_line_entrant, fig_pie_entrant = charge_entrant_sortant(filtrer_par_periode(df_support, periode_selectbox), agents_n1)
            fig_line_ticket, fig_pie_ticket = charge_ticket(filtrer_par_periode(df_tickets_processed, periode_selectbox), agents_n1_tickets)

            st.markdown("### 📞 Activité téléphonique")
            col1, col2 = st.columns(2)
            col1.plotly_chart(fig_line_entrant, use_container_width=True)
            col2.plotly_chart(fig_pie_entrant, use_container_width=True)

            st.markdown("### 🧾 Activité tickets")
            col1, col2 = st.columns(2)
            col1.plotly_chart(fig_line_ticket, use_container_width=True)
            col2.plotly_chart(fig_pie_ticket, use_container_width=True)

            # ----------- GRAPHIQUE N1 PAR SEMAINE -------------------
            st.markdown("## 📊 Graphique N1 par semaine")
            
            # Obtenir la liste des agents N1 disponibles
            agents_n1_disponibles = get_n1_agents_list(filtrer_par_periode(df_tickets_processed, periode_selectbox))
            
            # Sélecteur d'agent pour le graphique N1
            if agents_n1_disponibles:
                agent_n1_selection = st.selectbox(
                    "Sélectionner un agent N1 pour afficher ses tickets ouverts :",
                    ["Aucun"] + agents_n1_disponibles,
                    key="agent_n1_selection"
                )
                
                # Créer le graphique N1
                if agent_n1_selection != "Aucun":
                    fig_n1 = graph_tickets_n1_par_semaine(
                        filtrer_par_periode(df_tickets_processed, periode_selectbox), 
                        selected_agent=agent_n1_selection
                    )
                else:
                    fig_n1 = graph_tickets_n1_par_semaine(
                        filtrer_par_periode(df_tickets_processed, periode_selectbox)
                    )
                
                st.plotly_chart(fig_n1, use_container_width=True)
            else:
                st.info("Aucun agent N1 trouvé pour la période sélectionnée.")

            # ----------- HISTORIQUE SCORES -------------------
            st.markdown("## 📊 Historique des scores")
            st.plotly_chart(historique_scores_total(agents_n1, filtrer_par_periode(df_tickets_processed, periode_selectbox), filtrer_par_periode(df_support, periode_selectbox)), use_container_width=True)

            # ----------- TABLEAU DE BORD N1 -------------------
            st.markdown("## Tableau de bord des agents N2")

            # Graphique 1
            fig1 = graph_tickets_n2_par_semaine(filtrer_par_periode(df_tickets_processed, periode_selectbox))
            st.plotly_chart(fig1, use_container_width=True)

            # Graphique 2
            #fig2 = graph_tickets_n2_resolus_par_agent(filtrer_par_periode(df_tickets_processed, periode_selectbox))
            #st.plotly_chart(fig2, use_container_width=True)

            # Graphique 3
            #fig3 = graph_tickets_pierre_goupillon(filtrer_par_periode(df_tickets_processed, periode_selectbox))
            #st.plotly_chart(fig3, use_container_width=True)

            # Graphique des tickets ouverts de Pierre Goupillon (comparatif)
            from data_processing.kpi_generation import graph_tickets_ouverts_pierre_goupillon
            #fig3_ouverts = graph_tickets_ouverts_pierre_goupillon(filtrer_par_periode(df_tickets_processed, periode_selectbox))
           # st.plotly_chart(fig3_ouverts, use_container_width=True)

            # ----------- GRAPHIQUE CUMULATIF N1 -------------------
            st.markdown("## 📈 Cumulatif tickets N1 (tous les tickets SSI/SSIA/SPSA)")
            from data_processing.kpi_generation import graph_tickets_n1_cumulatif
            
            # Obtenir le graphique et les agents disponibles
            fig_cumul_n1, agents_disponibles, tickets_n1_en_cours = graph_tickets_n1_cumulatif(filtrer_par_periode(df_tickets_processed, periode_selectbox))
            
            # Filtre par agent
            if agents_disponibles:
                agent_selection = st.selectbox(
                    "Filtrer par agent :",
                    ["Tous les agents"] + agents_disponibles,
                    key="agent_filter_cumul_n1"
                )
                
                # Appliquer le filtre si un agent spécifique est sélectionné
                if agent_selection != "Tous les agents":
                    fig_cumul_n1, _, tickets_n1_en_cours = graph_tickets_n1_cumulatif(
                        filtrer_par_periode(df_tickets_processed, periode_selectbox),
                        agent_selection
                    )
            
            st.plotly_chart(fig_cumul_n1, use_container_width=True)
            
            # Bouton pour afficher les tickets N1 en cours
            if not tickets_n1_en_cours.empty:
                nb_tickets_n1_en_cours = len(tickets_n1_en_cours)
                st.markdown(f"**🔴 {nb_tickets_n1_en_cours} ticket(s) N1 en cours**")
                
                # Bouton pour télécharger le DataFrame des tickets N1 en cours
                if st.button(f"📥 Télécharger la liste des {nb_tickets_n1_en_cours} tickets N1 en cours", key="btn_download_tickets_n1_en_cours"):
                    # Préparer les données pour le téléchargement
                    tickets_n1_download = tickets_n1_en_cours[['Ticket ID', 'Date', 'Statut du ticket', 'Propriétaire du ticket']].copy()
                    tickets_n1_download = tickets_n1_download.rename(columns={
                        'Ticket ID': 'ID_Ticket',
                        'Date': 'Date_Derniere_Mise_A_Jour',
                        'Statut du ticket': 'Statut',
                        'Propriétaire du ticket': 'Agent'
                    })
                    tickets_n1_download = tickets_n1_download.sort_values('Date_Derniere_Mise_A_Jour', ascending=False)
                    
                    # Convertir en CSV
                    csv_data_n1 = tickets_n1_download.to_csv(index=False, encoding='utf-8-sig')
                    
                    # Bouton de téléchargement
                    st.download_button(
                        label="📄 Télécharger CSV des tickets N1 en cours",
                        data=csv_data_n1,
                        file_name=f"tickets_n1_en_cours_{periode_selectbox.replace(' ', '_')}.csv",
                        mime="text/csv"
                    )
                    
                    # Afficher un aperçu du DataFrame
                    st.markdown("### 📋 Aperçu des tickets N1 en cours")
                    st.dataframe(
                        tickets_n1_download,
                        use_container_width=True,
                        column_config={
                            "ID_Ticket": st.column_config.TextColumn("ID Ticket", width="medium"),
                            "Date_Derniere_Mise_A_Jour": st.column_config.DateColumn("Date dernière mise à jour", width="medium"),
                            "Statut": st.column_config.TextColumn("Statut", width="small"),
                            "Agent": st.column_config.TextColumn("Agent", width="medium")
                        }
                    )
            else:
                st.success("✅ Aucun ticket N1 en cours !")

            # ----------- GRAPHIQUE CUMULATIF N2 -------------------
            st.markdown("## 📈 Cumulatif tickets N2 (passés vs résolus)")
            from data_processing.kpi_generation import graph_tickets_n2_cumulatif
            
            # Obtenir le graphique et les pipelines disponibles
            fig_cumul_n2, pipelines_disponibles, tickets_en_cours = graph_tickets_n2_cumulatif(filtrer_par_periode(df_tickets_processed, periode_selectbox))
            
            # Filtre par pipeline
            if pipelines_disponibles:
                pipeline_selection = st.selectbox(
                    "Filtrer par pipeline :",
                    ["Tous les pipelines"] + pipelines_disponibles,
                    key="pipeline_filter_cumul"
                )
                
                # Appliquer le filtre si un pipeline spécifique est sélectionné
                if pipeline_selection != "Tous les pipelines":
                    fig_cumul_n2, _, tickets_en_cours = graph_tickets_n2_cumulatif(
                        filtrer_par_periode(df_tickets_processed, periode_selectbox),
                        pipeline_selection
                    )
            
            st.plotly_chart(fig_cumul_n2, use_container_width=True)
            
            # Bouton pour afficher les tickets en cours
            if not tickets_en_cours.empty:
                nb_tickets_en_cours = len(tickets_en_cours)
                st.markdown(f"**🔴 {nb_tickets_en_cours} ticket(s) N2 en cours**")
                
                # Bouton pour télécharger le DataFrame des tickets en cours
                if st.button(f"📥 Télécharger la liste des {nb_tickets_en_cours} tickets en cours", key="btn_download_tickets_en_cours"):
                    # Préparer les données pour le téléchargement
                    tickets_download = tickets_en_cours[['Ticket ID', 'Date', 'Statut du ticket']].copy()
                    tickets_download = tickets_download.rename(columns={
                        'Ticket ID': 'ID_Ticket',
                        'Date': 'Date_Derniere_Mise_A_Jour',
                        'Statut du ticket': 'Statut'
                    })
                    tickets_download = tickets_download.sort_values('Date_Derniere_Mise_A_Jour', ascending=False)
                    
                    # Convertir en CSV
                    csv_data = tickets_download.to_csv(index=False, encoding='utf-8-sig')
                    
                    # Bouton de téléchargement
                    st.download_button(
                        label="📄 Télécharger CSV des tickets en cours",
                        data=csv_data,
                        file_name=f"tickets_n2_en_cours_{periode_selectbox.replace(' ', '_')}.csv",
                        mime="text/csv"
                    )
                    
                    # Afficher un aperçu du DataFrame
                    st.markdown("### 📋 Aperçu des tickets en cours")
                    st.dataframe(
                        tickets_download,
                        use_container_width=True,
                        column_config={
                            "ID_Ticket": st.column_config.TextColumn("ID Ticket", width="medium"),
                            "Date_Derniere_Mise_A_Jour": st.column_config.DateColumn("Date dernière mise à jour", width="medium"),
                            "Statut": st.column_config.TextColumn("Statut", width="small")
                        }
                    )
                
                if st.button(f"📋 Voir la liste complète des {nb_tickets_en_cours} tickets en cours", key="btn_tickets_en_cours"):
                    st.markdown("### 📋 Liste complète des tickets N2 en cours")
                    
                    # Préparer les données pour l'affichage complet
                    tickets_display = tickets_en_cours[['Ticket ID', 'Propriétaire du ticket', 'Pipeline', 'Statut du ticket', 'Date', 'Semaine']].copy()
                    tickets_display = tickets_display.sort_values('Date', ascending=False)
                    
                    # Afficher le tableau
                    st.dataframe(
                        tickets_display,
                        use_container_width=True,
                        column_config={
                            "Ticket ID": st.column_config.TextColumn("ID Ticket", width="medium"),
                            "Propriétaire du ticket": st.column_config.TextColumn("Propriétaire", width="medium"),
                            "Pipeline": st.column_config.TextColumn("Pipeline", width="small"),
                            "Statut du ticket": st.column_config.TextColumn("Statut", width="small"),
                            "Date": st.column_config.DateColumn("Date", width="small"),
                            "Semaine": st.column_config.TextColumn("Semaine", width="small")
                        }
                    )
                    
                    # Option pour télécharger la liste complète
                    csv_data_complete = tickets_display.to_csv(index=False, encoding='utf-8-sig')
                    st.download_button(
                        label="📥 Télécharger la liste complète en CSV",
                        data=csv_data_complete,
                        file_name=f"tickets_n2_en_cours_complet_{periode_selectbox.replace(' ', '_')}.csv",
                        mime="text/csv"
                    )
            else:
                st.success("✅ Aucun ticket N2 en cours !")

            # ----------- ÉVALUATIONS MANAGERIALES ------------------
            
            if submit:
                df_evaluation_filre = filter_evaluation(df_evaluation, agents, periodes_eval)
                
                st.markdown("## 📝 Évaluations managériales")
                if not df_evaluation_filre.empty:
                    moyenne_score = round(df_evaluation_filre["average_score"].mean(), 2)
                    st.markdown(f"**🎯 Score moyen des évaluations : `{moyenne_score}`**")
                else:
                    st.warning("Aucune donnée disponible pour la sélection.")
                
                st.dataframe(df_evaluation_filre, use_container_width=True)

            import matplotlib.pyplot as plt

            def render_evaluation_as_image(df, image_path="evaluation_image.png"):
                # Calcul de la taille de l'image selon le nombre de lignes
                height = max(2, len(df) * 0.4 + 2)
                fig, ax = plt.subplots(figsize=(12, height))
                ax.axis("off")

                # Titre
                plt.text(0, 1.05, "📝 Évaluations managériales", fontsize=16, fontweight='bold', transform=ax.transAxes)

                # Score moyen
                if not df.empty:
                    moyenne_score = round(df["average_score"].mean(), 2)
                    plt.text(0, 1.00, f"🎯 Score moyen des évaluations : {moyenne_score}", fontsize=13, transform=ax.transAxes)
                else:
                    plt.text(0, 1.00, "⚠️ Aucune donnée disponible pour la sélection.", fontsize=13, color="red", transform=ax.transAxes)

                # Tableau
                if not df.empty:
                    # Conversion en liste pour matplotlib
                    table_data = df.reset_index().values.tolist()
                    column_labels = df.reset_index().columns
                    table = ax.table(cellText=table_data, colLabels=column_labels, loc='center', cellLoc='center')
                    table.auto_set_font_size(False)
                    table.set_fontsize(10)
                    table.scale(1, 1.3)

                plt.tight_layout()
                plt.savefig(image_path, dpi=300, bbox_inches='tight')
                plt.close()

                return image_path
            
            df_eval_filtered = filter_evaluation(df_evaluation, agents, periodes_eval)
            render_evaluation_as_image(df_eval_filtered)


            # ----------- EXPORT -------------------
            st.markdown("## 📤 Export des KPI")

            col1, col2 = st.columns(2)

            with col1:
                if st.button("📥 Exporter les KPI en HTML et PDF"):
                    try:
                        html_file = generate_html_report_agent(df_support, df_tickets_processed, periode_selectbox, df_evaluation_filre)
                        st.success(f"Les KPI ont été exportés en HTML sous '{html_file}'.")

                        with open(html_file, "r", encoding="utf-8") as file:
                            html_content = file.read()

                        st.download_button("📄 Télécharger le HTML", data=html_content, file_name=os.path.basename(html_file), mime="text/html")

                        pdf_file = html_file.replace(".html", ".pdf")
                        if os.path.exists(pdf_file):
                            with open(pdf_file, "rb") as f:
                                pdf_bytes = f.read()
                            st.download_button("📄 Télécharger le PDF", data=pdf_bytes, file_name=os.path.basename(pdf_file), mime="application/pdf")
                        else:
                            st.warning("Le fichier PDF n'a pas été généré.")
                    except Exception as e:
                        st.error(f"Erreur lors de la génération du rapport : {e}")

            with col2:
                if st.button("📥 Exporter en PowerPoint"):
                    try:
                        st.info("Début génération PowerPoint (debug)")
                        st.write("KPIs:", kpis)
                        st.write("df_support shape:", df_support.shape)
                        st.write("Période:", periode)
                        pptx_io = create_powerpoint(kpis, df_support, periode)
                        st.info("PowerPoint généré (debug)")
                        st.download_button(
                            label="📥 Télécharger la présentation PowerPoint",
                            data=pptx_io,
                            file_name="dashboard_support.pptx",
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                        )
                        st.success("La présentation PowerPoint a été générée avec succès!")
                    except Exception as e:
                        st.error(f"Erreur lors de la génération de la présentation : {e}")

            
        agents()

    elif selection_page == "Tickets":
        import pandas as pd
        from data_processing.kpi_generation import filtrer_par_periode, generate_kpis, convert_to_sixtieth, filtrer_par_agent, charge_entrant_sortant
        from data_processing.hubspot_processing import process_hubspot_data
        from data_processing.aircall_processing import process_aircall_data, def_df_support, agents_all, line_tous

        def tickets():
            

            periode = st.selectbox("Sélectionnez une période :", 
                    ["1 an", "Toute la période", "6 derniers mois", "3 derniers mois", "Dernier mois"],
                    index=0)
            
            
            df_support = def_df_support(process_aircall_data(df_aircall), process_aircall_data(df_aircall), line_tous, agents_all)
            df_tickets = process_hubspot_data(df_tickets)
            kpis = generate_kpis(filtrer_par_periode(df_support, periode), filtrer_par_periode(df_tickets, periode), 'agents_all', None)

            st.plotly_chart(kpis['fig_activite_ticket'])

            st.plotly_chart(kpis['activite_ticket_semaine'])
            st.plotly_chart(kpis['activite_categorie'])

            partenaire = st.selectbox("Sélectionnez un partenaire :", ["FOLLOW", "Odaiji", "Help Info", "Oppysoft"], key="unique_partenaire_selection")

            kpis_partenaire = generate_kpis(filtrer_par_periode(df_support, periode), filtrer_par_periode(df_tickets, periode), 'agents_all', partenaire)

            col1, col2, col3= st.columns(3) 
            #col1.metric(f'% sla < 2h',int(kpis_partenaire['%_sla_2']))
            col2.metric(f'Délai moyen de réponse (heure)',round(kpis_partenaire['delai_moyen_reponse'], 2))

            st.plotly_chart(kpis_partenaire['sla_fig'])
            st.plotly_chart(kpis_partenaire['activite_categorie'])

        tickets()

    # --- RAPPORT POWERPOINT AGENTS ---
    from utils.powerpoint_helpers import create_powerpoint_agents_report
    from data_processing.kpi_generation import charge_entrant_sortant, charge_ticket, df_compute_ticket_appels_metrics
    st.sidebar.markdown("---")
    if st.sidebar.button("Générer rapport PowerPoint Agents"):
        # Texte méthodologie scoring (markdown)
        markdown_methodo = """## Paramètres de scoring\n\n- Objectif total de demandes par jour : 25\n- Ratio d'appels : 70%\n- Ratio de tickets : 30%\n- Objectif taux de service : 70%\n\nLe score est calculé selon le volume total traité, la répartition appels/tickets, la comparaison à la moyenne et le taux d'appels entrants."""
        # Préparation des données 6 mois
        periode_6m = "6 derniers mois"
        periode_3m = "3 derniers mois"
        agents_n1 = ['Olivier Sainte-Rose','Mourad HUMBLOT', 'Archimede KESSI', 'Celine Crendal']
        df_support_6m = filtrer_par_periode(def_df_support(process_aircall_data(df_aircall), process_aircall_data(df_aircall), line_tous, agents_all), periode_6m)
        df_tickets_6m = filtrer_par_periode(df_tickets_processed, periode_6m)
        df_scores_6m = df_compute_ticket_appels_metrics(agents_n1, df_tickets_6m, df_support_6m)
        # Préparation des données 3 mois
        df_support_3m = filtrer_par_periode(def_df_support(process_aircall_data(df_aircall), process_aircall_data(df_aircall), line_tous, agents_all), periode_3m)
        df_tickets_3m = filtrer_par_periode(df_tickets_processed, periode_3m)
        df_scores_3m = df_compute_ticket_appels_metrics(agents_n1, df_tickets_3m, df_support_3m)
        # Graphiques activité téléphone et tickets (6 mois)
        fig_tel_6m, _ = charge_entrant_sortant(df_support_6m, agents_n1)
        fig_ticket_6m, _ = charge_ticket(df_tickets_6m, agents_n1)
        # Génération du rapport
        pptx_io = create_powerpoint_agents_report(df_scores_6m, df_scores_3m, fig_tel_6m, fig_ticket_6m, markdown_methodo)
        st.sidebar.success("Rapport PowerPoint Agents généré !")
        st.sidebar.download_button(
            label="📥 Télécharger le rapport Agents",
            data=pptx_io,
            file_name="rapport_agents.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
    # --- FIN RAPPORT POWERPOINT AGENTS ---


def create_powerpoint(kpis, df_support, periode):
    # Créer une nouvelle présentation
    prs = Presentation()
    
    # Slide 1 - Titre
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Dashboard Support"
    subtitle.text = f"Période : {periode}"
    
    # Slide 2 - KPIs principaux
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = 'KPIs Principaux'
    
    tf = body_shape.text_frame
    tf.text = f"Taux de service : {kpis['Taux_de_service']}%"
    p = tf.add_paragraph()
    p.text = f"Appels entrants par jour : {kpis['Entrant']}"
    p = tf.add_paragraph()
    p.text = f"Numéros uniques par jour : {kpis['Numero_unique']}"
    p = tf.add_paragraph()
    p.text = f"Temps moyen par appel : {kpis['temps_moy_appel']}"
    
    # Slide 3 - Graphiques
    for fig in [kpis['charge_affid_stellair_%'], kpis['charge_affid_stellair_v'], kpis['evo_appels_tickets']]:
        img_slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(img_slide_layout)
        
        # Sauvegarder le graphique temporairement
        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
            fig.write_image(tmp.name)
            slide.shapes.add_picture(tmp.name, Inches(1), Inches(1), width=Inches(8))
    
    # Sauvegarder la présentation en mémoire
    pptx_io = io.BytesIO()
    prs.save(pptx_io)
    pptx_io.seek(0)
    
    return pptx_io

